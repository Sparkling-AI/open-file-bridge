#!/usr/bin/env python3
"""
File Bridge — single-file, stdlib-only local file service for Open WebUI.
No dependencies. Works with Python 3.8+. Package with PyInstaller for
double-click binaries on Windows/macOS/Linux.

Usage: python file_bridge.py [folder]      (default: ~/file-bridge-shared)
Then:  pick a folder at http://127.0.0.1:8765 (opens in browser)

Security model (v2, roadmap P0):
  Tier 1 (default): CORS origin lock. Only the configured Open WebUI origin
      may call the bridge from a browser. Set on first run via the picker.
  Tier 2 (opt-in):  org-wide bearer token (X-Bridge-Token header), set by
      the admin (installer / setup script / picker). Checked in addition.
  "Production mode" hard-fail: file endpoints refuse to serve while BOTH
      tiers are unconfigured — the local picker stays reachable so setup
      can be completed. (Stance adopted from Open Terminal v0.11.30.)

State lives in a per-OS state dir (env FILE_BRIDGE_STATE_DIR overrides):
  Linux: ~/.local/state/file-bridge   macOS: ~/Library/Application Support/file-bridge
  Windows: %APPDATA%\\file-bridge     (pattern: OpenWorker coworker/secrets.py)
  state.json    — root, ocr_lang, allowed_origin, token hash (0600)
  bridge-token  — plaintext token (0600, owner-only; never sent in responses)
"""
import base64
import binascii
import fnmatch
import hashlib
import hmac
import http.server
import io
import json
import os
import re
import secrets as _secrets
import shutil
import socketserver
import subprocess
import stat
import sys
import tempfile
import threading
import time
import webbrowser
import zipfile
from html import escape as _hesc
from pathlib import Path
from urllib.parse import urlparse, parse_qs, unquote

PORT = int(os.environ.get("FILE_BRIDGE_PORT", "8765"))
MAX_LIST = 500
MAX_READ = 200_000      # chars (text)
MAX_BINARY = 8_000_000  # bytes (base64 endpoints)

VERSION = "2.6"
SKILL_VERSION = "2.6"   # keep in sync with skill/open-file-bridge/SKILL.md


# ------------------------------------------- risk classes (P3, openworker risk.py)
# Every endpoint declares its INTRINSIC side-effect class — the substrate
# future confirmation/audit gating reads (policy asks ENDPOINT_RISK[path],
# never a hardcoded name list; pattern: openworker coworker/risk.py).
#
#   read        — no side effects: never mutates the workspace
#   write_local — mutates files in a shared root (confirm + snapshot +
#                 rate-breaker pipeline applies)
#   ui          — interacts with the user's desktop (consent-gated)
#   meta        — no file access at all (health/versions/settings)
#
# Subprocess note: /ocr, /ocr_pdf and /convert spawn local helper binaries
# (tesseract / LibreOffice) with fixed argument shapes. Their PRIMARY class
# stays read/write_local by workspace effect; the exec surface is documented
# here instead of a separate class because argument shape is fixed (no
# model-chosen commands — unlike openworker's run_shell).

class RiskClass:
    READ = "read"
    WRITE_LOCAL = "write_local"
    UI = "ui"
    META = "meta"


ENDPOINT_RISK = {
    # meta — no file access
    "/health": RiskClass.META,
    "/version": RiskClass.META,
    "/state": RiskClass.META,
    "/wheels": RiskClass.META,
    "/picker": RiskClass.META,
    "/api/root": RiskClass.META,
    "/api/preview": RiskClass.META,
    "/api/shutdown": RiskClass.META,   # local-only: stops the process
    # desktop interaction — consent = the local user clicking the button
    "/api/pick_folder": RiskClass.UI,  # local-only: native choose-folder dialog
    # reads
    "/list": RiskClass.READ,
    "/read": RiskClass.READ,
    "/peek": RiskClass.READ,
    "/read_b64": RiskClass.READ,
    "/stat": RiskClass.READ,
    "/pdf_text": RiskClass.READ,
    "/ocr": RiskClass.READ,          # spawns tesseract; writes nothing
    "/ocr/config": RiskClass.READ,
    "/image_info": RiskClass.READ,
    "/image_b64": RiskClass.READ,
    "/search": RiskClass.READ,
    "/html_text": RiskClass.READ,
    "/csv_head": RiskClass.READ,
    "/csv_stats": RiskClass.READ,
    "/xlsx_read": RiskClass.READ,
    "/docx_read": RiskClass.READ,
    "/pptx_read": RiskClass.READ,
    "/eml_read": RiskClass.READ,
    "/directory_tree": RiskClass.READ,
    "/versions/list": RiskClass.READ,
    "/trash/list": RiskClass.READ,
    # desktop interaction — consent-gated (allow_reveal)
    "/reveal": RiskClass.UI,
    # workspace mutation — confirm/snapshot/rate pipeline
    "/write": RiskClass.WRITE_LOCAL,
    "/write_b64": RiskClass.WRITE_LOCAL,
    "/edit": RiskClass.WRITE_LOCAL,
    "/delete": RiskClass.WRITE_LOCAL,
    "/write_many": RiskClass.WRITE_LOCAL,
    "/versions/restore": RiskClass.WRITE_LOCAL,
    "/trash/restore": RiskClass.WRITE_LOCAL,
    "/trash/purge": RiskClass.WRITE_LOCAL,   # 403 by design, but declared
    "/zip": RiskClass.WRITE_LOCAL,
    "/unzip": RiskClass.WRITE_LOCAL,
    "/ocr_pdf": RiskClass.WRITE_LOCAL,       # spawns tesseract → writes a PDF
    "/pdf_op": RiskClass.WRITE_LOCAL,
    "/pdf_from_text": RiskClass.WRITE_LOCAL,
    "/docx_merge": RiskClass.WRITE_LOCAL,
    "/docx_mailmerge": RiskClass.WRITE_LOCAL,
    "/docx_write": RiskClass.WRITE_LOCAL,
    "/pptx_from_template": RiskClass.WRITE_LOCAL,
    "/xlsx_append": RiskClass.WRITE_LOCAL,
    "/convert": RiskClass.WRITE_LOCAL,       # spawns soffice → writes a file
}

_IS_WINDOWS = sys.platform == "win32"


# ---------------------------------------------------------------- state dir

def state_dir() -> Path:
    """Per-OS state directory (pattern borrowed from OpenWorker secrets.py:
    env override > native app-data location)."""
    base = os.environ.get("FILE_BRIDGE_STATE_DIR")
    if base:
        p = Path(base).expanduser()
    elif _IS_WINDOWS:
        appdata = os.environ.get("APPDATA")
        p = Path(appdata) / "file-bridge" if appdata else Path.home() / ".file-bridge"
    elif sys.platform == "darwin":
        p = Path.home() / "Library" / "Application Support" / "file-bridge"
    else:
        p = Path.home() / ".local" / "state" / "file-bridge"
    # Canonicalize (realpath) so the state-inside-root guards compare like
    # for like: on macOS /var is a symlink to /private/var, and an unresolved
    # state path (e.g. FILE_BRIDGE_STATE_DIR=/var/folders/...) would never
    # match a resolved root path — the containment check would be bypassable.
    p = p.resolve()
    p.mkdir(parents=True, exist_ok=True)
    return p


def _restrict_to_user(path: Path, *, is_dir: bool) -> None:
    """Owner-only permissions. POSIX: mode bits (0700/0600). Windows has no
    mode bits — strip inherited ACLs and grant the current user alone
    (best-effort; pattern from OpenWorker secrets.py)."""
    if _IS_WINDOWS:
        user = os.environ.get("USERNAME")
        if not user:
            return
        domain = os.environ.get("USERDOMAIN")
        account = f"{domain}\\{user}" if domain else user
        grant = f"{account}:(OI)(CI)F" if is_dir else f"{account}:F"
        try:
            subprocess.run(["icacls", str(path), "/inheritance:r", "/grant:r", grant],
                           capture_output=True, check=False)
        except OSError:
            pass
        return
    try:
        os.chmod(path, 0o700 if is_dir else 0o600)
    except OSError:
        pass


STATE_DIR = state_dir()
STATE_FILE = STATE_DIR / "state.json"
TOKEN_FILE = STATE_DIR / "bridge-token"

_STATE_LOCK = threading.RLock()  # reentrant: set_token holds it while calling _state_update


def _state_load() -> dict:
    try:
        return json.loads(STATE_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _state_update(**kv) -> dict:
    """Merge-update state atomically (v1 lost keys by rewriting the whole
    file — e.g. save_root wiped a previously saved ocr_lang)."""
    with _STATE_LOCK:
        st = _state_load()
        st.update(kv)
        tmp = STATE_FILE.with_suffix(".tmp")
        tmp.write_text(json.dumps(st, indent=2), encoding="utf-8")
        _restrict_to_user(tmp, is_dir=False)   # chmod BEFORE the rename
        os.replace(tmp, STATE_FILE)
        return st


# ------------------------------------------------- audit log (P2, JSONL)
# Every file-touching call gets one line in state_dir/audit.log: ts,
# endpoint, method, path, size, status + SECRET-SCRUBBED args. Pattern:
# openworker coworker/audit.py `_SECRET_KEYS` scrub. Append-only, 0600,
# best-effort (audit failure must never break serving). Not an API — the
# file is for the human owner; models can't read it (state dir is
# structurally outside every root).

AUDIT_FILE = STATE_DIR / "audit.log"
_AUDIT_LOCK = threading.Lock()
_AUDIT_MAX_BYTES = 5 * 1024 * 1024   # rotate to audit.log.1 once past 5 MB

_SECRET_KEYS = (
    "token", "secret", "password", "api_key", "apikey", "authorization",
    "access_token", "refresh_token", "bearer", "credential", "private_key",
)
_BODY_KEYS = ("content", "b64", "body", "text", "old_text", "new_text", "html")
_ARG_TRUNCATE = 200


def _audit_scrub(value):
    """Recursive copy with secrets → [redacted] and payload keys → size only."""
    if isinstance(value, dict):
        out = {}
        for k, v in value.items():
            lk = str(k).lower()
            if any(s in lk for s in _SECRET_KEYS):
                out[k] = "[redacted]"
            elif any(b == lk or lk.endswith("_" + b) for b in _BODY_KEYS):
                # never log file CONTENT — record shape, not bytes
                out[k] = f"<{type(v).__name__}:{len(v) if hasattr(v, '__len__') else '?'}>"
            else:
                out[k] = _audit_scrub(v)
        return out
    if isinstance(value, list):
        return [_audit_scrub(v) for v in value[:10]]
    if isinstance(value, str):
        s = value.replace("\n", "\\n")
        return s if len(s) <= _ARG_TRUNCATE else s[:_ARG_TRUNCATE - 3] + "..."
    return value


def audit_log(endpoint: str, *, method: str = "GET", path: str | None = None,
              args: dict | None = None, status: int = 200,
              size: int | None = None, **extra) -> None:
    rec = {"ts": round(time.time(), 3), "endpoint": endpoint, "method": method,
           "status": status}
    if path is not None:
        rec["path"] = path[-400:]
    if args:
        rec["args"] = _audit_scrub(args)
    if size is not None:
        rec["size"] = size
    rec.update(extra)
    line = json.dumps(rec, default=str, ensure_ascii=False)
    try:
        with _AUDIT_LOCK:
            if AUDIT_FILE.exists() and AUDIT_FILE.stat().st_size > _AUDIT_MAX_BYTES:
                # simple rotation: shift .1 → dropped, current → .1
                old = AUDIT_FILE.with_suffix(".log.1")
                try:
                    old.unlink()
                except OSError:
                    pass
                try:
                    AUDIT_FILE.rename(old)
                except OSError:
                    pass
            with open(AUDIT_FILE, "a", encoding="utf-8") as fh:
                fh.write(line + "\n")
            try:
                os.chmod(AUDIT_FILE, 0o600)
            except OSError:
                pass
    except Exception:
        pass  # audit is best-effort: never fail a user request over it


# ------------------------------------------------------------- token store

def _hash_token(tok: str) -> str:
    return "sha256:" + hashlib.sha256(tok.encode()).hexdigest()


def get_configured_token() -> str | None:
    """Plaintext token for internal checks only (env override > token file).
    NEVER returned in any HTTP response."""
    env = os.environ.get("FILE_BRIDGE_TOKEN")
    if env:
        return env
    try:
        t = TOKEN_FILE.read_text(encoding="utf-8").strip()
        # ignore stale token file if state has no matching hash
        if t and _state_load().get("token_hash") == _hash_token(t):
            return t
    except Exception:
        pass
    return None


def set_token(tok: str | None) -> bool:
    """Store (or clear) the org-wide token. Returns success."""
    with _STATE_LOCK:
        if tok:
            tok = tok.strip()
            if not 8 <= len(tok) <= 256:
                return False
            TOKEN_FILE.write_text(tok, encoding="utf-8")
            _restrict_to_user(TOKEN_FILE, is_dir=False)
            _state_update(token_hash=_hash_token(tok))
        else:
            TOKEN_FILE.unlink(missing_ok=True)
            _state_update(token_hash=None)
    return True


def generate_token() -> str:
    tok = _secrets.token_urlsafe(24)
    set_token(tok)
    return tok


# ---------------------------------------------------------- origin (tier 1)

def _normalize_origin(o: str) -> str | None:
    """scheme://host[:port] with default ports made explicit, or None."""
    try:
        u = urlparse(o.strip())
        if u.scheme not in ("http", "https") or not u.hostname:
            return None
        port = u.port or (443 if u.scheme == "https" else 80)
        return f"{u.scheme}://{u.hostname.lower()}:{port}"
    except Exception:
        return None


def get_allowed_origin() -> str | None:
    env = os.environ.get("FILE_BRIDGE_ALLOWED_ORIGIN")
    raw = env or _state_load().get("allowed_origin")
    return _normalize_origin(raw) if raw else None


def set_allowed_origin(o: str | None) -> tuple[bool, str | None]:
    """Validate + persist. Returns (ok, normalized)."""
    if not o:
        _state_update(allowed_origin=None)
        return True, None
    n = _normalize_origin(o)
    if not n:
        return False, None
    _state_update(allowed_origin=o.strip())
    return True, n


def _request_origin(headers) -> str | None:
    """Normalize the Origin header (empty/`null` → None)."""
    raw = headers.get("Origin")
    if not raw or raw.strip().lower() == "null":
        return None
    return _normalize_origin(raw)


# ------------------------------------------------------------ security gate

def security_mode() -> str:
    """'token+origin' | 'token' | 'origin' | 'UNLOCKED' (both tiers off)."""
    has_tok = get_configured_token() is not None
    has_org = get_allowed_origin() is not None
    if has_tok and has_org:
        return "token+origin"
    if has_tok:
        return "token"
    if has_org:
        return "origin"
    return "UNLOCKED"


def check_request(headers) -> tuple[bool, int | None, str]:
    """Gate for FILE endpoints (not the local picker).

    Returns (allowed, http_status_if_denied, reason).
      - production hard-fail: both tiers off → deny 503
      - token tier: token configured → every request must carry it (401)
      - origin tier: browser Origin must match when present (403);
        non-browser callers without Origin are served (local tools) —
        that residual risk is exactly what the token tier closes.
    """
    mode = security_mode()
    if mode == "UNLOCKED":
        return False, 503, ("bridge unlocked — open the File Bridge settings page "
                            "(picker) and set the allowed Open WebUI origin "
                            "(and optionally a token) before serving files")
    tok = get_configured_token()
    if tok:
        supplied = headers.get("X-Bridge-Token", "")
        if not supplied:
            supplied = ""
            auth = headers.get("Authorization", "")
            if auth.lower().startswith("bearer "):
                supplied = auth[7:]
        if not hmac.compare_digest(supplied.encode(), tok.encode()):
            return False, 401, "missing or invalid bridge token"
    allowed = get_allowed_origin()
    origin = _request_origin(headers)
    if allowed and origin and origin != allowed:
        return False, 403, (f"origin {origin} is not the configured Open WebUI "
                            f"origin")
    return True, None, ""


# ------------------------------------------------------------- add-on setup

# Wheel hosting: serve pure-Python wheels bundled next to this file in ./wheels/
# so Pyodide installs them from localhost instead of PyPI (fast + offline).
# (WHEELS_DIR is set after _app_dir() is defined, for frozen-build support)

# Optional PDF/OCR add-on. Two parts, each auto-detected:
#   PDF: `pip install pymupdf` (importable)          -> /pdf_text, + /ocr for PDFs
#   OCR: tesseract binary on PATH (or TESSERACT_CMD) -> /ocr
# Bridge stays fully functional without them; /health reports availability.
try:
    import fitz  # pymupdf
    HAVE_PYMUPDF = True
except ImportError:
    fitz = None
    HAVE_PYMUPDF = False

try:
    import pypdfium2  # optional: /pdf_text?mode=images raster mode
    HAVE_PDFIUM = True
except ImportError:
    pypdfium2 = None
    HAVE_PDFIUM = False


# ------------------------------------------- /pdf_text mode=images (P2)
# Rasterize pages for VISION models: each page → PNG data URL, 144 dpi
# (scale 2.0 of 72 dpi base), ≤100 pages (openworker pdf_support pattern).
# Hand-rolled PNG encoder — no Pillow dependency on the bridge side.

RASTER_SCALE = 2.0    # ~144 dpi; readable text without giant payloads
RASTER_MAX_PAGES = 100


def _encode_png(width: int, height: int, pixels: bytes, stride: int,
                channels: int) -> bytes:
    """Minimal PNG writer (RGB/RGBA 8-bit), stdlib zlib only (MIT,
    openworker coworker/pdf_support.py)."""
    import struct
    import zlib
    color_type = 6 if channels == 4 else 2
    row_bytes = width * channels
    scanlines = bytearray()
    for y in range(height):
        scanlines.append(0)  # filter: None
        start = y * stride
        scanlines.extend(pixels[start:start + row_bytes])

    def chunk(tag: bytes, payload: bytes) -> bytes:
        return (struct.pack(">I", len(payload)) + tag + payload
                + struct.pack(">I", zlib.crc32(tag + payload) & 0xFFFFFFFF))

    header = struct.pack(">IIBBBBB", width, height, 8, color_type, 0, 0, 0)
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", header)
            + chunk(b"IDAT", zlib.compress(bytes(scanlines), 6))
            + chunk(b"IEND", b""))


def pdf_rasterize(p: Path, pages_param: str, max_pages: int) -> tuple[int, dict]:
    """/pdf_text?mode=images backend: pages as PNG data URLs. Returns
    (http_code, response)."""
    try:
        doc = pypdfium2.PdfDocument(str(p))
    except Exception as e:
        return 500, {"error": f"pdfium could not open the document: {e}"}
    try:
        total = len(doc)
        wanted = list(range(total))
        if pages_param:
            wanted = []
            for part in pages_param.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    wanted.extend(range(int(a) - 1, min(int(b), total)))
                elif part:
                    wanted.append(int(part) - 1)
            wanted = sorted(set(w for w in wanted if 0 <= w < total))
        wanted = wanted[:max_pages]
        pages_out = []
        budget = MAX_BINARY  # guard total payload like the text path guards chars
        for i in wanted:
            bitmap = doc[i].render(scale=RASTER_SCALE, rev_byteorder=True)
            png = _encode_png(bitmap.width, bitmap.height, bytes(bitmap.buffer),
                              bitmap.stride, bitmap.n_channels)
            budget -= len(png)
            pages_out.append({"page": i + 1, "size": [bitmap.width, bitmap.height],
                              "png_b64": base64.b64encode(png).decode("ascii")})
            if budget <= 0:
                pages_out.append({"page": i + 1, "note": "truncated (byte budget)"})
                break
        return 200, {"path": str(p.name), "mode": "images", "page_count": total,
                     "pages": pages_out, "rendered": len(pages_out),
                     "note": "each page is a PNG data URL — decode png_b64 and "
                             "show to the vision model, or upload to the chat"}
    finally:
        doc.close()


# ------------------------------------------- /pdf_op split|merge|rotate (P3)
# pymupdf page surgery. Output atomically written + snapshotted, rate-breaker
# counted. Confirmation only when the OUT file already exists (mirrors /write).

def _parse_pages(spec: str, total: int) -> list:
    """'1-3,5' (1-indexed) -> sorted 0-based page indices; [] = all."""
    if not spec:
        return list(range(total))
    out = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            out.extend(range(int(a) - 1, min(int(b), total)))
        else:
            out.append(int(part) - 1)
    return sorted(set(w for w in out if 0 <= w < total))


def pdf_op(op: str, srcs: list, out: Path, out_root: Path, spec: str,
           angle: int) -> tuple[int, dict]:
    """split/merge/rotate. srcs = resolved abs Paths (1 for split/rotate,
    >=1 for merge)."""
    if op == "split":
        src = srcs[0]
        doc = fitz.open(src)
        try:
            total = doc.page_count
            wanted = _parse_pages(spec, total)
            if not wanted:
                return 400, {"error": f"no pages selected (doc has {total})"}
            base = out.with_suffix("")
            results = []
            nbytes = 0
            for i in wanted:
                nd = fitz.open()
                try:
                    nd.insert_pdf(doc, from_page=i, to_page=i)
                    data = nd.tobytes(deflate=True, garbage=3)
                finally:
                    nd.close()
                tf = Path(f"{base}.p{i + 1}.pdf")
                if nbytes + len(data) > MAX_BINARY:
                    return 413, {"error": f"split output too large (> {MAX_BINARY} bytes)"}
                okr, err = rate_check(len(data))
                if not okr:
                    return 429, {"error": err, "rate_limited": True}
                snap = snapshot_before_write(out_root, tf) if tf.exists() else None
                atomic_write_bytes(tf, data)
                nbytes += len(data)
                results.append({"file": tf.name, "page": i + 1, "bytes": len(data),
                                "snapshot": snap})
            return 200, {"ok": True, "pages_split": len(results),
                         "files": results,
                         "note": "one PDF per selected page: <out-base>.pN.pdf"}
        finally:
            doc.close()
    if op == "merge":
        merged = fitz.open()
        try:
            pages_info = []
            for s in srcs:
                d = fitz.open(s)
                try:
                    merged.insert_pdf(d)
                    pages_info.append({"file": s.name, "pages": d.page_count})
                finally:
                    d.close()
            data = merged.tobytes(deflate=True, garbage=3)
        finally:
            merged.close()
        if len(data) > MAX_BINARY:
            return 413, {"error": f"merged PDF too large (> {MAX_BINARY} bytes)"}
        snap = snapshot_before_write(out_root, out) if out.exists() else None
        okr, err = rate_check(len(data))
        if not okr:
            return 429, {"error": err, "rate_limited": True}
        atomic_write_bytes(out, data)
        return 200, {"ok": True, "written": str(out), "bytes": len(data),
                     "pages": sum(p["pages"] for p in pages_info),
                     "sources": pages_info, "snapshot": snap}
    if op == "rotate":
        src = srcs[0]
        doc = fitz.open(src)
        try:
            total = doc.page_count
            wanted = _parse_pages(spec, total)
            if not wanted:
                return 400, {"error": f"no pages selected (doc has {total})"}
            for i in wanted:
                doc[i].set_rotation((doc[i].rotation + angle) % 360)
            data = doc.tobytes(deflate=True, garbage=3)
        finally:
            doc.close()
        if len(data) > MAX_BINARY:
            return 413, {"error": f"rotated PDF too large (> {MAX_BINARY} bytes)"}
        snap = snapshot_before_write(out_root, out) if out.exists() else None
        okr, err = rate_check(len(data))
        if not okr:
            return 429, {"error": err, "rate_limited": True}
        atomic_write_bytes(out, data)
        return 200, {"ok": True, "written": str(out), "bytes": len(data),
                     "pages_rotated": len(wanted), "angle": angle,
                     "snapshot": snap}
    return 400, {"error": f"unknown op {op!r} (split|merge|rotate)"}


def _app_dir() -> Path:
    """Directory where the app's bundled assets live (wheels/, tessdata/,
    tesseract/). Under PyInstaller --onefile, __file__ points into a throwaway
    extraction dir, so frozen apps must use the executable's directory."""
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


def _find_tesseract():
    """Locate the tesseract binary. Order:
    1. TESSERACT_CMD env override
    2. bundled: <app>/tesseract/tesseract.exe      (Windows installer layout)
       or      <app>/tesseract/bin/tesseract        (unix layout)
    3. `tesseract` on PATH
    Returns (path, version_str) or (None, None)."""
    import shutil as _sh
    cands = []
    if os.environ.get("TESSERACT_CMD"):
        cands.append(os.environ["TESSERACT_CMD"])
    cands += [str(_app_dir() / "tesseract" / "tesseract.exe"),
              str(_app_dir() / "tesseract" / "bin" / "tesseract")]
    wh = _sh.which("tesseract")
    if wh:
        cands.append(wh)
    for cand in cands:
        if cand and os.path.isfile(cand) and os.access(cand, os.X_OK):
            try:
                out = subprocess.run([cand, "--version"], capture_output=True,
                                     text=True, timeout=10)
                ver = (out.stdout or out.stderr).splitlines()[0] if (out.stdout or out.stderr) else ""
                return cand, ver
            except Exception:
                continue
    return None, None


TESSERACT_BIN, TESSERACT_VER = _find_tesseract()
# tessdata dir override (bundled langs next to the app, or system default).
# Order: env TESSDATA_PREFIX > merged(bundled + user drop-in) > <app>/tessdata
# > <app>/tesseract/tessdata (bundled-engine layout) > tesseract default.
_ad = _app_dir()
WHEELS_DIR = _ad / "wheels"
# user drop-in dir: extra .traineddata files (tessdata_fast) land here and
# are picked up WITHOUT rebuilding the app — survives app updates too.
USER_TESSDATA_DIR = STATE_DIR / "tessdata"
MERGED_TESSDATA_DIR = STATE_DIR / "tessdata-merged"


def _tessdata_refresh() -> None:
    """If the user dropped .traineddata files into USER_TESSDATA_DIR, mirror
    the bundled tessdata + those files into MERGED_TESSDATA_DIR (tesseract
    takes exactly one tessdata dir; merge = copy, ~20 MB, once per start)."""
    if os.environ.get("TESSDATA_PREFIX"):
        return  # explicit env wins; user drop-in not consulted
    bundled = _ad / "tessdata"
    if not bundled.is_dir():
        return
    user_files = sorted(USER_TESSDATA_DIR.glob("*.traineddata")) \
        if USER_TESSDATA_DIR.is_dir() else []
    if not user_files:
        return
    import shutil as _sh
    MERGED_TESSDATA_DIR.mkdir(parents=True, exist_ok=True)
    # refresh when the bundled tree or the user file set changed
    sig = "|".join([str(f.stat().st_mtime) for f in bundled.rglob("*") if f.is_file()]
                   + [f.name for f in user_files])
    marker = MERGED_TESSDATA_DIR / ".merged-sig"
    try:
        if marker.read_text() == sig:
            return
    except OSError:
        pass
    try:
        _sh.copytree(bundled, MERGED_TESSDATA_DIR, dirs_exist_ok=True)
        for f in user_files:
            _sh.copy2(f, MERGED_TESSDATA_DIR / f.name)
        marker.write_text(sig)
    except OSError:
        pass  # best-effort; the bundled set still works


_tessdata_refresh()
TESSDATA_DIR = (Path(os.environ["TESSDATA_PREFIX"]) if os.environ.get("TESSDATA_PREFIX")
                else MERGED_TESSDATA_DIR if (MERGED_TESSDATA_DIR / "eng.traineddata").exists()
                else _ad / "tessdata" if (_ad / "tessdata" / "eng.traineddata").exists()
                else _ad / "tesseract" / "tessdata" if (_ad / "tesseract" / "tessdata" / "eng.traineddata").exists()
                else None)  # None = let tesseract use its compiled-in default


def _ocr_langs_available():
    """List installed language codes from the active tessdata dir."""
    if TESSDATA_DIR and TESSDATA_DIR.is_dir():
        return sorted(f.stem for f in TESSDATA_DIR.glob("*.traineddata"))
    if TESSERACT_BIN:
        try:
            out = subprocess.run([TESSERACT_BIN, "--list-langs"],
                                 capture_output=True, text=True, timeout=10)
            langs = [l.strip() for l in (out.stdout or "").splitlines()[1:]
                     if l.strip() and ":" not in l]
            return langs
        except Exception:
            pass
    return []


def _get_ocr_lang():
    """Configured OCR language(s), e.g. 'swe+eng'. Saved in state file."""
    return _state_load().get("ocr_lang", "eng")


def _set_ocr_lang(lang: str):
    _state_update(ocr_lang=lang)


def _safe_walk(root: Path, *, max_entries: int = 5000,
               deadline_s: float = 10.0) -> tuple[list, bool]:
    """Recursive walk over a USER-SHARED folder, hardened:
    - never follows or lists symlinked dirs/files (pathlib's rglob FOLLOWS
      directory symlinks when recursing — a symlink cycle in a shared
      folder would loop a request thread forever);
    - stops at max_entries or deadline_s and reports truncated=True, so a
      pathological folder returns a partial answer instead of hanging.
      (A readdir blocked at the kernel level — stalled iCloud/mount —
      cannot be preempted in-thread; the deadline bounds the walk once
      syscalls start returning.)
    Returns (paths incl. dirs, truncated)."""
    out: list[Path] = []
    truncated = False
    t0 = time.time()
    for dirpath, dirnames, filenames in os.walk(root, followlinks=False):
        keep = []
        for d in sorted(dirnames):
            dp = Path(dirpath) / d
            if dp.is_symlink():
                continue          # prune: never descend, never list
            keep.append(d)
            out.append(dp)
            if len(out) >= max_entries:
                return out, True
        dirnames[:] = keep
        for name in sorted(filenames):
            f = Path(dirpath) / name
            if f.is_symlink():
                continue
            out.append(f)
            if len(out) >= max_entries:
                return out, True
        if time.time() - t0 > deadline_s:
            truncated = True
            break
    return out, truncated


# ------------------------------------------------------------------- paths

def safe_child(base: Path, rel: str) -> Path:
    """Resolve rel under base with traversal protection (for tmp files)."""
    p = (base / rel).resolve()
    if os.path.commonpath([str(base.resolve()), str(p)]) != str(base.resolve()):
        raise PermissionError("path escapes allowed root")
    return p


def load_root():
    try:
        p = _state_load()["root"]
        p = Path(p)
        return p if p.is_dir() else None
    except Exception:
        return None


def save_root(p: Path):
    _state_update(root=str(p.resolve()))


def resolve_safe(root: Path, rel: str) -> Path:
    p = (root / rel).resolve()
    if os.path.commonpath([str(root.resolve()), str(p)]) != str(root.resolve()):
        raise PermissionError("path escapes shared root")
    return p


# ------------------------------------------- atomic writes (P2, temp+rename)
# A crash mid-write must never leave a torn file at a real path: write to a
# same-dir temp file, fsync, then os.replace (atomic on POSIX + Win on same
# volume). Final-path symlink protection: resolve_safe() rejects escape, but
# the LAST hop could still be a symlink swapped in after resolution — for
# overwrite targets we refuse to replace through one.

def atomic_write_bytes(target: Path, data: bytes, *, preserve_mode: bool = True):
    """Atomic binary write: same-dir tmp → fsync → chmod → os.replace."""
    target.parent.mkdir(parents=True, exist_ok=True)
    if target.is_symlink():
        raise PermissionError("refusing to write through a symlink")
    mode = None
    if preserve_mode and target.exists():
        mode = stat.S_IMODE(target.stat().st_mode)
    fd, tmp = tempfile.mkstemp(dir=str(target.parent), prefix=".fb-tmp-")
    try:
        with os.fdopen(fd, "wb") as fh:
            fh.write(data)
            fh.flush()
            os.fsync(fh.fileno())
        os.chmod(tmp, mode if mode is not None else 0o644)
        os.replace(tmp, target)
    except BaseException:
        try:
            os.unlink(tmp)
        except OSError:
            pass
        raise


def atomic_write_text(target: Path, text: str, *, preserve_mode: bool = True):
    """Atomic text write (UTF-8). See atomic_write_bytes."""
    atomic_write_bytes(target, text.encode("utf-8"), preserve_mode=preserve_mode)


# ----------------------------------------------------- read safety (P0 #3)

# /read serves ONLY these extensions — fail-closed on anything else (unknown
# ext → 415 + hint to /peek). Stance validated by openworker readonly.py.
TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json", ".jsonl",
    ".xml", ".yml", ".yaml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".log", ".py", ".pyi", ".js", ".mjs", ".ts", ".jsx", ".tsx", ".html",
    ".htm", ".css", ".scss", ".less", ".svg", ".sql", ".sh", ".bash",
    ".zsh", ".fish", ".ps1", ".bat", ".cmd", ".make", ".mk", ".dockerfile",
    ".tex", ".srt", ".vtt", ".sub", ".c", ".h", ".cpp", ".hpp", ".cc",
    ".java", ".rs", ".go", ".rb", ".php", ".pl", ".pm", ".lua", ".r",
    ".swift", ".kt", ".kts", ".scala", ".vue", ".svelte", ".diff", ".patch",
}
# Extension-less files that are safely text (unix conventions).
KNOWN_BASENAMES = {
    "makefile", "dockerfile", "license", "licence", "readme", "changelog",
    "contributing", "authors", "notice", "codeowners", ".gitignore",
    ".gitattributes", ".dockerignore", ".editorconfig", ".env", ".npmrc",
    ".gitmodules",
}
OFFICE_ZIP_EXTS = {".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp", ".epub"}

# Sensitive-name blacklist (P2): matched against the BASENAME, case-insensitive,
# for reads AND writes. Complements per-root ignore lists (which the user
# configures) with a floor the model can never talk its way past. Entries
# deliberately broad — a false positive costs one "ask the user", a false
# negative leaks credentials into model context.
SENSITIVE_NAMES = {
    ".env", ".env.local", ".env.production", ".env.development",
    "id_rsa", "id_dsa", "id_ecdsa", "id_ed25519", "id_ecdsa_sk", "id_ed25519_sk",
    "authorized_keys", "known_hosts", "ssh_config", "config.ssh",
    "credentials", "credentials.json", "credentials.xml", "secrets",
    "secrets.json", "secrets.yaml", "secrets.yml", "secrets.toml",
    ".npmrc", ".pypirc", ".netrc", ".htpasswd", ".git-credentials",
    "serviceaccount.json", "firebase-adminsdk",
}
SENSITIVE_EXTS = {".pem", ".key", ".p12", ".pfx", ".keystore", ".kdbx", ".env"}
SENSITIVE_PATTERNS = re.compile(
    r"(^|[-_])(id_rsa|id_dsa|id_ecdsa|id_ed25519|authorized_keys|"
    r"credentials|secret|token|\.htpasswd|\.netrc|\.pypirc|\.npmrc)",
    re.IGNORECASE,
)


def sensitive_name(name: str) -> bool:
    """True if the basename looks like a credential/secret carrier."""
    base = name.rsplit("/", 1)[-1].lower()
    if base in SENSITIVE_NAMES:
        return True
    stem, dot, ext = base.rpartition(".")
    if dot and ("." + ext) in SENSITIVE_EXTS:
        return True
    return bool(SENSITIVE_PATTERNS.search(base))

# magic-byte signatures → coarse kind (checked BEFORE the extension whitelist:
# a .txt that is really a PNG gets rejected with the right routing hint)
_MAGIC = [
    (b"%PDF-", "pdf"),
    (b"PK\x03\x04", "zip"), (b"PK\x05\x06", "zip"), (b"PK\x07\x08", "zip"),
    (b"\x89PNG", "image"), (b"\xff\xd8\xff", "image"), (b"GIF87a", "image"),
    (b"GIF89a", "image"), (b"BM", "image"), (b"II*\x00", "image"),
    (b"MM\x00*", "image"), (b"RIFF", "binary"), (b"\x7fELF", "binary"),
    (b"MZ", "binary"), (b"\xca\xfe\xba\xbe", "binary"),
    (b"SQLite format 3\x00", "binary"), (b"OggS", "binary"),
    (b"fLaC", "binary"), (b"ID3", "binary"), (b"\x1f\x8b", "binary"),
    (b"\xfd7zXZ\x00", "binary"), (b"BZh", "binary"),
]


def sniff_kind(head: bytes):
    for sig, kind in _MAGIC:
        if head.startswith(sig):
            return kind
    return None


def _routing_hint(kind: str, ext: str) -> str:
    if kind == "pdf":
        return "PDF detected — use /pdf_text (text layer) or /ocr (scanned)"
    if kind == "zip" and ext in OFFICE_ZIP_EXTS:
        return (f"{ext} is a zip-based Office file — use /read_b64 and the "
                f"Pyodide office stack (see skill)")
    if kind == "zip":
        return "zip archive — use /read_b64 (Pyodide can open it via zipfile)"
    if kind == "image":
        return "image file — use /ocr (for text in it) or /read_b64"
    return "binary file — /read_b64 if you truly need the bytes"


# /read windowing (pattern: openworker coworker/tools/files.py)
DEFAULT_MAX_LINES = 2000
MAX_LINE_CHARS = 500


def _windowed_read(p: Path, start_line: int, max_lines: int) -> dict:
    start = start_line if start_line > 0 else 1
    n = max_lines if max_lines > 0 else DEFAULT_MAX_LINES
    n = min(n, DEFAULT_MAX_LINES)
    selected, total, budget = [], 0, MAX_READ
    with open(p, "r", encoding="utf-8", errors="replace") as fh:
        for i, line in enumerate(fh, 1):
            total = i
            if i < start or len(selected) >= n:
                continue
            text = line.rstrip("\n")
            if len(text) > MAX_LINE_CHARS:
                text = text[:MAX_LINE_CHARS] + "… (line truncated)"
            if budget - (len(text) + 8) < 0:  # overall char backstop
                selected.append(f"{i:>6}\t…(char budget reached — narrow the window)")
                budget = -1
                break
            budget -= len(text) + 8
            selected.append(f"{i:>6}\t{text}")
    end = start + len(selected) - 1 if selected else start - 1
    out = {
        "path": str(p.name), "start_line": start, "end_line": end,
        "total_lines": total, "content": "\n".join(selected),
    }
    if end < total:
        out["note"] = (f"showing lines {start}-{end} of {total}; "
                       f"call again with start_line={end + 1} to continue")
    return out


# ------------------------------------------------- versions + confirmations

CONFIRM_TTL_SECONDS = 60
_CONFIRM_FILE = STATE_DIR / "pending-confirmations.json"
_CONFIRM_LOCK = threading.Lock()

# writes that hit an EXISTING target snapshot it first (P0/P0b: 'makes model
# edits reversible'). Versions live OUTSIDE all roots — structurally
# unreachable via the path resolver (user decision 2026-08-27).
VERSIONS_DIR = STATE_DIR / "versions"


def _root_key(root: Path) -> str:
    return hashlib.sha256(str(root).encode()).hexdigest()[:12]


def _versions_root_for(root: Path) -> Path:
    d = VERSIONS_DIR / _root_key(root)
    return d


def snapshot_before_write(root: Path, target: Path) -> dict | None:
    """Copy current contents to versions store (root-hash scoped). Returns
    metadata or None (new file / unreadable). Cap: 8 MB per snapshot,
    512 MB total store (oldest pruned)."""
    try:
        if not target.is_file():
            return None
        size = target.stat().st_size
        if size > MAX_BINARY:
            return {"skipped": "too_large", "size": size}
        ts = time.strftime("%Y%m%d-%H%M%S") + "-" + _secrets.token_hex(2)
        rel = target.relative_to(root).as_posix()
        vdir = _versions_root_for(root)
        dest_dir = vdir / ts / rel.lstrip("/")
        dest_dir.parent.mkdir(parents=True, exist_ok=True)
        _restrict_to_user(vdir, is_dir=True)
        shutil.copy2(target, dest_dir)
        # manifest entry (append-only JSONL, one line per snapshot)
        with open(vdir / "manifest.jsonl", "a", encoding="utf-8") as fh:
            fh.write(json.dumps({"ts": ts, "path": rel, "size": size,
                                 "snapshot": str(dest_dir)}) + "\n")
        _prune_versions(vdir)
        return {"ts": ts, "path": rel, "size": size}
    except Exception as e:
        return {"skipped": f"error: {e}"}


def _prune_versions(vdir: Path, max_bytes: int = 512 * 1024 * 1024):
    """Keep the store under ~max_bytes by dropping oldest snapshot dirs."""
    try:
        snaps = sorted((d for d in vdir.glob("*/") if d.is_dir()))
        total = sum(f.stat().st_size for d in snaps for f in d.rglob("*") if f.is_file())
        i = 0
        while total > max_bytes and i < len(snaps):
            sz = sum(f.stat().st_size for f in snaps[i].rglob("*") if f.is_file())
            shutil.rmtree(snaps[i], ignore_errors=True)
            total -= sz
            i += 1
    except Exception:
        pass


def versions_list(root: Path, rel_filter: str = "") -> list:
    """METADATA ONLY (path/ts/size) — contents are never served (user
    decision: old/removed content must not re-enter model reach)."""
    out = []
    mf = _versions_root_for(root) / "manifest.jsonl"
    try:
        for line in open(mf, encoding="utf-8").read().splitlines():
            try:
                e = json.loads(line)
            except Exception:
                continue
            if rel_filter and rel_filter not in e.get("path", ""):
                continue
            e["restore"] = "POST /versions/restore {path: %r, ts: %r}" % (e["path"], e["ts"])
            out.append(e)
    except FileNotFoundError:
        pass
    out.reverse()  # newest first
    return out[:200]


def version_restore(root: Path, rel: str, ts: str) -> tuple[bool, str]:
    """Restore = a write: snapshot-first if target exists now, then copy back."""
    vdir = _versions_root_for(root)
    snap_root = vdir / ts / rel.lstrip("/")
    if not snap_root.is_file():
        # search by ts + suffix (snapshot dirs embed the tree)
        cands = list((vdir / ts).rglob(snap_root.name))
        if len(cands) != 1:
            return False, f"no unique snapshot for {rel} @ {ts}"
        snap_root = cands[0]
    target = (root / rel).resolve()  # same-root guard done by caller
    if target.exists():
        snapshot_before_write(root, target)  # don't lose current state either
    target.parent.mkdir(parents=True, exist_ok=True)
    # verified copy-then-rename (snapshot store may be another filesystem)
    fd, tmp = tempfile.mkstemp(dir=str(target.parent), prefix=".fb-restore-")
    os.close(fd)
    shutil.copy2(snap_root, tmp)
    os.chmod(tmp, stat.S_IMODE(snap_root.stat().st_mode))
    os.replace(tmp, target)
    return True, str(target)


# ---- two-step confirmation tokens (pattern: openapi-servers filesystem) ----

def _confirm_load() -> dict:
    try:
        data = json.loads(_CONFIRM_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {}
    now = time.time()
    out = {}
    for tok, det in data.items():
        try:
            if float(det["expiry"]) > now:
                out[tok] = det
        except Exception:
            continue
    return out


def _confirm_save(data: dict):
    tmp = _CONFIRM_FILE.with_suffix(".tmp")
    tmp.write_text(json.dumps(data), encoding="utf-8")
    _restrict_to_user(tmp, is_dir=False)
    os.replace(tmp, _CONFIRM_FILE)


def confirmation_issue(params: dict) -> dict:
    """Create a pending confirmation; returns the token payload (TTL 60 s)."""
    with _CONFIRM_LOCK:
        allc = _confirm_load()
        tok = _secrets.token_hex(4)  # 8 hex chars — relayed via chat
        allc[tok] = {"params": params, "expiry": time.time() + CONFIRM_TTL_SECONDS}
        _confirm_save(allc)
    return {"confirmation_token": tok,
            "expires_in": CONFIRM_TTL_SECONDS}


def confirmation_consume(tok: str, params: dict) -> tuple[bool, str]:
    """Validate token + params match. The token is burned on ANY consume
    attempt — mismatch included (stricter than the openapi-servers original,
    which allowed retries: a mismatch means the model changed its request)."""
    with _CONFIRM_LOCK:
        allc = _confirm_load()
        det = allc.pop(tok, None)
        _confirm_save(allc)
        if not det:
            return False, "invalid or expired confirmation token"
        if det["params"] != params:
            return False, ("request parameters do not match the original request "
                           "for this token (token consumed — request a new one)")
        return True, ""






# -------------------------------------------------- result cache (P2)
# /pdf_text + /ocr are the expensive endpoints (raster + tesseract ≈ 1s per
# page). History replays and chat re-asks hit the same file+params over and
# over — cache the computed result keyed by (sha256(content), op+params).
# Pattern: openworker coworker/pdf_support.py `_cached` LRU.

_CACHE: dict[tuple, dict] = {}
_CACHE_MAX = 16          # entries (each entry = one full result doc)
_CACHE_LOCK = threading.Lock()


def _file_digest(p: Path) -> str | None:
    try:
        h = hashlib.sha256()
        with open(p, "rb") as fh:
            for chunk in iter(lambda: fh.read(1 << 20), b""):
                h.update(chunk)
        return h.hexdigest()
    except OSError:
        return None


def _cache_key(p: Path, op: str, params: dict) -> tuple | None:
    d = _file_digest(p)
    if d is None:
        return None
    return (d, op, json.dumps(params, sort_keys=True, default=str))


def cache_get(p: Path, op: str, params: dict):
    """Returns (hit, value). hit=False → compute, then cache_put()."""
    key = _cache_key(p, op, params)
    if key is None:
        return False, None
    with _CACHE_LOCK:
        val = _CACHE.get(key)
        return val is not None, val


def cache_put(p: Path, op: str, params: dict, value: dict) -> None:
    key = _cache_key(p, op, params)
    if key is None:
        return
    with _CACHE_LOCK:
        if len(_CACHE) >= _CACHE_MAX and key not in _CACHE:
            _CACHE.pop(next(iter(_CACHE)))   # FIFO eviction, openworker-style
        _CACHE[key] = value


def cache_stats() -> dict:
    with _CACHE_LOCK:
        return {"entries": len(_CACHE), "max": _CACHE_MAX}


# ------------------------------------- zip / tree endpoints (P2, stdlib)

def zip_create(root: Path, cfg: dict, body: dict) -> tuple[int, dict]:
    """Zip files/dirs (root-relative) into <out> (root-relative .zip).
    Uses resolve_guarded for EVERY member → ignore rules + traversal
    protection apply to each path. Overwrites are atomic + snapshotted."""
    members = body.get("members")   # list of rel paths; missing = 400
    out_rel = body.get("out", "")
    if not isinstance(members, list) or not members or len(members) > 200:
        return 400, {"error": "members must be a 1-200 list of root-relative paths"}
    if not out_rel or Path(out_rel).suffix.lower() != ".zip":
        return 400, {"error": "out must end in .zip (root-relative)"}
    if Path(out_rel).name.startswith("."):
        return 400, {"error": "refusing to write hidden/dot file names"}

    # resolve ALL members first (fail before writing anything)
    resolved = []
    for m in members:
        if not isinstance(m, str) or not m:
            return 400, {"error": f"bad member: {m!r}"}
        try:
            p, r, c = resolve_guarded(unquote(m))
        except (PermissionError, ExcludedPath) as e:
            return 400, {"error": f"member not allowed: {m} ({e})"}
        if not p.exists():
            return 404, {"error": f"member not found: {m}"}
        resolved.append((m, p, r))
    op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)

    total = 0
    nfiles = 0
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for m, p, r in resolved:
            if p.is_file():
                data = p.read_bytes()
                total += len(data)
                if total > MAX_BINARY:
                    return 413, {"error": f"zipped payload too large (> {MAX_BINARY} bytes)"}
                zf.writestr(p.name, data)
                nfiles += 1
            else:
                for f in _safe_walk(p)[0]:
                    frel = f.relative_to(r).as_posix()
                    if _ignore_match(frel, f.is_dir(), _all_ignore(cfg)):
                        continue
                    if f.is_file():
                        data = f.read_bytes()
                        total += len(data)
                        if total > MAX_BINARY:
                            return 413, {"error": f"zipped payload too large (> {MAX_BINARY} bytes)"}
                        zf.writestr(f.name, data)
                        nfiles += 1
    snap = snapshot_before_write(oroot, op) if op.exists() else None
    okr, err = rate_check(total)
    if not okr:
        return 429, {"error": err, "rate_limited": True}
    atomic_write_bytes(op, buf.getvalue())
    return 200, {"ok": True, "written": str(op), "files": nfiles,
                 "bytes": buf.tell(), "snapshot": snap,
                 "note": "members stored FLAT (basename only) — directories recurse"}


def zip_extract(root: Path, cfg: dict, body: dict) -> tuple[int, dict]:
    """Unzip <path> (.zip, root-relative) under <dest>/ (root-relative dir).
    Every member name is sanitized: no absolute, no .., no drive letters,
    no symlink attrs — then extracted under resolve_guarded(dest)."""
    rel = body.get("path", "")
    dest_rel = body.get("dest", "")
    if not rel or not dest_rel:
        return 400, {"error": "need path (a .zip) + dest (directory)"}
    p, _r, _c = resolve_guarded(unquote(rel))
    if not p.is_file():
        return 404, {"error": f"no such file: {rel}"}
    if p.suffix.lower() != ".zip":
        return 400, {"error": f"not a zip: {p.suffix}"}
    dp, droot, dcfg = resolve_guarded(unquote(dest_rel), for_write=True)

    def safe_name(name: str) -> str | None:
        """Reject absolute/drive/../members outright (zip-slip), don't
        silently rewrite — a member that lies about its path aborts the
        whole extraction."""
        name = name.replace("\\", "/")
        if not name or name.startswith("/") or (len(name) > 1 and name[1] == ":"):
            return None
        parts = [x for x in name.split("/") if x not in ("", ".")]
        if not parts or any(x == ".." for x in parts):
            return None
        return "/".join(parts)

    try:
        with zipfile.ZipFile(p) as zf:
            bad = zf.testzip()
            if bad is not None:
                return 400, {"error": f"corrupt member: {bad}"}
            infos = zf.infolist()
            if len(infos) > 1000:
                return 413, {"error": f"too many entries: {len(infos)} > 1000"}
            total = sum(i.file_size for i in infos)
            if total > MAX_BINARY:
                return 413, {"error": f"unzipped payload too large (> {MAX_BINARY} bytes)"}
            okr, err = rate_check(total)
            if not okr:
                return 429, {"error": err, "rate_limited": True}
            count = 0
            for info in infos:
                if info.is_dir():
                    continue
                nm = safe_name(info.filename)
                if nm is None:
                    return 400, {"error": f"unsafe member name: {info.filename!r}"}
                target = safe_child(dp, nm)
                # ignore rules apply to the EXTRACTED path too
                frel = target.relative_to(droot).as_posix()
                if _ignore_match(frel, False, _all_ignore(dcfg)):
                    continue
                target.parent.mkdir(parents=True, exist_ok=True)
                with zf.open(info) as src, open(target, "wb") as dst:
                    shutil.copyfileobj(src, dst)
                if info.external_attr >> 16 & 0o170000 == 0o120000:
                    return 400, {"error": "refusing to extract symlink member"}
                count += 1
    except zipfile.BadZipFile:
        return 400, {"error": "not a valid zip archive"}
    return 200, {"ok": True, "dest": str(dp), "files": count}


def directory_tree(p: Path, cfg: dict, q: dict) -> dict:
    """Recursive tree (openapi-servers /directory_tree pattern), respecting
    ignore lists, with entry caps + depth limit for huge folders, plus a
    wall-clock budget: a single huge directory (100k entries) is fully
    readdir'd before any entry cap can bite, so once the budget elapses we
    stop DESCENDING further (children already listed still render) and
    report truncated. Caps pathological-but-allowed shares — mounting the
    whole disk is separately rejected by the state-dir containment rule."""
    max_entries = min(int(q.get("max_entries", "500") or 500), 2000)
    max_depth = min(int(q.get("max_depth", "6") or 6), 12)
    budget_s = min(max(float(q.get("budget_s", "1.5") or 1.5), 0.5), 10.0)
    t0 = time.time()
    pats = _all_ignore(cfg)
    truncated = False
    count = 0

    def build(cur: Path, depth: int, rel_prefix: str):
        nonlocal truncated, count
        entries = []
        try:
            children = sorted(cur.iterdir(), key=lambda c: (c.is_file(), c.name.lower()))
        except OSError:
            return entries
        for item in children:
            if count >= max_entries:
                truncated = True
                return entries
            rel = f"{rel_prefix}/{item.name}" if rel_prefix else item.name
            if _ignore_match(rel, item.is_dir(), pats):
                continue
            if item.is_symlink():
                continue  # symlinks never appear in the tree
            count += 1
            entry = {"name": item.name,
                     "type": "directory" if item.is_dir() else "file"}
            if item.is_file():
                try:
                    entry["size"] = item.stat().st_size
                except OSError:
                    pass
            elif depth < max_depth:
                if time.time() - t0 > budget_s:
                    truncated = True
                    entry["truncated"] = True   # listed, not expanded
                else:
                    entry["children"] = build(item, depth + 1, rel)
            else:
                entry["truncated"] = True
            entries.append(entry)
        return entries

    return {"path": q.get("path", "."), "entries": build(p, 0, ""),
            "entry_count": count, "truncated": truncated}


# ------------------------------------- /ocr_pdf (P2, searchable PDF)
# Scan → archive flow: rasterize pages, OCR each with tesseract's `pdf`
# renderer (page image + INVISIBLE text layer), merge the parts into one
# searchable PDF, write atomically. Result is a real file write → snapshot
# + confirmation flow applies when the output already exists.

_OCR_PDF_MAX_PAGES = 50


def _tesseract_env() -> dict:
    env = dict(os.environ)
    if TESSDATA_DIR:
        env["TESSDATA_PREFIX"] = str(TESSDATA_DIR)
    return env


def ocr_pdf(src: Path, out: Path, out_root: Path, lang: str, dpi: int,
            max_pages: int) -> tuple[int, dict]:
    """Create a searchable PDF from a scanned PDF or image. Returns
    (http_code, response). Only a %PDF- output file counts as tesseract
    success (arg errors print usage to stdout)."""
    tmpdir = tempfile.mkdtemp(prefix="fb-ocrpdf-")
    try:
        page_imgs: list[str] = []
        ext = src.suffix.lower()
        skipped = 0
        if ext in {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"}:
            page_imgs.append(str(src))
        elif ext == ".pdf":
            doc = fitz.open(src)
            try:
                n = min(doc.page_count, max_pages)
                for i in range(n):
                    pix = doc[i].get_pixmap(dpi=dpi)
                    tf = os.path.join(tmpdir, f"in-{i}.png")
                    pix.save(tf)
                    page_imgs.append(tf)
                skipped = doc.page_count - n
            finally:
                doc.close()
        else:
            return 400, {"error": f"ocr_pdf supports images and PDF, not {ext}"}
        if not page_imgs:
            return 400, {"error": "nothing to OCR (empty document?)"}
        parts: list[str] = []
        for idx, img in enumerate(page_imgs):
            base = os.path.join(tmpdir, f"pg-{idx:04d}")
            r = subprocess.run([TESSERACT_BIN, img, base, "-l", lang,
                                "--dpi", str(dpi), "pdf"],
                               capture_output=True, text=True, timeout=300,
                               env=_tesseract_env())
            outpdf = base + ".pdf"
            head = b""
            try:
                with open(outpdf, "rb") as fh:
                    head = fh.read(5)
            except OSError:
                pass
            if r.returncode != 0 or head != b"%PDF-":
                return 500, {"error": f"tesseract failed on page {idx + 1}: "
                                      + (r.stderr.strip() or r.stdout.strip())[:300]}
            parts.append(outpdf)
        merged = fitz.open()
        try:
            for part in parts:
                pdoc = fitz.open(part)
                try:
                    merged.insert_pdf(pdoc)
                finally:
                    pdoc.close()
            pdf_bytes = merged.tobytes(deflate=True, garbage=3)
        finally:
            merged.close()
        if len(pdf_bytes) > MAX_BINARY:
            return 413, {"error": f"searchable PDF too large: {len(pdf_bytes)} "
                                  f"> {MAX_BINARY} bytes"}
        snap = snapshot_before_write(out_root, out) if out.exists() else None
        okr, err = rate_check(len(pdf_bytes))
        if not okr:
            return 429, {"error": err, "rate_limited": True}
        atomic_write_bytes(out, pdf_bytes)
        resp = {"ok": True, "written": str(out), "pages": len(parts),
                "bytes": len(pdf_bytes), "lang": lang, "dpi": dpi,
                "snapshot": snap,
                "note": "searchable PDF: page images + invisible text layer — "
                        "/pdf_text and copy-paste search work on it now"}
        if skipped:
            resp["pages_skipped"] = skipped
        return 200, resp
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


# --------------------------------------------- /image_info (P3, stdlib)
# Dimensions + format + EXIF orientation for png/jpg/gif/webp/bmp —
# no Pillow needed (stdlib struct walk). Reports effective size after
# orientation so vision-model planning knows the real aspect ratio.

def _image_info(path: Path) -> tuple[int, dict]:
    IMG_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
    if path.suffix.lower() not in IMG_EXTS:
        return 400, {"error": f"not a supported image type: {path.suffix} "
                              f"(png/jpg/gif/webp/bmp)"}
    import struct

    def be(b, off, n):
        return int.from_bytes(b[off:off + n], "big")

    head = b""
    try:
        with open(path, "rb") as fh:
            head = fh.read(64 * 1024)
    except OSError as e:
        return 500, {"error": f"read failed: {e}"}
    if len(head) < 12:
        return 400, {"error": "file too small to be an image"}

    fmt = w = h = None
    orientation = None
    if head.startswith(b"\x89PNG\r\n\x1a\n"):
        fmt = "png"
        if head[12:16] == b"IHDR":
            w, h = be(head, 16, 4), be(head, 20, 4)
    elif head.startswith(b"GIF87a") or head.startswith(b"GIF89a"):
        fmt = "gif"
        # GIF logical screen descriptor is little-endian
        w, h = struct.unpack("<HH", head[6:10])
    elif head[:2] == b"BM":
        fmt = "bmp"
        w, h = struct.unpack("<ii", head[18:26])
    elif head[:4] == b"RIFF" and head[8:12] == b"WEBP":
        fmt = "webp"
        if head[12:16] == b"VP8 " and head[23:26] == b"\x9d\x01\x2a":
            w, h = struct.unpack("<HH", head[26:30])
            w &= 0x3fff
            h &= 0x3fff
        elif head[12:16] == b"VP8L":
            b0, b1, b2, b3 = head[21:25]
            w = 1 + (((b1 & 0x3F) << 8) | b0)
            h = 1 + (((b3 & 0xF) << 10) | (b2 << 2) | ((b1 & 0xC0) >> 6))
        elif head[12:16] == b"VP8X":
            w = 1 + (int.from_bytes(head[24:27], "little") << 0)
            h = 1 + (int.from_bytes(head[27:30], "little") << 0)
    elif head[:2] in (b"\xff\xd8",):
        fmt = "jpeg"
        off = 2
        # walk markers for SOFn + EXIF orientation
        while off + 4 < len(head):
            if head[off] != 0xFF:
                off += 1
                continue
            marker = head[off + 1]
            if marker in (0xD8, 0x01) or 0xD0 <= marker <= 0xD7:
                off += 2
                continue
            seglen = be(head, off + 2, 2)
            if 0xC0 <= marker <= 0xCF and marker not in (0xC4, 0xC8, 0xCC):
                h, w = be(head, off + 5, 2), be(head, off + 7, 2)
                break
            if marker == 0xE1 and head[off + 4:off + 10] == b"Exif\x00\x00":
                # TIFF header inside EXIF; find tag 0x0112 (Orientation)
                tiff = off + 10
                endian = head[tiff:tiff + 2]
                en = "<" if endian == b"II" else ">"
                ifd0 = struct.unpack(en + "I", head[tiff + 4:tiff + 8])[0]
                cnt_at = tiff + ifd0
                n = struct.unpack(en + "H", head[cnt_at:cnt_at + 2])[0]
                for k in range(n):
                    ent = cnt_at + 2 + k * 12
                    tag, typ, cnt = struct.unpack(en + "HHI", head[ent:ent + 8])
                    if tag == 0x0112:
                        val = struct.unpack(en + "H", head[ent + 8:ent + 10])[0]
                        orientation = val
                        break
            off += 2 + seglen
    if fmt is None or w is None:
        return 400, {"error": "could not parse image header (corrupt or "
                              "unsupported variant)"}
    out = {"path": path.name, "format": fmt, "width": w, "height": h,
           "megapixels": round(w * h / 1e6, 2),
           "size": path.stat().st_size}
    if orientation:
        out["exif_orientation"] = orientation
        swap = orientation in (5, 6, 7, 8)
        out["effective_width"], out["effective_height"] = (h, w) if swap else (w, h)
        out["note"] = "EXIF orientation present — effective dims swap w/h; " \
                      "viewers auto-rotate but raw decoders may not"
    return 200, out


# ------------------------------------------------- office reads (P1)

# Optional native readers; stdlib zipfile+XML fallbacks keep the bridge
# dependency-free (same pattern as the pymupdf add-on).
try:
    import openpyxl
    HAVE_OPENPYXL = True
except ImportError:
    openpyxl = None
    HAVE_OPENPYXL = False


def _cell_ref_to_idx(ref: str):
    """'B3' -> (col_idx_1based, row_idx_1based)."""
    import re as _r
    m = _r.match(r"([A-Z]+)([0-9]+)", ref.strip().upper())
    if not m:
        return None
    col = 0
    for ch in m.group(1):
        col = col * 26 + (ord(ch) - 64)
    return col, int(m.group(2))


def _xlsx_read_stdlib(path: Path, sheet: str | None, rng: str | None, max_rows: int):
    """Minimal xlsx reader: sharedStrings + sheet XML -> grid."""
    import zipfile
    import xml.etree.ElementTree as ET
    NS = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
    out = {"engine": "stdlib-xml"}
    with zipfile.ZipFile(path) as z:
        # shared strings
        shared = []
        if "xl/sharedStrings.xml" in z.namelist():
            root = ET.fromstring(z.read("xl/sharedStrings.xml"))
            for si in root.findall(f"{NS}si"):
                shared.append("".join(t.text or "" for t in si.iter(f"{NS}t")))
        # workbook sheet map
        wb = ET.fromstring(z.read("xl/workbook.xml"))
        sheets = [{"name": sh.get("name"),
                   "rid": sh.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")}
                  for sh in wb.iter(f"{NS}sheet")]
        rels = ET.fromstring(z.read("xl/_rels/workbook.xml.rels"))
        relmap = {r.get("Id"): r.get("Target") for r in rels}
        target_sheet = None
        if sheet:
            for sh in sheets:
                if sh["name"] == sheet:
                    target_sheet = sh
                    break
            if not target_sheet:
                return {"error": f"no such sheet: {sheet}", "sheets": [s["name"] for s in sheets]}
        else:
            target_sheet = sheets[0] if sheets else None
        if not target_sheet:
            return {"error": "workbook has no sheets"}
        t = relmap.get(target_sheet["rid"], "")
        t = t.lstrip("/")
        if not t.startswith("xl/"):
            t = "xl/" + t
        # absolute rel targets ("/xl/worksheets/...") and "worksheets/..." both
        # normalize to "xl/worksheets/..." — never double-prefix
        out["sheet"] = target_sheet["name"]
        out["sheets"] = [s["name"] for s in sheets]
        ws = ET.fromstring(z.read(t))
        # merged cells
        merged = [mc.get("ref") for mc in ws.iter(f"{NS}mergeCell")]
        out["merged_cells"] = merged[:200]
        # range bounds
        col_a = row_a = 1
        col_b = row_b = None
        if rng:
            try:
                a, b = rng.split(":")
                ca, ra = _cell_ref_to_idx(a)
                cb, rb = _cell_ref_to_idx(b)
                col_a, row_a, col_b, row_b = ca, ra, cb, rb
            except Exception:
                return {"error": f"bad range: {rng} (want A1:C10)"}
        rows_out, nrow = [], 0
        for row in ws.iter(f"{NS}row"):
            ri = int(row.get("r") or len(rows_out) + 1)
            if ri < row_a:
                continue
            if row_b and ri > row_b:
                break
            if nrow >= max_rows:
                out["truncated"] = True
                break
            cells = {}
            for c in row.findall(f"{NS}c"):
                ref = c.get("r") or ""
                t = c.get("t")
                v = c.find(f"{NS}v")
                is_el = c.find(f"{NS}is")
                if t == "s" and v is not None:
                    val = shared[int(v.text)]
                elif t == "inlineStr" and is_el is not None:
                    val = "".join(x.text or "" for x in is_el.iter(f"{NS}t"))
                elif v is not None:
                    val = v.text
                    if val is not None:
                        try:
                            f = float(val)
                            val = int(f) if f == int(f) else f
                        except ValueError:
                            pass
                else:
                    continue
                idx = _cell_ref_to_idx(ref)
                if idx:
                    ci = idx[0]
                    if ci < col_a or (col_b and ci > col_b):
                        continue
                    cells[ci] = val
            if cells:
                width = max(cells) - col_a + 1
                rows_out.append([cells.get(col_a + i) for i in range(width)])
            else:
                rows_out.append([])
            nrow += 1
    while rows_out and not any(rows_out[-1]):
        rows_out.pop()
    out["row_count"] = len(rows_out)
    out["data"] = rows_out
    return out


def _xlsx_read(path: Path, q: dict):
    sheet = q.get("sheet") or None
    rng = q.get("range") or None
    max_rows = min(int(q.get("max_rows", "200") or 200), 5000)
    if HAVE_OPENPYXL:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            names = wb.sheetnames
            if sheet and sheet not in names:
                return {"error": f"no such sheet: {sheet}", "sheets": names}
            ws = wb[sheet] if sheet else wb.active
            out = {"engine": "openpyxl", "sheet": ws.title, "sheets": names,
                   "merged_cells": [str(m) for m in
                                    (ws.merged_cells.ranges if hasattr(ws, "merged_cells") else [])][:200]}
            col_a = row_a = 1
            col_b = row_b = None
            if rng:
                try:
                    a, b = rng.split(":")
                    ca, ra = _cell_ref_to_idx(a)
                    cb, rb = _cell_ref_to_idx(b)
                    col_a, row_a, col_b, row_b = ca, ra, cb, rb
                except Exception:
                    return {"error": f"bad range: {rng} (want A1:C10)"}
            rows_out = []
            for ri, row in enumerate(ws.iter_rows(), 1):
                if ri < row_a:
                    continue
                if row_b and ri > row_b:
                    break
                if len(rows_out) >= max_rows:
                    out["truncated"] = True
                    break
                vals = []
                for ci, cell in enumerate(row, 1):
                    if ci < col_a:
                        continue
                    if col_b and ci > col_b:
                        break
                    vals.append(cell.value)
                rows_out.append(vals)
            while rows_out and not any(v is not None for v in rows_out[-1]):
                rows_out.pop()
            out["row_count"] = len(rows_out)
            out["data"] = rows_out
            return out
        finally:
            wb.close()
    return _xlsx_read_stdlib(path, sheet, rng, max_rows)


def _docx_read(path: Path, q: dict):
    """python-docx if present, else stdlib XML walk → markdown-ish text."""
    import zipfile
    import xml.etree.ElementTree as ET
    W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    out_lines = []
    with zipfile.ZipFile(path) as z:
        root = ET.fromstring(z.read("word/document.xml"))
        body = root.find(f"{W}body")
        for el in body.iter():
            if el.tag == f"{W}p":
                # heading?
                style = ""
                pPr = el.find(f"{W}pPr")
                if pPr is not None:
                    s = pPr.find(f"{W}pStyle")
                    if s is not None:
                        style = s.get(f"{W}val", "")
                text = "".join(t.text or "" for t in el.iter(f"{W}t")).strip()
                if not text:
                    continue
                sm = re.match(r"(?i)heading(\d)$", style) or re.match(r"(?i)titre(\d)$", style)
                if style.lower() == "title":
                    out_lines.append("# " + text)
                elif sm:
                    lvl = min(int(sm.group(1)), 6)
                    out_lines.append("#" * lvl + " " + text)
                elif style.startswith("List") or style == "BulletList":
                    out_lines.append("- " + text)
                else:
                    out_lines.append(text)
            elif el.tag == f"{W}tbl":
                # table: rows -> pipe rows
                rows = []
                for tr in el.findall(f"{W}tr"):
                    cells = []
                    for tc in tr.findall(f"{W}tc"):
                        cells.append("".join(t.text or "" for t in tc.iter(f"{W}t")).strip())
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    sep = "|" + "---|" * (rows[0].count("|") - 1)
                    rows.insert(1, sep)
                    out_lines.extend(rows)
    text = "\n".join(out_lines)
    max_chars = int(q.get("max_chars", "60000") or 60000)
    return {"path": q.get("path"), "chars": len(text),
            "truncated": len(text) > max_chars,
            "markdown": text[:max_chars]}


def _pptx_read(path: Path, q: dict):
    """Per-slide titles + text boxes via stdlib XML."""
    import zipfile
    import xml.etree.ElementTree as ET
    import re as _r
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    slides = []
    with zipfile.ZipFile(path) as z:
        names = sorted((n for n in z.namelist()
                        if _r.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
                       key=lambda n: int(_r.search(r"(\d+)", n).group(1)))
        for i, name in enumerate(names, 1):
            root = ET.fromstring(z.read(name))
            texts = []
            for p in root.iter(f"{A}p"):
                t = "".join(x.text or "" for x in p.iter(f"{A}t")).strip()
                if t:
                    texts.append(t)
            slides.append({"slide": i, "texts": texts,
                           "title": texts[0] if texts else None})
            if i >= 200:
                break
    return {"path": q.get("path"), "slide_count": len(slides), "slides": slides}


# -------------------------- office writes: docx_merge + pptx_from_template (P2)
# Real office demand (openworker issue #454): fill {{placeholders}} in a
# .docx template, and build a deck from a .potx/.pptx layout template.
# These WRITE office files → they need the native libs (stdlib can't
# author OOXML safely), get 501 without them, and run the standard
# snapshot + confirmation + rate-breaker write pipeline.

_PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Za-z0-9_.\-]+)\s*\}\}")


def _iter_docx_paragraphs(doc):
    """Paragraphs in the body + all tables (python-docx leaves table cells
    out of doc.paragraphs)."""
    for p in doc.paragraphs:
        yield p
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    yield p


def _fill_run_text(text: str, values: dict, missing: list):
    def sub(m):
        key = m.group(1)
        if key in values and key not in missing:
            return str(values[key])
        missing.append(key)
        return m.group(0)
    return _PLACEHOLDER_RE.sub(sub, text)


def _docx_placeholder_fill(data: bytes, values: dict) -> tuple[bytes, list]:
    from docx import Document
    bio = io.BytesIO(data)
    doc = Document(bio)
    missing = []
    for para in _iter_docx_paragraphs(doc):
        # python-docx merges adjacent same-format runs; replace per-run so
        # formatting is preserved. Multi-run placeholders are caught by the
        # joined-text pass below.
        for run in para.runs:
            if "{" in run.text:
                run.text = _fill_run_text(run.text, values, missing)
        joined = "".join(r.text for r in para.runs)
        if _PLACEHOLDER_RE.search(joined):
            # placeholder split across runs → rewrite whole paragraph text
            filled = _fill_run_text(joined, values, missing)
            for i, run in enumerate(para.runs):
                run.text = filled if i == 0 else ""
    # headers + footers
    for sec in doc.sections:
        for hf in (sec.header, sec.footer):
            for para in hf.paragraphs:
                for run in para.runs:
                    if "{" in run.text:
                        run.text = _fill_run_text(run.text, values, missing)
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue(), missing


def docx_merge(src: Path, out: Path, out_root: Path, values: dict,
               strict: bool) -> tuple[int, dict]:
    try:
        filled, missing = _docx_placeholder_fill(src.read_bytes(), values)
    except ImportError:
        return 501, {"error": "docx_merge needs python-docx "
                              "(pip install python-docx)"}
    except Exception as e:
        return 500, {"error": f"docx parse failed: {e}"}
    if strict and missing:
        return 400, {"error": "unresolved placeholders (strict mode)",
                     "missing": sorted(set(missing))}
    snap = snapshot_before_write(out_root, out) if out.exists() else None
    okr, err = rate_check(len(filled))
    if not okr:
        return 429, {"error": err, "rate_limited": True}
    atomic_write_bytes(out, filled)
    resp = {"ok": True, "written": str(out), "bytes": len(filled),
            "placeholders_filled": len(values), "snapshot": snap}
    if missing:
        resp["missing"] = sorted(set(missing))
    return 200, resp


def _pptx_fill_frame(frame, values: dict, missing: list):
    """Fill every paragraph of a text frame, preserving runs."""
    for para in frame.paragraphs:
        for run in para.runs:
            if "{" in run.text:
                run.text = _fill_run_text(run.text, values, missing)
        joined = "".join(r.text for r in para.runs)
        if _PLACEHOLDER_RE.search(joined):
            filled = _fill_run_text(joined, values, missing)
            for i, run in enumerate(para.runs):
                run.text = filled if i == 0 else ""


def pptx_from_template(src: Path, out: Path, out_root: Path, slides: list,
                       values: dict) -> tuple[int, dict]:
    """Copy template → optional global {{placeholder}} fill → optional
    per-slide dict list appended via template LAYOUTS (new slides copy the
    layout's placeholder frames, so corporate design survives)."""
    try:
        from pptx import Presentation
    except ImportError:
        return 501, {"error": "pptx_from_template needs python-pptx "
                              "(pip install python-pptx)"}
    try:
        prs = Presentation(str(src))
    except Exception as e:
        return 500, {"error": f"pptx parse failed: {e}"}
    missing = []
    # global fill over existing slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _pptx_fill_frame(shape.text_frame, values, missing)
    added = 0
    # append new slides from per-slide dicts. Fresh add_slide() placeholders
    # carry layout PROMPT text, not {{placeholders}} — so per-slide specs set
    # title/body directly (values-fill only applies to template's own slides).
    for spec in slides:
        layout_idx = spec.get("layout", 1)
        layouts = list(prs.slide_layouts)
        if not isinstance(layout_idx, int) or not (0 <= layout_idx < len(layouts)):
            return 400, {"error": f"bad layout index {layout_idx!r} — template "
                                  f"has {len(layouts)} layouts (0-"
                                  f"{len(layouts) - 1})"}
        slide = prs.slides.add_slide(layouts[layout_idx])
        if spec.get("title") is not None and slide.shapes.title is not None:
            slide.shapes.title.text = str(spec["title"])
        if spec.get("body") is not None:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    ph.text = str(spec["body"])
                    break
        added += 1
    bio = io.BytesIO()
    prs.save(bio)
    data = bio.getvalue()
    snap = snapshot_before_write(out_root, out) if out.exists() else None
    okr, err = rate_check(len(data))
    if not okr:
        return 429, {"error": err, "rate_limited": True}
    atomic_write_bytes(out, data)
    resp = {"ok": True, "written": str(out), "bytes": len(data),
            "slides_added": added, "layout_count": len(prs.slide_layouts),
            "snapshot": snap}
    if missing:
        resp["missing"] = sorted(set(missing))
    return 200, resp






# ---------------------- structured writes: pdf_from_text/docx_write/xlsx_append (P3)
# Bridge-native document WRITES (roadmap P3). Optional add-on libs (fpdf2,
# python-docx, openpyxl) — 501 when missing, same as the P2 office writes.
# All three run the standard pipeline: validate-cheap-first → confirmation
# (always, they write NEW files) → build → atomic write + snapshot + breaker.

def _sanitize_pdf_text(s: str) -> str:
    """fpdf2 core fonts are latin-1: map the common Windows-1252 range and
    drop anything else (replacement '?') so a stray emoji can't 500 the
    write. Honest lossiness beats a broken document."""
    repl = {"→": "->", "←": "<-", "–": "-", "—": "-", "‘": "'", "’": "'",
            "“": '"', "”": '"', "•": "*", "…": "...", "×": "x", "≈": "~",
            "\u00a0": " ", "\t": "    "}
    for k, v in repl.items():
        s = s.replace(k, v)
    return s.encode("latin-1", "replace").decode("latin-1")


def _pdf_font_for(style: str):
    from fpdf import FPDF
    return {"title": ("helvetica", "B", 20),
            "h1": ("helvetica", "B", 16),
            "h2": ("helvetica", "B", 13),
            "body": ("helvetica", "", 11)}.get(style)


def pdf_from_text(out: Path, out_root: Path, title: str, blocks: list,
                  page_size: str) -> tuple[int, dict]:
    try:
        from fpdf import FPDF
    except ImportError:
        return 501, {"error": "pdf_from_text needs the fpdf2 add-on "
                              "(pip install fpdf2)"}
    if not blocks:
        return 400, {"error": "no blocks given"}
    fmt = "A4" if page_size == "a4" else "letter"
    pdf = FPDF(format=fmt)
    pdf.set_auto_page_break(True, margin=18)
    pdf.set_title(_sanitize_pdf_text(title)[:200] if title else "")
    used = 0
    try:
        for blk in blocks:
            style = blk.get("style", "body")
            if style == "pagebreak":
                pdf.add_page(); used += 1
                continue
            if style not in ("title", "h1", "h2", "body"):
                return 400, {"error": f"unknown style {style!r} "
                                      f"(title|h1|h2|body|pagebreak)"}
            text = _sanitize_pdf_text(str(blk.get("text", "")))
            if not text.strip() and style == "body":
                text = ""
            if used == 0:
                pdf.add_page()
            font = _pdf_font_for(style)
            pdf.set_font(font[0], font[1], font[2])
            if style == "body":
                pdf.multi_cell(0, font[2] * 0.55, text)
                pdf.ln(2)
            else:
                pdf.multi_cell(0, font[2] * 0.6, text)
                pdf.ln(3)
            used += 1
            if used > 5000:
                return 413, {"error": "too many blocks (5000 cap)"}
        data = bytes(pdf.output())
    except Exception as e:
        return 500, {"error": f"pdf build failed: {e}"}
    if len(data) > MAX_BINARY:
        return 413, {"error": f"PDF too large (> {MAX_BINARY} bytes)"}
    snap = snapshot_before_write(out_root, out) if out.exists() else None
    okr, err = rate_check(len(data))
    if not okr:
        return 429, {"error": err, "rate_limited": True}
    atomic_write_bytes(out, data)
    return 200, {"ok": True, "written": str(out), "bytes": len(data),
                 "blocks": len(blocks), "page_size": fmt, "snapshot": snap}


_DOCX_LIST_STYLES = ("list", "bullets", "bullet", "numbered", "numbers", "ul", "ol")


def docx_write(out: Path, out_root: Path, title: str, sections: list) -> tuple[int, dict]:
    try:
        from docx import Document
    except ImportError:
        return 501, {"error": "docx_write needs python-docx "
                              "(pip install python-docx)"}
    if not sections:
        return 400, {"error": "no sections given"}
    doc = Document()
    if title:
        doc.add_heading(title, 0)
    counters = {"h1": 0, "h2": 0}
    try:
        for sec in sections:
            st = sec.get("style", "paragraph")
            text = str(sec.get("text", ""))
            if st == "h1":
                doc.add_heading(text, 1); counters["h1"] += 1
            elif st == "h2":
                doc.add_heading(text, 2); counters["h2"] += 1
            elif st == "pagebreak":
                doc.add_page_break()
            elif st in _DOCX_LIST_STYLES:
                numbered = st in ("numbered", "numbers", "ol")
                for item in ([text] if text and not isinstance(sec.get("items"), list)
                             else sec.get("items", [])):
                    doc.add_paragraph(str(item),
                                      style="List Number" if numbered else "List Bullet")
            else:
                doc.add_paragraph(text)
        bio = io.BytesIO()
        doc.save(bio)
        data = bio.getvalue()
    except Exception as e:
        return 500, {"error": f"docx build failed: {e}"}
    if len(data) > MAX_BINARY:
        return 413, {"error": f"docx too large (> {MAX_BINARY} bytes)"}
    snap = snapshot_before_write(out_root, out) if out.exists() else None
    okr, err = rate_check(len(data))
    if not okr:
        return 429, {"error": err, "rate_limited": True}
    atomic_write_bytes(out, data)
    return 200, {"ok": True, "written": str(out), "bytes": len(data),
                 "sections": len(sections), "headings": counters, "snapshot": snap}


def xlsx_append(path: Path, rows: list, sheet: str | None, out_root: Path,
                header: list | None) -> tuple[int, dict]:
    """Append rows to an existing .xlsx (or create it). openpyxl add-on."""
    try:
        import openpyxl
    except ImportError:
        return 501, {"error": "xlsx_append needs openpyxl (pip install openpyxl)"}
    if not rows:
        return 400, {"error": "no rows given"}
    created = not path.exists()
    try:
        wb = openpyxl.load_workbook(path) if not created else openpyxl.Workbook()
        ws = wb[sheet] if (sheet and sheet in wb.sheetnames) else \
             (wb.create_sheet(sheet) if sheet else wb.active)
        if created and header:
            ws.append(header)
        start_row = ws.max_row or 0
        for r in rows:
            ws.append(r)
        bio = io.BytesIO()
        wb.save(bio)
        data = bio.getvalue()
    except Exception as e:
        return 500, {"error": f"xlsx append failed: {e}"}
    if len(data) > MAX_BINARY:
        return 413, {"error": f"workbook too large (> {MAX_BINARY} bytes)"}
    snap = snapshot_before_write(out_root, path) if not created else None
    okr, err = rate_check(len(data))
    if not okr:
        return 429, {"error": err, "rate_limited": True}
    atomic_write_bytes(path, data)
    resp = {"ok": True, "written": str(path), "bytes": len(data),
            "rows_appended": len(rows), "sheet": ws.title,
            "created": created, "row_count": ws.max_row, "snapshot": snap}
    if created and header:
        resp["header"] = header
    return 200, resp


# --------------------------------- mail-merge: docx template + rows -> batch (P3)
# Natural extension of /docx_merge (openworker #454): same {{placeholder}}
# engine, but one document PER ROW. Rows come from an .xlsx (bridge reader),
# a .csv (stdlib), or inline JSON. Output: N docx files following an out
# pattern, or ONE zip when out ends in .zip.

def _mailmerge_rows(src: Path, q_cols: str | None) -> tuple[int, dict]:
    """Load merge rows: xlsx (header row = keys) | csv | inline handled by
    caller. Returns (code, {rows|error})."""
    ext = src.suffix.lower()
    if ext == ".xlsx":
        res = _xlsx_read(src, {"max_rows": "501"})
        if "error" in res:
            return 400, res
        rows_raw = res.get("data") or []
        if not rows_raw:
            return 400, {"error": "xlsx has no rows"}
        header = [str(h) if h is not None else "" for h in rows_raw[0]]
        rows = []
        for r in rows_raw[1:]:
            row = {}
            for i, h in enumerate(header):
                if h and i < len(r) and r[i] is not None:
                    row[h] = r[i]
            if row:
                rows.append(row)
        return 200, {"rows": rows}
    if ext == ".csv":
        import csv as _csv
        with open(src, newline="", encoding="utf-8", errors="replace") as fh:
            reader = list(_csv.reader(fh))
        if not reader:
            return 400, {"error": "csv is empty"}
        header = [h.strip() for h in reader[0]]
        rows = []
        for r in reader[1:]:
            row = {}
            for i, h in enumerate(header):
                if h and i < len(r):
                    row[h] = r[i]
            if row:
                rows.append(row)
        return 200, {"rows": rows}
    return 400, {"error": f"row source must be .xlsx or .csv, not {ext}"}


def _safe_name(s: str) -> str:
    """Filename-safe placeholder value (merged into out patterns)."""
    keep = []
    for ch in str(s):
        if ch.isalnum() or ch in "-_ ":
            keep.append(ch)
        else:
            keep.append("_")
    return "".join(keep).strip()[:80] or "row"


def docx_mailmerge(tpl: Path, out_rel: str, out_root: Path, root: Path,
                   rows: list, to_zip: bool) -> tuple[int, dict]:
    try:
        filled_list, all_missing = [], set()
        for i, row in enumerate(rows[:50]):
            if not isinstance(row, dict):
                return 400, {"error": f"row {i} is not an object"}
            vals = {str(k): str(v) for k, v in row.items()}
            data, missing = _docx_placeholder_fill(tpl.read_bytes(), vals)
            all_missing.update(missing)
            filled_list.append((vals, data))
    except ImportError:
        return 501, {"error": "docx_mailmerge needs python-docx "
                              "(pip install python-docx)"}
    except Exception as e:
        return 500, {"error": f"mail-merge failed: {e}"}
    total = 0
    results = []
    tpl_bytes = tpl.read_bytes()
    del filled_list
    # zip mode: write every docx into one archive (atomic, one breaker hit)
    if to_zip:
        import zipfile as _zf
        bio = io.BytesIO()
        with _zf.ZipFile(bio, "w", _zf.ZIP_DEFLATED) as z:
            for row in rows[:50]:
                vals = {str(k): str(v) for k, v in row.items()}
                data, missing = _docx_placeholder_fill(tpl_bytes, vals)
                all_missing.update(missing)
                name = _out_name_for(out_rel, vals)
                z.writestr(name, data)
                total += len(data)
                results.append({"file": name, "bytes": len(data)})
        data_out = bio.getvalue()
        if len(data_out) > MAX_BINARY:
            return 413, {"error": f"zip too large (> {MAX_BINARY} bytes)"}
        snap = snapshot_before_write(out_root, out_root / out_rel) \
            if (out_root / out_rel).exists() else None
        okr, err = rate_check(len(data_out))
        if not okr:
            return 429, {"error": err, "rate_limited": True}
        atomic_write_bytes(out_root / out_rel, data_out)
        resp = {"ok": True, "written": str(out_root / out_rel),
                "bytes": len(data_out), "documents": len(results),
                "files": results, "snapshot": snap}
        if all_missing:
            resp["missing"] = sorted(all_missing)
        return 200, resp
    # loose-file mode: one docx per row, out pattern names them
    for row in rows[:50]:
        vals = {str(k): str(v) for k, v in row.items()}
        data, missing = _docx_placeholder_fill(tpl_bytes, vals)
        all_missing.update(missing)
        rel = _out_name_for(out_rel, vals)
        tf = out_root / rel
        if total + len(data) > MAX_BINARY:
            return 413, {"error": f"output too large (> {MAX_BINARY} bytes)"}
        okr, err = rate_check(len(data))
        if not okr:
            return 429, {"error": err, "rate_limited": True,
                         "hint": "batch aborted — files written so far are "
                                 "listed"}
        snap = snapshot_before_write(out_root, tf) if tf.exists() else None
        atomic_write_bytes(tf, data)
        total += len(data)
        results.append({"file": rel, "bytes": len(data), "snapshot": snap})
    resp = {"ok": True, "documents": len(results), "bytes": total,
            "files": results}
    if all_missing:
        resp["missing"] = sorted(all_missing)
    return 200, resp


_PLACEHOLDER_NAME_RE = None


def _out_name_for(pattern: str, vals: dict) -> str:
    """Replace {{placeholders}} in the out pattern; default 'row-N.docx'."""
    global _PLACEHOLDER_NAME_RE
    if _PLACEHOLDER_NAME_RE is None:
        _PLACEHOLDER_NAME_RE = re.compile(r"\{\{\s*([A-Za-z0-9_.\-]+)\s*\}\}")
    def sub(m):
        k = m.group(1)
        return _safe_name(vals.get(k, "")) if k in vals else m.group(0)
    name = _PLACEHOLDER_NAME_RE.sub(sub, pattern)
    if _PLACEHOLDER_NAME_RE.search(name):
        # unresolved placeholder in the NAME → numbered fallback
        name = re.sub(r"\{\{[^}]*\}\}", "", name)
        name = name.strip("_- ")
    stem, dot, ext = name.rpartition(".")
    if not stem or not dot:
        return pattern
    return name


# ------------------------------------------------ .eml parsing (P3, stdlib)
# RFC-822 email files: headers + text body + attachment METADATA. stdlib
# email package — no add-on needed. Attachments are listed (name/size/type),
# never decoded into the response: raw attachment bytes can hold anything,
# and dumping binaries into model context is a /read_b64 decision.

def _eml_read(path: Path, q: dict) -> tuple[int, dict]:
    import email
    from email import policy
    try:
        with open(path, "rb") as fh:
            msg = email.message_from_binary_file(fh, policy=policy.default)
    except Exception as e:
        return 500, {"error": f"eml parse failed: {e}"}

    def _hdr(name):
        v = msg.get(name)
        return str(v) if v is not None else None

    out = {
        "subject": _hdr("Subject"),
        "from": _hdr("From"),
        "to": _hdr("To"),
        "cc": _hdr("Cc"),
        "date": _hdr("Date"),
        "message_id": _hdr("Message-ID"),
    }
    # parse the Date header into an ISO timestamp when possible
    if out["date"]:
        try:
            from email.utils import parsedate_to_datetime
            out["date_iso"] = parsedate_to_datetime(out["date"]).isoformat()
        except Exception:
            pass

    body_text, body_html = [], []
    attachments = []
    max_chars = min(int(q.get("max_chars", "20000") or 20000), MAX_READ)
    for part in msg.walk():
        if part.is_multipart():
            continue
        disp = (part.get_content_disposition() or "")
        ctype = part.get_content_type()
        fname = part.get_filename()
        if disp == "attachment" or (fname and disp != "inline"):
            payload = part.get_payload(decode=True) or b""
            attachments.append({"filename": fname or "unnamed",
                                "content_type": ctype,
                                "bytes": len(payload)})
            continue
        if ctype == "text/plain":
            try:
                body_text.append(part.get_content())
            except Exception:
                raw = part.get_payload(decode=True) or b""
                body_text.append(raw.decode(part.get_content_charset() or
                                            "utf-8", errors="replace"))
        elif ctype == "text/html":
            try:
                body_html.append(part.get_content())
            except Exception:
                raw = part.get_payload(decode=True) or b""
                body_html.append(raw.decode(part.get_content_charset() or
                                            "utf-8", errors="replace"))
    text = "\n\n".join(t for t in body_text if t and t.strip())
    if not text.strip() and body_html:
        # html-only mail: reuse the stdlib tag-stripper from /html_text
        import html.parser as _hp

        class _Strip(_hp.HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts = []
                self._skip = 0

            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip += 1

            def handle_endtag(self, tag):
                if tag in ("script", "style") and self._skip:
                    self._skip -= 1

            def handle_data(self, data):
                if not self._skip and data.strip():
                    self.parts.append(data)

        s = _Strip()
        s.feed("\n".join(body_html))
        text = "\n".join(s.parts)
        out["body_was_html"] = True
    if len(text) > max_chars:
        out["truncated"] = True
        text = text[:max_chars]
    out["body"] = text
    out["attachment_count"] = len(attachments)
    if attachments:
        out["attachments"] = attachments
    return 200, out


# ------------------------------------------------ LibreOffice /convert (P3)
# Headless conversion between document formats (roadmap P3): legacy
# .doc/.xls/.ppt -> modern OOXML, docx/xlsx/pptx -> pdf, plus image export.
# Needs a local soffice binary (env SOFFICE_CMD > bundled/adjacent
# tesseract-style locations > PATH). 501 with a clear hint when absent.
#
# Subprocess posture (mirrors /ocr): FIXED argument shape — the model may
# choose the format pair, never a raw command. Conversions run in a temp
# outdir (soffice derives its own output name), the product is verified by
# magic bytes (soffice rc alone lies on some builds), then moved atomically
# to the confirmed out path. LibreOffice cannot run two instances off one
# user profile at once — a module lock serializes conversions.

_CONV_MATRIX = {
    # from-ext -> {to-ext: (soffice family, filter-or-None, magic)}
    ".doc": {".docx": ("writer", None, b"PK\x03\x04"),
             ".pdf": ("writer", None, b"%PDF-")},
    ".docx": {".pdf": ("writer", None, b"%PDF-"),
              ".doc": ("writer", None, b"\xd0\xcf\x11\xe0"),
              ".txt": ("writer", "Text (encoded):UTF8", None),
              ".png": ("writer", "writer_png_Export", b"\x89PNG"),
              ".html": ("writer", "html", None)},
    ".rtf": {".docx": ("writer", None, b"PK\x03\x04"),
             ".pdf": ("writer", None, b"%PDF-")},
    ".odt": {".docx": ("writer", None, b"PK\x03\x04"),
             ".pdf": ("writer", None, b"%PDF-")},
    ".xls": {".xlsx": ("calc", None, b"PK\x03\x04"),
             ".pdf": ("calc", None, b"%PDF-"),
             ".csv": ("calc", "Text - txt - csv (StarCalc):44,34,76,1",
                      None)},
    ".xlsx": {".pdf": ("calc", None, b"%PDF-"),
              ".xls": ("calc", None, b"\xd0\xcf\x11\xe0"),
              ".csv": ("calc", "Text - txt - csv (StarCalc):44,34,76,1",
                       None)},
    ".ods": {".xlsx": ("calc", None, b"PK\x03\x04"),
             ".pdf": ("calc", None, b"%PDF-")},
    ".ppt": {".pptx": ("impress", None, b"PK\x03\x04"),
             ".pdf": ("impress", None, b"%PDF-")},
    ".pptx": {".pdf": ("impress", None, b"%PDF-"),
              ".png": ("impress", "impress_png_Export", b"\x89PNG")},
    ".odp": {".pptx": ("impress", None, b"PK\x03\x04"),
             ".pdf": ("impress", None, b"%PDF-")},
}
_CONV_LOCK = threading.Lock()   # one soffice profile at a time
_CONV_TIMEOUT = 120             # seconds per conversion


def _soffice_bin() -> str | None:
    """SOFFICE_CMD env > soffice next to the app/binary dir > PATH."""
    cands = []
    if os.environ.get("SOFFICE_CMD"):
        cands.append(os.environ["SOFFICE_CMD"])
    here = Path(sys.executable).parent if getattr(sys, "frozen", False) \
        else Path(__file__).resolve().parent
    cands += [here / "soffice",
              here / "libreoffice" / "program" / "soffice",
              Path("/usr/bin/soffice"), Path("/usr/local/bin/soffice"),
              Path("/opt/libreoffice/program/soffice")]
    for c in cands:
        if c and Path(c).exists():
            return str(c)
    return shutil.which("soffice") or shutil.which("libreoffice")


def convert_file(src: Path, out: Path, out_root: Path, to_ext: str) -> tuple[int, dict]:
    soffice = _soffice_bin()
    if not soffice:
        return 501, {"error": "convert needs LibreOffice (soffice). Install it "
                              "or set SOFFICE_CMD to the binary.",
                     "hint": "windows/mac installers usually bundle it; on "
                             "linux: libreoffice package or a portable "
                             "AppImage extracted next to the bridge"}
    src_ext = src.suffix.lower()
    family, filt, magic = _CONV_MATRIX.get(src_ext, {}).get(to_ext,
                                                            (None, None, None))
    if family is None:
        pairs = ", ".join(sorted(k for k in _CONV_MATRIX.get(src_ext, {})))
        return 400, {"error": f"cannot convert {src_ext} -> {to_ext}",
                     "supported_targets": pairs or "none"}
    with tempfile.TemporaryDirectory(prefix="fb-conv-") as td:
        cmd = [soffice, "--headless", "--norestore", "--nolockcheck",
               "--convert-to", to_ext.lstrip(".") if not filt else
               f"{to_ext.lstrip('.')}:{filt}",
               "--outdir", td, str(src)]
        with _CONV_LOCK:  # soffice user-profile lock: serialize
            try:
                proc = subprocess.run(cmd, capture_output=True, timeout=_CONV_TIMEOUT)
            except subprocess.TimeoutExpired:
                return 504, {"error": f"conversion timed out after "
                                      f"{_CONV_TIMEOUT}s"}
        produced = list(Path(td).glob("*"))
        produced = [f for f in produced if f.is_file() and not f.name.startswith(".")]
        if not produced:
            tail = (proc.stderr or proc.stdout or b"").decode(errors="replace")[-300:]
            return 500, {"error": "soffice produced no output", "soffice": tail}
        blob = produced[0].read_bytes()
        if magic and not blob.startswith(magic):
            return 500, {"error": f"conversion output failed magic check "
                                  f"(expected {magic[:4]!r})"}
        if len(blob) > MAX_BINARY:
            return 413, {"error": f"converted file too large "
                                  f"(> {MAX_BINARY} bytes)"}
        okr, err = rate_check(len(blob))
        if not okr:
            return 429, {"error": err, "rate_limited": True}
        snap = snapshot_before_write(out_root, out) if out.exists() else None
        atomic_write_bytes(out, blob)
        return 200, {"ok": True, "written": str(out), "bytes": len(blob),
                     "from": src.name, "to_ext": to_ext,
                     "soffice": "yes", "snapshot": snap}


# ------------------------------------------------ search + edit (P1)

def _search(root: Path, cfg: dict, q: dict):
    """Cross-file grep with context lines, glob filter, exclusions,
    case-insensitivity (param design from openapi-servers /search_content)."""
    import fnmatch
    query = q.get("q", "")
    if not query:
        return {"error": "need q= search term"}
    glob_pat = q.get("glob") or "*"
    case_sensitive = q.get("case", "insensitive") == "sensitive"
    needle = query if case_sensitive else query.lower()
    ctx = min(int(q.get("context", "1") or 1), 5)
    max_matches = min(int(q.get("max", "50") or 50), 200)
    excl = [x for x in (q.get("exclude") or "").split(",") if x.strip()]
    pats = _all_ignore(cfg) + excl
    matches, scanned = [], 0
    for f in _safe_walk(root)[0]:
        if len(matches) >= max_matches:
            break
        if not f.is_file():
            continue
        rel = f.relative_to(root).as_posix()
        if not fnmatch.fnmatch(rel, glob_pat):
            continue
        if _ignore_match(rel, False, pats):
            continue
        if f.suffix.lower() not in TEXT_EXTS and f.name.lower() not in KNOWN_BASENAMES:
            continue  # text files only (binary grep = noise)
        try:
            lines = f.read_text(encoding="utf-8", errors="replace").splitlines()
        except OSError:
            continue
        scanned += 1
        for i, line in enumerate(lines):
            hay = line if case_sensitive else line.lower()
            if needle in hay:
                lo, hi = max(0, i - ctx), min(len(lines), i + ctx + 1)
                matches.append({
                    "path": rel, "line": i + 1,
                    "context": [f"{j+1}: {lines[j][:MAX_LINE_CHARS]}"
                                for j in range(lo, hi)],
                })
                if len(matches) >= max_matches:
                    break
    return {"query": query, "scanned_files": scanned, "matches": matches,
            "truncated": len(matches) >= max_matches}


def _edit_file(root: Path, cfg: dict, body: dict, confirm_func):
    """Surgical replacements with dry-run unified diff (pattern:
    openapi-servers /edit_file). Applies via the standard write path
    (snapshot + confirmation) — never a raw overwrite."""
    import difflib
    rel = body.get("path", "")
    if not rel:
        return 400, {"error": "need path"}
    p, _root, _cfg = resolve_guarded(unquote(rel), for_write=True)
    if not p.is_file():
        return 404, {"error": f"no such file: {rel}"}
    if p.suffix.lower() not in TEXT_EXTS and p.name.lower() not in KNOWN_BASENAMES:
        return 415, {"error": "/edit is for text files (see /read whitelist)"}
    if _readonly_for(cfg):
        return 403, {"error": "read-only mode is active"}
    edits = body.get("edits", [])
    if not isinstance(edits, list) or not edits:
        return 400, {"error": "need edits: [{old_text, new_text}, ...]"}
    original = p.read_text(encoding="utf-8", errors="replace")
    modified = original
    for e in edits:
        old, new = e.get("old_text", ""), e.get("new_text", "")
        count = e.get("count", 1)
        if old not in modified:
            return 400, {"error": f"old_text not found: {old[:60]!r}"}
        modified = modified.replace(old, new, int(count) if count else -1)
    if body.get("dry_run"):
        diff = "".join(difflib.unified_diff(
            original.splitlines(keepends=True),
            modified.splitlines(keepends=True),
            fromfile=f"a/{rel}", tofile=f"b/{rel}"))
        return 200, {"dry_run": True, "diff": diff,
                     "changed": modified != original}
    # real write: same guarded path as /write (confirmation + snapshot)
    code, resp = confirm_func({"op": "edit", "path": rel,
                               "bytes": len(modified)})
    if code:
        return code, resp
    return 200, {"ok": True, "path": rel, "edited": modified != original}


def _html_text(path: Path, q: dict):
    from html.parser import HTMLParser

    class Extractor(HTMLParser):
        SKIP = {"script", "style", "noscript", "head", "meta", "title"}
        BLOCK = {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5",
                 "h6", "section", "article", "table", "ul", "ol"}

        def __init__(self):
            super().__init__()
            self.parts = []
            self.skip_depth = 0

        def handle_starttag(self, tag, attrs):
            if tag in self.SKIP:
                self.skip_depth += 1
            elif tag in self.BLOCK and self.parts and not self.parts[-1].endswith("\n"):
                self.parts.append("\n")

        def handle_endtag(self, tag):
            if tag in self.SKIP and self.skip_depth:
                self.skip_depth -= 1
            elif tag in self.BLOCK:
                self.parts.append("\n")

        def handle_data(self, data):
            if not self.skip_depth and data.strip():
                self.parts.append(data)

    ex = Extractor()
    ex.feed(path.read_text(encoding="utf-8", errors="replace"))
    text = "".join(ex.parts)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    max_chars = int(q.get("max_chars", "60000") or 60000)
    return {"path": q.get("path"), "chars": len(text),
            "truncated": len(text) > max_chars, "text": text[:max_chars]}


def _csv_head_stats(path: Path, q: dict, want_stats: bool):
    import csv
    import io as _io
    nrows = min(int(q.get("rows", "20") or 20), 200)
    delim = q.get("delim") or ","
    if delim in ("tab", "\\t", "t"):
        delim = "\t"
    rows, widths = [], set()
    with open(path, newline="", encoding="utf-8", errors="replace") as fh:
        reader = csv.reader(fh, delimiter=delim)
        for i, row in enumerate(reader):
            widths.add(len(row))
            if i < nrows:
                rows.append(row)
    out = {"path": q.get("path"), "delimiter": delim,
           "ragged": len(widths) > 1, "widths": sorted(widths)[:10],
           "head": rows}
    if want_stats:
        with open(path, newline="", encoding="utf-8", errors="replace") as fh:
            all_rows = list(csv.reader(fh, delimiter=delim))
        total = len(all_rows)
        ncols = max(widths) if widths else 0
        cols = []
        for c in range(min(ncols, 100)):
            vals = [r[c] for r in all_rows if len(r) > c and r[c] != ""]
            nums = []
            for v in vals:
                try:
                    nums.append(float(v))
                except ValueError:
                    pass
            col = {"col": c, "non_empty": len(vals)}
            if nums and len(nums) >= max(1, len(vals) // 2):
                col["type"] = "numeric"
                col["min"] = min(nums)
                col["max"] = max(nums)
            else:
                col["type"] = "text"
                col["samples"] = vals[:3]
            cols.append(col)
        out.update({"row_count": total, "column_count": ncols, "columns": cols})
    return out




# ------------------------------------------- trash + write guards (P0b)

TRASH_DIR = STATE_DIR / "trash"
TRASH_PURGE_DAYS = 30

# ---- rate circuit breaker: max writes per rolling window ----
_WRITE_LOG_LOCK = threading.Lock()
_WRITE_LOG: list = []          # (monotonic_ts, nbytes)
RATE_MAX_WRITES = 20          # per 60 s window (env-tunable)
RATE_MAX_BYTES = 50 * 1024 * 1024


def _rate_limits():
    mw = os.environ.get("FILE_BRIDGE_MAX_WRITES")
    mb = os.environ.get("FILE_BRIDGE_MAX_WRITE_MB")
    return (int(mw) if mw and mw.isdigit() else RATE_MAX_WRITES,
            int(mb) * 1024 * 1024 if mb and mb.isdigit() else RATE_MAX_BYTES)


def rate_check(nbytes: int) -> tuple[bool, str]:
    """Record a prospective write; False when the rolling window is over
    budget. A runaway model hits the brake, not the folder."""
    now = time.monotonic()
    max_w, max_b = _rate_limits()
    with _WRITE_LOG_LOCK:
        global _WRITE_LOG
        _WRITE_LOG = [t for t in _WRITE_LOG if now - t[0] < 60]
        w_count = sum(1 for t in _WRITE_LOG if t[1] >= 0)
        w_bytes = sum(max(t[1], 0) for t in _WRITE_LOG)
        if w_count + 1 > max_w or w_bytes + nbytes > max_b:
            return False, (f"write-rate circuit breaker: {w_count} writes / "
                           f"{w_bytes // 1024} KiB in the last 60 s (limits "
                           f"{max_w} / {max_b // 1024 // 1024} MiB). STOP and ask "
                           f"the user to confirm the mass edit before continuing.")
        _WRITE_LOG.append((now, nbytes))
    return True, ""


def _readonly_for(cfg: dict) -> bool:
    if _is_readonly():
        return True
    return bool(cfg.get("readonly"))


def trash_store_for(root: Path) -> Path:
    d = TRASH_DIR / _root_key(root)
    d.mkdir(parents=True, exist_ok=True)
    _restrict_to_user(d, is_dir=True)
    return d


def trash_move(root: Path, target: Path) -> dict:
    """Move target into the trash store preserving the tree. Cross-device
    fallback: verified copy-then-unlink (user decision note in roadmap)."""
    rel = target.relative_to(root).as_posix()
    ts = time.strftime("%Y%m%d-%H%M%S") + "-" + _secrets.token_hex(2)
    dest = trash_store_for(root) / ts / rel.lstrip("/")
    dest.parent.mkdir(parents=True, exist_ok=True)
    try:
        os.rename(target, dest)            # same device: atomic
    except OSError:
        if target.is_dir():
            shutil.copytree(target, dest)
            n_src = sum(1 for _ in target.rglob("*"))
            n_dst = sum(1 for _ in dest.rglob("*"))
            if n_src != n_dst:
                shutil.rmtree(dest, ignore_errors=True)
                raise RuntimeError("trash copy verification failed — aborting, "
                                   "original untouched")
            shutil.rmtree(target)
        else:
            shutil.copy2(target, dest)
            if dest.stat().st_size != target.stat().st_size:
                dest.unlink(missing_ok=True)
                raise RuntimeError("trash copy verification failed — aborting, "
                                   "original untouched")
            target.unlink()
    mf = trash_store_for(root) / "manifest.jsonl"
    with open(mf, "a", encoding="utf-8") as fh:
        fh.write(json.dumps({"ts": ts, "path": rel,
                             "type": "dir" if dest.is_dir() else "file",
                             "size": _tree_size(dest)}) + "\n")
    return {"ts": ts, "path": rel, "trashed_to": str(dest)}


def _tree_size(p: Path) -> int:
    try:
        return sum(f.stat().st_size for f in p.rglob("*") if f.is_file()) if p.is_dir() \
            else p.stat().st_size
    except OSError:
        return 0


def trash_purge(root: Path | None = None):
    """Drop entries older than TRASH_PURGE_DAYS (opportunistic, called on
    trash ops; manual purge = settings page only)."""
    import datetime
    cutoff = time.time() - TRASH_PURGE_DAYS * 86400
    bases = [trash_store_for(root)] if root else \
            [d for d in TRASH_DIR.glob("*/") if d.is_dir()]
    for base in bases:
        for snap in base.glob("*/"):
            try:
                if snap.stat().st_mtime < cutoff:
                    shutil.rmtree(snap, ignore_errors=True)
            except OSError:
                pass


def trash_list(root: Path, rel_filter: str = "") -> list:
    """METADATA ONLY — same never-re-enter-model-reach rule as versions."""
    out = []
    mf = trash_store_for(root) / "manifest.jsonl"
    try:
        for line in open(mf, encoding="utf-8").read().splitlines():
            try:
                e = json.loads(line)
            except Exception:
                continue
            if rel_filter and rel_filter not in e.get("path", ""):
                continue
            e["restore"] = "POST /trash/restore {path: %r, ts: %r}" % (e["path"], e["ts"])
            out.append(e)
    except FileNotFoundError:
        pass
    out.reverse()
    return out[:200]


def trash_restore(root: Path, rel: str, ts: str) -> tuple[bool, str]:
    src = trash_store_for(root) / ts / rel.lstrip("/")
    if not src.exists():
        cands = list((trash_store_for(root) / ts).rglob(Path(rel).name))
        if len(cands) != 1:
            return False, f"no unique trash entry for {rel} @ {ts}"
        src = cands[0]
    target = (root / rel).resolve()
    if os.path.commonpath([str(root), str(target)]) != str(root):
        return False, "restore target escapes root"
    if target.exists():
        return False, ("target exists — trash/restore refuses to overwrite; "
                       "delete or rename it first")
    target.parent.mkdir(parents=True, exist_ok=True)
    try:
        os.rename(src, target)
    except OSError:
        if src.is_dir():
            shutil.copytree(src, target)
            shutil.rmtree(src)
        else:
            shutil.copy2(src, target)
            src.unlink()
    return True, str(target)



# ------------------------------------------------------- multi-root (P0b)

class ExcludedPath(Exception):
    """Path matches an ignore pattern → endpoints return 404-with-hint
    (roadmap: 'excluded by settings'), writes are refused outright."""
    def __init__(self, pattern, rel):
        self.pattern, self.rel = pattern, rel
        if pattern in DEFAULT_IGNORE:
            msg = (f"'{rel}' is OS metadata junk (e.g. .DS_Store) — always "
                   f"excluded by default, nothing to configure")
        else:
            msg = (f"'{rel}' is excluded by ignore settings (pattern: {pattern}); "
                   f"tell the user — ignore patterns are editable in the File "
                   f"Bridge settings page (the local picker)")
        super().__init__(msg)


def roots_config() -> list:
    """All configured roots. Legacy single-'root' state migrates on the fly."""
    st = _state_load()
    if isinstance(st.get("roots"), list) and st["roots"]:
        return st["roots"]
    r = st.get("root")
    if r:
        return [{"id": "main", "path": r, "alias": "main",
                 "enabled": True, "ignore": []}]
    return []


def enabled_roots() -> list:
    out = []
    for r in roots_config():
        if not r.get("enabled", True):
            continue
        p = Path(r.get("path", "")).expanduser()
        if not p.is_dir():
            continue  # fail-closed: unlocatable roots simply don't resolve
        out.append(r)
    return out


def _sanitize_root_entry(e: dict, i: int) -> dict:
    rid = str(e.get("id") or f"root{i}").strip().lower()
    rid = re.sub(r"[^a-z0-9_-]", "-", rid)[:32] or f"root{i}"
    path = str(e.get("path", "")).strip()
    return {"id": rid, "path": path,
            "alias": str(e.get("alias") or rid)[:64],
            "enabled": bool(e.get("enabled", True)),
            "ignore": [str(x) for x in (e.get("ignore") or [])][:200]}


def set_roots(entries: list) -> tuple[bool, str]:
    """Validate + persist the root list. Rejects: empty paths, dup ids,
    roots inside the bridge state dir (self-protection)."""
    if not isinstance(entries, list) or not entries:
        return False, "roots must be a non-empty list"
    cleaned, seen = [], set()
    for i, e in enumerate(entries):
        c = _sanitize_root_entry(e, i)
        p = Path(c["path"]).expanduser()
        if not c["path"]:
            return False, f"root {i}: empty path"
        if not p.is_dir():
            return False, f"root '{c['id']}': not a folder: {p}"
        rp = p.resolve()
        try:
            if os.path.commonpath([str(STATE_DIR), str(rp)]) == str(STATE_DIR):
                return False, "bridge state dir cannot live inside a shared root"
            if os.path.commonpath([str(rp), str(STATE_DIR)]) == str(rp):
                return False, "shared root cannot contain the bridge state dir"
        except ValueError:
            pass  # different drives (Windows)
        if c["id"] in seen:
            return False, f"duplicate root id: {c['id']}"
        seen.add(c["id"])
        c["path"] = str(rp)
        cleaned.append(c)
    _state_update(roots=cleaned, root=cleaned[0]["path"])  # root=key compat
    return True, ""


_ROOT_ID_RE = re.compile(r"^([a-z0-9_-]{1,32})(?:/|$)")


def resolve_any(rel: str):
    """Resolve a request path to (abs, root_path, root_cfg).

    Addressing: '<root-id>/sub/path' picks that root; anything else resolves
    under the FIRST enabled root (default). Unknown leading segment = ordinary
    subdirectory of the default root."""
    ros = enabled_roots()
    if not ros:
        raise PermissionError("no shared folder configured — open the File Bridge "
                              "settings page")
    rel = (rel or "").strip()
    if rel.startswith("/") or (len(rel) > 1 and rel[1] == ":"):
        raise PermissionError("absolute paths are not allowed — address files "
                              "relative to a shared root")
    rel = rel.lstrip("/").replace("\\", "/")
    cfg = ros[0]
    m = _ROOT_ID_RE.match(rel)
    if m and m.group(1) in {r["id"] for r in ros}:
        cfg = next(r for r in ros if r["id"] == m.group(1))
        rel = rel[m.end():]
    root = Path(cfg["path"]).resolve()
    p = (root / rel).resolve()
    try:
        if os.path.commonpath([str(root), str(p)]) != str(root):
            raise PermissionError("path escapes shared root")
    except ValueError:
        # commonpath ValueError = different drives (Windows UNC 'C:\' vs
        # '\\server\share') — never inside the root, so it's an escape
        raise PermissionError("path escapes shared root")
    # sensitive-name floor (P2): .env / id_rsa / *.pem / credentials… never
    # readable or writable via the API, regardless of ignore lists
    if sensitive_name(p.name):
        raise PermissionError(
            f"'{p.name}' looks like a credential/secret file — the bridge "
            f"refuses to serve it. Ask the user to handle it manually.")
    # symlink guard: no component of p may be a symlink (a link to a dir
    # INSIDE the same root still resolves inside, but mixed trust — refuse
    # uniformly; the model can address the real path instead)
    cur = root
    for comp in rel.split("/"):
        if not comp:
            continue
        cur = cur / comp
        if cur.is_symlink():
            raise PermissionError(f"symlink in path is not allowed: {comp}")
    return p, root, cfg


def _global_ignore() -> list:
    return [str(x) for x in (_state_load().get("ignore_global") or [])]


# Built-in ignore floor: Finder/Explorer metadata junk that real shared
# folders accumulate on macOS/Windows (.DS_Store in EVERY folder Finder
# touches, AppleDouble ._ files on non-HFS volumes, Thumbs.db/desktop.ini
# from Explorer). Always active for every root, without configuration —
# never listed, zipped, extracted, read or written. Users who truly need
# one of these names should pick a different name for their real file.
DEFAULT_IGNORE = [".DS_Store", "._*", "Thumbs.db", "desktop.ini"]


def _all_ignore(cfg: dict) -> list:
    """Full pattern set enforced by the bridge: per-root ignore +
    built-in junk floor + user-global ignore. Every consumer of ignore
    patterns composes through here so the floor cannot be bypassed."""
    return list(cfg.get("ignore", [])) + DEFAULT_IGNORE + _global_ignore()


def _ignore_match(rel: str, is_dir: bool, patterns: list):
    """gitignore-style subset: 'dir/' dir-only, '/x' anchored to root,
    '*' within/across segments, '#' comments. A path is excluded if it OR
    ANY ancestor directory matches. A pattern containing '/' is matched
    against the joined path; a bare name matches that name at ANY depth
    (gitignore semantics — .DS_Store or .git/ must prune nested copies
    too, not just top-level ones)."""
    import fnmatch
    parts = [x for x in rel.split("/") if x]
    cands = ["/".join(parts[:i + 1]) for i in range(len(parts))]
    for pat in patterns:
        p = pat.strip()
        if not p or p.startswith("#"):
            continue
        dir_only = p.endswith("/")
        if dir_only:
            p = p[:-1]
        anchored = p.startswith("/")
        if anchored:
            p = p[1:]
        if not p:
            continue
        if "/" in p:
            tests = cands          # path-shaped: match joined path prefixes
        elif anchored:
            tests = parts[:1]      # '/x': root segment only
        else:
            tests = parts          # bare name: any segment, any depth
        for i, cand in enumerate(tests):
            if i == len(tests) - 1 and dir_only and not is_dir:
                continue
            if fnmatch.fnmatch(cand, p):
                return pat
    return None


def resolve_guarded(rel: str, *, for_write: bool = False):
    """resolve_any + self-protection floor + ignore enforcement.
    Used by EVERY file endpoint (roadmap: enforcement lives in the BRIDGE)."""
    p, root, cfg = resolve_any(rel)
    # --- self-protection floor (P0b): bridge-owned files are never
    # accessible via the file API, in any mode, even when a root overlaps —
    # prevents 'one approved write quietly widens future permissions'.
    try:
        if os.path.commonpath([str(STATE_DIR), str(p)]) == str(STATE_DIR):
            raise PermissionError("bridge state/storage is never accessible via the file API")
    except ValueError:
        pass
    rel_in_root = p.relative_to(root).as_posix()
    if rel_in_root != ".":
        pats = _all_ignore(cfg)
        hit = _ignore_match(rel_in_root, p.is_dir(), pats)
        if hit:
            if for_write:
                raise ExcludedPath(hit, rel_in_root)
            raise ExcludedPath(hit, rel_in_root)
    return p, root, cfg


def _legacy_root() -> Path | None:
    ros = enabled_roots()
    return Path(ros[0]["path"]) if ros else None


# ------------------------------------------------------------------ server

class Handler(http.server.BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):
        sys.stderr.write(f"[bridge] {self.address_string()} {fmt % args}\n")

    # ---- CORS / auth plumbing ----

    def _add_matching_cors(self):
        """Emit CORS headers ONLY when the request origin matches the lock
        (or no lock is configured — picker/setup phase). v1 echoed `*`.

        Opaque origins (literal `Origin: null` — Open WebUI's sandboxed
        Pyodide iframe) are granted CORS only while the TOKEN tier is
        active: a preflight grants nothing by itself and every real file
        request still has to carry the token (check_request), so responses
        become readable only to callers who know it. Origin-only mode keeps
        refusing null (an unmatched sandbox must read nothing)."""
        raw_origin = (self.headers.get("Origin") or "").strip()
        origin = _request_origin(self.headers)      # null/absent → None
        allowed = get_allowed_origin()
        # Opaque origins are granted CORS whenever the token tier is active.
        # Enumerating null-origin requests: public endpoints, token-valid
        # calls, and 401 token failures — every response body reachable this
        # way is either public or an error message with no secrets, and a
        # READABLE 401 ("missing or invalid bridge token") lets the model
        # self-correct in one retry instead of hitting an opaque CORS block.
        # Origin-only mode keeps refusing null (an unmatched sandbox must
        # read nothing).
        opaque_ok = (raw_origin.lower() == "null"
                     and get_configured_token() is not None)
        if allowed is None or (origin and origin == allowed) or opaque_ok:
            echo = "null" if opaque_ok else (origin or "*")
            self.send_header("Access-Control-Allow-Origin", echo)
            self.send_header("Access-Control-Allow-Private-Network", "true")
            self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
            self.send_header("Access-Control-Allow-Headers",
                             "Content-Type, X-Bridge-Token, Authorization")
        # mismatched/foreign origins get NO CORS headers — the browser
        # refuses to expose the response even before our 403 body.

    def _json(self, code, obj, cors=True):
        body = json.dumps(obj).encode()
        if getattr(self, "_audit_ep", None):
            audit_log(self._audit_ep, method=self.command,
                      path=getattr(self, "_audit_path", None),
                      args=getattr(self, "_audit_args", None), status=code,
                      size=len(body), risk=self._risk())
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        if cors:
            self._add_matching_cors()
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _audit(self, ep, path=None, args=None):
        """Mark this request for audit logging (idempotent per response)."""
        self._audit_ep = ep
        self._audit_path = path if path is None else str(path)
        self._audit_args = args

    def _risk(self):
        """Declared risk class for this request's endpoint (P3, openworker
        risk.py pattern). Unknown paths report 'read' — the security gate
        has already run by the time this is consulted."""
        return ENDPOINT_RISK.get(urlparse(self.path).path, RiskClass.READ)

    def do_OPTIONS(self):
        self.send_response(204)
        self._add_matching_cors()
        self.send_header("Access-Control-Max-Age", "86400")
        self.end_headers()

    # ---- GET ----

    def do_GET(self):
        root = load_root()
        u = urlparse(self.path)
        q = {k: v[0] for k, v in parse_qs(u.query).items()}

        if u.path == "/health":
            ros = [{"id": r["id"], "path": r["path"], "alias": r.get("alias")}
                   for r in enabled_roots()]
            info = {"ok": bool(ros), "roots": ros, "version": VERSION} if ros else \
                   {"ok": False, "hint": "no folder chosen yet",
                    "roots": [], "version": VERSION}
            if ros:
                info["root"] = ros[0]["path"]
            info["addons"] = {"pdf": HAVE_PYMUPDF, "ocr": bool(TESSERACT_BIN)}
            info["ocr_lang"] = _get_ocr_lang()
            info["ocr_langs_available"] = _ocr_langs_available()
            info["wheels"] = len(list(WHEELS_DIR.glob("*.whl"))) if WHEELS_DIR.is_dir() else 0
            mode = security_mode()
            info["security"] = mode
            info["locked"] = mode != "UNLOCKED"
            return self._json(200, info)

        if u.path == "/version":
            # token-free like /health: just versions, no fs info. Skill 2.4+
            # reads the version from /health instead (one fewer round trip);
            # /version stays for older skills and manual debugging.
            return self._json(200, {"bridge": VERSION, "skill": SKILL_VERSION,
                                    "skill_expected": SKILL_VERSION,
                                    "note": "bridge and skill versions should "
                                            "match; re-run scripts/setup_owui.py "
                                            "to sync the skill"})

        if u.path == "/api/preview":
            # LOCAL picker preview ("What the AI can see") — loopback-only
            # and token-free: the owner UI never sees the token (it is only
            # ever stored hashed), so the preview cannot send it. GET + no
            # CORS headers ⇒ cross-origin pages can trigger it but never
            # read the response; locals can read the folder from disk
            # anyway. Same data as /directory_tree for the first root.
            if not self._is_loopback():
                return self._json(403, {"error": "picker API is local only"}, cors=False)
            self._audit("/api/preview")
            ros = enabled_roots()
            if not ros:
                return self._json(200, {"ok": False,
                                        "hint": "no folder chosen yet"}, cors=False)
            return self._json(200, directory_tree(Path(ros[0]["path"]), ros[0],
                                                  {"max_entries": "500",
                                                   "max_depth": "6"}), cors=False)

        if u.path == "/state":
            return self._json(200, {
                "root": str(root) if root else None,
                "roots": roots_config(), "port": PORT,
                "ocr_lang": _get_ocr_lang(),
                "allowed_origin": get_allowed_origin(),
                "security": security_mode(),
                "readonly": _is_readonly(),
                "allow_reveal": bool(_state_load().get("allow_reveal")),
                "ignore_global": _global_ignore(),
                "rate_limits": _rate_limits(),
                "endpoint_risk": {k: ENDPOINT_RISK[k] for k in
                                  sorted(ENDPOINT_RISK)}})

        if u.path == "/ocr/config":
            # UNGATED (loopback-bound server): lang list + engine version are
            # the same disclosure class as /health (which already reports
            # ocr_langs_available), and the local picker needs this on FIRST
            # RUN — before the origin lock exists, while file endpoints stay
            # hard-disabled (the settings page is the only thing reachable).
            return self._json(200, {
                "lang": _get_ocr_lang(),
                "available": _ocr_langs_available(),
                "engine": TESSERACT_VER or "tesseract",
                "user_dir": str(USER_TESSDATA_DIR),
                "hint": "drop extra .traineddata files (tessdata_fast) into "
                        "user_dir and restart the bridge to add languages",
            })

        # ---- wheel hosting (token-free by design: static public wheels) ----
        if u.path == "/wheels":
            if WHEELS_DIR.is_dir():
                names = sorted(f.name for f in WHEELS_DIR.glob("*.whl"))
            else:
                names = []
            return self._json(200, {"wheels": names,
                                    "urls": [f"http://127.0.0.1:{PORT}/wheels/{n}" for n in names]})

        if u.path.startswith("/wheels/"):
            name = os.path.basename(unquote(u.path[len("/wheels/"):]))
            wf = (WHEELS_DIR / name).resolve()
            if os.path.commonpath([str(WHEELS_DIR), str(wf)]) != str(WHEELS_DIR) or not wf.is_file():
                return self._json(404, {"error": f"no such wheel: {name}"})
            data = wf.read_bytes()
            self.send_response(200)
            self.send_header("Content-Type", "application/octet-stream")
            self._add_matching_cors()
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)
            return

        # ---- security gate: everything below serves files ----
        ok, status, reason = check_request(self.headers)
        if not ok:
            return self._json(status, {"error": reason})

        if root is None:
            return self._json(409, {"error": "No folder chosen. Open http://127.0.0.1:%d to pick one." % PORT})

        self._audit(u.path, path=unquote(q.get("path", "")) or None,
                    args={k: q[k] for k in q if k != "path"} or None)

        try:
            if u.path == "/list":
                p, root, cfg = resolve_guarded(unquote(q.get("path", ".")))
                if not p.is_dir():
                    return self._json(404, {"error": f"not a directory: {q.get('path')}"})
                entries = []
                pats = _all_ignore(cfg)
                walked, truncated = _safe_walk(p)
                for f in walked:
                    rel = f.relative_to(root).as_posix()
                    if _ignore_match(rel, f.is_dir(), pats):
                        continue
                    try:
                        st = f.stat()
                        entries.append({"path": rel, "type": "dir" if f.is_dir() else "file",
                                        "size": None if f.is_dir() else st.st_size,
                                        "mtime": int(st.st_mtime)})
                    except OSError:
                        continue
                    if len(entries) >= MAX_LIST:
                        truncated = True
                        break
                return self._json(200, {"root": str(root), "root_id": cfg.get("id"),
                                        "entries": entries, "truncated": truncated,
                                        **({"hint": "partial listing — folder exceeded the "
                                                    "entry/time cap; narrow with path="
                                            } if truncated else {})})

            if u.path == "/read":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                ext = p.suffix.lower()
                base = p.name.lower()
                if ext not in TEXT_EXTS and base not in KNOWN_BASENAMES and base != ".gitignore":
                    # fail-closed: unknown extension → refuse, offer /peek
                    return self._json(415, {
                        "error": f"/read is text-only; '.{ext.lstrip('.')}' is not on the "
                                 f"text whitelist (fail-closed)",
                        "hint": f"call /peek?path=…&bytes=512 to identify the file, or "
                                f"/read_b64 for raw bytes"})
                head = b""
                with open(p, "rb") as fh:
                    head = fh.read(16)
                kind = sniff_kind(head)
                if kind:
                    return self._json(415, {
                        "error": f"/read is text-only; this file sniffs as {kind}",
                        "hint": _routing_hint(kind, ext)})
                start_line = int(q.get("start_line", "1") or 1)
                max_lines = int(q.get("max_lines", str(DEFAULT_MAX_LINES)))
                out = _windowed_read(p, start_line, max_lines)
                out["path"] = q.get("path")
                return self._json(200, out)

            if u.path == "/peek":
                # identify a file cheaply: first bytes as preview + sniffing
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                nbytes = max(16, min(int(q.get("bytes", "512") or 512), 4096))
                with open(p, "rb") as fh:
                    head = fh.read(nbytes)
                kind = sniff_kind(head)
                ext = p.suffix.lower()
                guessed = kind or ("text" if ext in TEXT_EXTS or p.name.lower() in KNOWN_BASENAMES
                                   else "unknown")
                preview = head.decode("utf-8", errors="replace")
                preview = "".join(ch if ch.isprintable() or ch in "\n\r\t" else "·"
                                  for ch in preview)
                return self._json(200, {
                    "path": q.get("path"), "size": p.stat().st_size,
                    "ext": ext or None, "kind": guessed,
                    "printable_ratio": (sum(1 for b in head if 32 <= b < 127 or b in (9, 10, 13))
                                        / max(len(head), 1)),
                    "preview": preview,
                    "hint": _routing_hint(kind, ext) if kind else
                            ("looks like text — /read it" if guessed == "text" else
                             "unknown format — /read_b64 if you need the bytes"),
                })

            if u.path == "/read_b64":
                # binary-safe read: returns base64. For Office files, images,
                # PDFs. Capped at MAX_BINARY bytes.
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                raw = p.read_bytes()
                if len(raw) > MAX_BINARY:
                    return self._json(413, {"error": f"file too large: {len(raw)} > {MAX_BINARY} bytes"})
                return self._json(200, {
                    "path": q.get("path"),
                    "size": len(raw),
                    "b64": base64.b64encode(raw).decode("ascii"),
                })

            if u.path == "/stat":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.exists():
                    return self._json(404, {"error": f"not found: {q.get('path')}"})
                st = p.stat()
                ext = p.suffix.lower()
                kind = ("image" if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
                        else "pdf" if ext == ".pdf"
                        else "zip" if ext in {".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"}
                        else "legacy" if ext in {".doc", ".xls", ".ppt"}
                        else "text" if ext in {".txt", ".md", ".csv", ".json", ".xml", ".yml", ".yaml", ".log", ".py", ".js", ".html", ".css"}
                        else "other")
                return self._json(200, {
                    "path": q.get("path"), "size": st.st_size,
                    "mtime": int(st.st_mtime), "kind": kind, "ext": ext or None,
                })

            # ---------- optional PDF/OCR add-on endpoints ----------

            if u.path == "/pdf_text":
                # Extract embedded text layer from a PDF (pymupdf) — or, with
                # mode=images, rasterize pages for vision models (pypdfium2).
                if not HAVE_PYMUPDF:
                    return self._json(501, {"error": "PDF add-on not installed on this machine "
                                                     "(pip install pymupdf)"})
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                pages_param = q.get("pages", "")  # "1-3,5" optional, 1-indexed
                mode = q.get("mode", "text")
                if mode == "images":
                    if not HAVE_PDFIUM:
                        return self._json(501, {"error": "images mode needs pypdfium2 "
                                                         "(pip install pypdfium2)"})
                    try:
                        max_pages = max(1, min(int(q.get("max_pages", str(RASTER_MAX_PAGES))),
                                               RASTER_MAX_PAGES))
                    except ValueError:
                        max_pages = RASTER_MAX_PAGES
                    cache_params = {"mode": "images", "pages": pages_param,
                                    "max_pages": max_pages}
                    hit, cached = cache_get(p, "pdf_text", cache_params)
                    if hit:
                        return self._json(200, {**cached, "cached": True})
                    code, out = pdf_rasterize(p, pages_param, max_pages)
                    if code == 200:
                        cache_put(p, "pdf_text", cache_params, out)
                    return self._json(code, out)
                cache_params = {"pages": pages_param}
                hit, cached = cache_get(p, "pdf_text", cache_params)
                if hit:
                    return self._json(200, {**cached, "cached": True})
                doc = fitz.open(p)
                try:
                    total = doc.page_count
                    wanted = list(range(total))
                    if pages_param:
                        wanted = []
                        for part in pages_param.split(","):
                            part = part.strip()
                            if "-" in part:
                                a, b = part.split("-", 1)
                                wanted.extend(range(int(a) - 1, min(int(b), total)))
                            elif part:
                                wanted.append(int(part) - 1)
                        wanted = sorted(set(w for w in wanted if 0 <= w < total))
                    pages_out = []
                    char_budget = MAX_READ
                    for i in wanted:
                        txt = doc[i].get_text("text").strip()
                        pages_out.append({"page": i + 1, "text": txt})
                        char_budget -= len(txt)
                        if char_budget <= 0:
                            pages_out.append({"page": i + 1, "text": "…truncated (budget)"})
                            break
                    out = {"path": q.get("path"), "page_count": total,
                           "pages": pages_out}
                    cache_put(p, "pdf_text", cache_params, out)
                    return self._json(200, out)
                finally:
                    doc.close()

            if u.path == "/ocr":
                # OCR an image (png/jpg/bmp/webp/tiff) or a scanned PDF (per
                # page) using the tesseract binary. Language: request param
                # `lang` > saved setting (picker UI) > "eng".
                if not TESSERACT_BIN:
                    return self._json(501, {"error": "OCR unavailable: tesseract not "
                                                     "found. Install tesseract or set "
                                                     "TESSERACT_CMD."})
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                raw = q.get("lang") or _get_ocr_lang()
                # URL '+' decodes to space; treat space and ',' as separators
                parts = [p for p in re.split(r"[\s,+]+", raw) if p and re.fullmatch(r"[a-zA-Z_]{2,8}", p)]
                lang = "+".join(parts) if parts else "eng"
                ext = p.suffix.lower()
                max_pages = int(q.get("max_pages", "5"))
                dpi = q.get("dpi", "200")
                cache_params = {"lang": lang, "max_pages": max_pages, "dpi": dpi}
                hit, cached = cache_get(p, "ocr", cache_params)
                if hit:
                    return self._json(200, {**cached, "cached": True})
                results = []
                tmpdir = tempfile.mkdtemp(prefix="fb-ocr-")

                def run_tesseract(img_path):
                    cmd = [TESSERACT_BIN, img_path, "stdout", "-l", lang,
                           "--dpi", dpi]
                    env = dict(os.environ)
                    if TESSDATA_DIR:
                        env["TESSDATA_PREFIX"] = str(TESSDATA_DIR)
                    r = subprocess.run(cmd, capture_output=True, text=True,
                                       timeout=120, env=env)
                    if r.returncode != 0:
                        raise RuntimeError(r.stderr.strip()[:300])
                    return [l for l in r.stdout.splitlines() if l.strip()]

                try:
                    if ext in {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"}:
                        lines = run_tesseract(str(p))
                        results.append({"page": 1, "lines": lines})
                    elif ext == ".pdf":
                        if not HAVE_PYMUPDF:
                            return self._json(501, {"error": "PDF OCR needs the PDF add-on "
                                                             "(pip install pymupdf)"})
                        doc = fitz.open(p)
                        try:
                            for i in range(min(doc.page_count, max_pages)):
                                pix = doc[i].get_pixmap(dpi=int(dpi))
                                tf = os.path.join(tmpdir, f"p{i}.png")
                                pix.save(tf)
                                lines = run_tesseract(tf)
                                results.append({"page": i + 1, "lines": lines})
                        finally:
                            doc.close()
                    else:
                        return self._json(400, {"error": f"OCR supports images and PDF, not {ext}"})
                finally:
                    shutil.rmtree(tmpdir, ignore_errors=True)
                out = {"path": q.get("path"), "lang": lang,
                       "pages": results}
                cache_put(p, "ocr", cache_params, out)
                return self._json(200, out)

            if u.path == "/search":
                base = unquote(q.get("path", "."))
                p, root, cfg = resolve_guarded(base)
                if not p.is_dir():
                    return self._json(400, {"error": f"not a directory: {base}"})
                return self._json(200, _search(p, cfg, q))

            if u.path == "/directory_tree":
                base = unquote(q.get("path", "."))
                p, root, cfg = resolve_guarded(base)
                if not p.is_dir():
                    return self._json(400, {"error": f"not a directory: {base}"})
                return self._json(200, directory_tree(p, cfg, q))

            if u.path == "/html_text":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() not in {".html", ".htm"}:
                    return self._json(400, {"error": f"not html: {p.suffix}"})
                try:
                    return self._json(200, _html_text(p, q))
                except Exception as e:
                    return self._json(500, {"error": f"html parse failed: {e}"})

            if u.path == "/csv_head":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() not in {".csv", ".tsv"}:
                    return self._json(400, {"error": f"not a csv: {p.suffix}"})
                return self._json(200, _csv_head_stats(p, q, want_stats=False))

            if u.path == "/csv_stats":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() not in {".csv", ".tsv"}:
                    return self._json(400, {"error": f"not a csv: {p.suffix}"})
                return self._json(200, _csv_head_stats(p, q, want_stats=True))

            if u.path == "/xlsx_read":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() not in {".xlsx", ".xlsm"}:
                    return self._json(400, {"error": f"not an xlsx: {p.suffix}"})
                try:
                    out = _xlsx_read(p, q)
                except Exception as e:
                    return self._json(500, {"error": f"xlsx parse failed: {e}",
                                            "hint": "file may be legacy .xls — ask the "
                                                    "user to convert to .xlsx"})
                return self._json(200, out)

            if u.path == "/docx_read":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() != ".docx":
                    return self._json(400, {"error": f"not a docx: {p.suffix}"})
                try:
                    out = _docx_read(p, q)
                except Exception as e:
                    return self._json(500, {"error": f"docx parse failed: {e}"})
                return self._json(200, out)

            if u.path == "/pptx_read":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() != ".pptx":
                    return self._json(400, {"error": f"not a pptx: {p.suffix}"})
                try:
                    out = _pptx_read(p, q)
                except Exception as e:
                    return self._json(500, {"error": f"pptx parse failed: {e}"})
                return self._json(200, out)

            if u.path == "/eml_read":
                # RFC-822 email → headers + body + attachment metadata
                # (P3, stdlib email — no add-on)
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                if p.suffix.lower() not in (".eml", ".msg"):
                    return self._json(415, {
                        "error": f"eml_read needs an .eml file, not {p.suffix}",
                        "hint": ".msg (Outlook) needs the extract-msg add-on; "
                                "convert or export as .eml"})
                code, out = _eml_read(p, q)
                return self._json(code, out)

            if u.path == "/image_info":
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                code, out = _image_info(p)
                return self._json(code, out)

            if u.path == "/image_b64":
                # image → data URL for OWUI's display convention (print it,
                # echo as markdown) — size-capped, auto-downscaled with
                # pymupdf when available. Vision INPUT is a different path:
                # OWUI feeds code output to the model as TEXT; a vision model
                # truly "sees" a local file only if the user attaches it.
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                code, out = _image_info(p)
                if code != 200:
                    return self._json(code, out)
                try:
                    max_bytes = min(max(int(q.get("max_bytes", 4_000_000)),
                                        50_000), MAX_BINARY)
                except ValueError:
                    max_bytes = 4_000_000
                data = p.read_bytes()
                mime = {"png": "image/png", "jpeg": "image/jpeg", "gif": "image/gif",
                        "webp": "image/webp", "bmp": "image/bmp"}[out["format"]]
                w, hgt = out.get("width"), out.get("height")
                shrunk = False
                if len(data) > max_bytes and HAVE_PYMUPDF \
                        and out["format"] in ("png", "jpeg", "bmp"):
                    try:
                        pix = fitz.Pixmap(str(p))
                        if pix.colorspace and pix.colorspace.n > 3:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        tries = 0
                        while len(data) > max_bytes and min(pix.width, pix.height) > 64 \
                                and tries < 6:
                            pix.shrink(1)          # halve dimensions
                            data = pix.tobytes("png")
                            shrunk = True
                            tries += 1
                        if shrunk:
                            mime, w, hgt = "image/png", pix.width, pix.height
                    except Exception:
                        pass
                if len(data) > max_bytes:
                    return self._json(413, {
                        "error": f"image is {len(data)} bytes, over the "
                                 f"{max_bytes}-byte cap",
                        "hint": "pass a larger max_bytes (≤ 8 MB), or convert/"
                                f"downscale the image first"
                                + ("" if HAVE_PYMUPDF else
                                   " (this bridge has no pymupdf, so no auto-downscale)")})
                b64 = base64.b64encode(data).decode()
                return self._json(200, {
                    "path": q.get("path"), "mime": mime, "width": w, "height": hgt,
                    "bytes": len(data), "shrunk": shrunk,
                    "data_url": f"data:{mime};base64,{b64}"})

            if u.path == "/reveal":
                # open the OS file manager at the file's location. Consent-
                # gated: refuses unless the LOCAL user enabled it in the
                # picker (allow_reveal in state) — a remote model must never
                # pop windows on the user's desktop unasked.
                if not _state_load().get("allow_reveal"):
                    return self._json(403, {
                        "error": "reveal is disabled on this machine — ask the "
                                 "user to enable it in the File Bridge settings "
                                 "page if they want file-manager popups"})
                p, root, cfg = resolve_guarded(unquote(q.get("path", "")))
                if not p.exists():
                    return self._json(404, {"error": f"no such file: {q.get('path')}"})
                import subprocess as _sp
                if _IS_WINDOWS:
                    r = _sp.run(["explorer", "/select,", str(p)], capture_output=True)
                elif sys.platform == "darwin":
                    r = _sp.run(["open", "-R", str(p)], capture_output=True)
                else:
                    r = _sp.run(["xdg-open", str(p.parent)], capture_output=True)
                if r.returncode != 0:
                    return self._json(500, {"error": "file manager failed to open "
                                                     f"(rc={r.returncode})"})
                return self._json(200, {"ok": True, "revealed": str(p),
                                        "hint": "file manager opened on the user's "
                                                "screen"})

        except ExcludedPath as e:
            return self._json(404, {"error": str(e), "excluded": True,
                                    "hint": "excluded by settings — tell the user; "
                                            "ignore patterns are editable in the "
                                            "File Bridge settings page"})
        except PermissionError as e:
            return self._json(400, {"error": str(e)})
        except Exception as e:
            return self._json(500, {"error": str(e)})
        return self._json(404, {"error": "unknown endpoint"})

    # ---- POST ----

    def do_POST(self):
        u = urlparse(self.path)

        # /api/root is the LOCAL picker API — loopback only, no bridge token
        # (the picker is how you set the token in the first place).
        if u.path == "/api/root":
            if not self._is_loopback():
                return self._json(403, {"error": "picker API is local only"}, cors=False)
            return self._do_api_root()

        # native choose-folder dialog for the picker's Browse… button
        # (loopback + the click itself = the consent; RiskClass.UI)
        if u.path == "/api/pick_folder":
            if not self._is_loopback():
                return self._json(403, {"error": "picker API is local only"}, cors=False)
            self._audit("/api/pick_folder")
            path, err = pick_folder_dialog()
            if err:
                return self._json(500, {"ok": False, "error": err}, cors=False)
            if not path:
                return self._json(200, {"ok": True, "canceled": True}, cors=False)
            return self._json(200, {"ok": True, "path": path}, cors=False)

        # settings-page Stop button (loopback only; response is sent first,
        # the shutdown itself fires ~0.4 s later from a Timer thread)
        if u.path == "/api/shutdown":
            if not self._is_loopback():
                return self._json(403, {"error": "picker API is local only"}, cors=False)
            self._audit("/api/shutdown", args={"source": "picker"})
            if _SHUTDOWN_FN:
                threading.Timer(0.4, _SHUTDOWN_FN).start()
                return self._json(200, {"ok": True, "stopping": True,
                                        "note": "File Bridge is stopping — this "
                                                "page goes offline now"}, cors=False)
            return self._json(500, {"ok": False, "error": "shutdown not available "
                                    "in this mode"}, cors=False)

        root = load_root()
        ok, status, reason = check_request(self.headers)
        if not ok:
            return self._json(status, {"error": reason})
        if root is None:
            return self._json(409, {"error": "No folder chosen yet."})
        try:
            length = int(self.headers.get("Content-Length", 0))
            body = json.loads(self.rfile.read(length) or b"{}")
            # audit BEFORE handling: one line per call, scrubbed in _json
            self._audit(u.path, path=unquote(str(body.get("path", ""))) or None,
                        args={k: body[k] for k in body
                              if k not in ("path", "content", "b64")} or None)
            if u.path == "/write":
                p, root, cfg = resolve_guarded(unquote(body.get("path", "")), for_write=True)
                content = body.get("content", "")
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "write", "path": body.get("path", ""),
                                  "bytes": len(content)}
                if p.exists() and p.is_file() and not confirm_token:
                    # destructive overwrite → two-step confirmation (60 s)
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "target exists — overwrite needs confirmation",
                        "confirmation_required": True,
                        "confirm_op": ("reply to the user with the token; re-send the "
                                       "SAME write with confirmation_token set after "
                                       "the user approves"),
                        **iss})
                if confirm_token:
                    okc, err = confirmation_consume(confirm_token, confirm_params)
                    if not okc:
                        return self._json(400, {"error": err})
                if _readonly_for(cfg):
                    return self._json(403, {"error": "read-only mode is active — "
                                                     "writes are disabled"})
                okr, err = rate_check(len(content))
                if not okr:
                    return self._json(429, {"error": err, "rate_limited": True,
                                            "hint": "relay this to the user and wait "
                                                    "for their confirmation"})
                snap = snapshot_before_write(root, p) if p.exists() else None
                atomic_write_text(p, content)
                return self._json(200, {"ok": True, "written": str(p), "bytes": len(content),
                                        "snapshot": snap})

            if u.path == "/write_b64":
                # binary-safe write: takes base64. For Office/PDF/image files.
                p, root, cfg = resolve_guarded(unquote(body.get("path", "")), for_write=True)
                try:
                    raw = base64.b64decode(body.get("b64", ""), validate=True)
                except (binascii.Error, ValueError) as e:
                    return self._json(400, {"error": f"invalid base64: {e}"})
                if len(raw) > MAX_BINARY:
                    return self._json(413, {"error": f"payload too large: {len(raw)} > {MAX_BINARY} bytes"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "write_b64", "path": body.get("path", ""),
                                  "bytes": len(raw)}
                if p.exists() and p.is_file() and not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "target exists — overwrite needs confirmation",
                        "confirmation_required": True, **iss})
                if confirm_token:
                    okc, err = confirmation_consume(confirm_token, confirm_params)
                    if not okc:
                        return self._json(400, {"error": err})
                if _readonly_for(cfg):
                    return self._json(403, {"error": "read-only mode is active — "
                                                     "writes are disabled"})
                okr, err = rate_check(len(raw))
                if not okr:
                    return self._json(429, {"error": err, "rate_limited": True,
                                            "hint": "relay this to the user and wait "
                                                    "for their confirmation"})
                snap = snapshot_before_write(root, p) if p.exists() else None
                atomic_write_bytes(p, raw)
                return self._json(200, {"ok": True, "written": str(p), "bytes": len(raw),
                                        "snapshot": snap})

            if u.path == "/edit":
                _ep, _er, cfg = resolve_guarded(unquote(body.get("path", "")), for_write=True)
                def _confirm_edit(confirm_params):
                    # returns (http_code_or_0, resp) — mirrors /write flow
                    tok = body.get("confirmation_token")
                    if not tok:
                        iss = confirmation_issue(confirm_params)
                        return 409, {"error": "edit needs confirmation",
                                     "confirmation_required": True, **iss}
                    okc, err = confirmation_consume(tok, confirm_params)
                    if not okc:
                        return 400, {"error": err}
                    return 0, {}

                code, resp = _edit_file(root, cfg, body, _confirm_edit)
                if code >= 400 or resp.get("dry_run"):
                    return self._json(code, resp)
                # apply: snapshot + write via the guarded helper
                p, root, cfg = resolve_guarded(unquote(body.get("path", "")), for_write=True)
                edits = body.get("edits", [])
                modified = p.read_text(encoding="utf-8", errors="replace")
                for e in edits:
                    modified = modified.replace(e.get("old_text", ""),
                                                e.get("new_text", ""),
                                                int(e.get("count", 1) or 1) if e.get("count") else -1)
                okr, err = rate_check(len(modified))
                if not okr:
                    return self._json(429, {"error": err, "rate_limited": True})
                snap = snapshot_before_write(root, p)
                atomic_write_text(p, modified)
                return self._json(200, {"ok": True, "path": body.get("path"),
                                        "snapshot": snap})

            if u.path == "/versions/list":
                rel = body.get("path", "")
                rel = unquote(rel) if isinstance(rel, str) else ""
                _rp, _rr, _rc = resolve_guarded(rel or ".")
                return self._json(200, {"versions": versions_list(_rr, rel)})

            if u.path == "/versions/restore":
                rel = body.get("path", "")
                ts = str(body.get("ts", ""))
                if not rel or not ts:
                    return self._json(400, {"error": "need path + ts (from /versions/list)"})
                _rp, _rr, _rc = resolve_guarded(rel)
                rel = _rp.relative_to(_rr).as_posix()
                # restore is itself a write over an existing file → confirm
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "restore", "path": rel, "ts": ts}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {"error": "restore overwrites the current file — "
                                                     "needs confirmation",
                                            "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                okr, mes = version_restore(_rr, rel, ts)
                if not okr:
                    return self._json(404, {"error": mes})
                return self._json(200, {"ok": True, "restored": mes})

            if u.path == "/delete":
                # NO unlink ever: delete = move to trash (structural isolation
                # outside all roots). Two-step confirmation like overwrites.
                rel = body.get("path", "")
                if not rel:
                    return self._json(400, {"error": "need path"})
                p, root, cfg = resolve_guarded(unquote(rel), for_write=True)
                if not p.exists():
                    return self._json(404, {"error": f"not found: {rel}"})
                if _readonly_for(cfg):
                    return self._json(403, {"error": "read-only mode is active"})
                confirm_params = {"op": "delete", "path": rel}
                confirm_token = body.get("confirmation_token")
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    desc = f"{p.stat().st_size} bytes" if p.is_file() else "directory"
                    return self._json(409, {
                        "error": f"deletion needs confirmation ({desc}) — the file "
                                 f"goes to the trash store, not unrecoverable",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                okr, err = rate_check(0)
                if not okr:
                    return self._json(429, {"error": err, "rate_limited": True})
                info = trash_move(root, p)
                trash_purge(root)
                return self._json(200, {"ok": True, "trashed": info,
                                        "recover": "POST /trash/restore"})

            if u.path == "/trash/list":
                rel = body.get("path", "")
                rel = unquote(rel) if isinstance(rel, str) else ""
                _rp, _rr, _rc = resolve_guarded(rel or ".")
                return self._json(200, {"trash": trash_list(_rr, rel)})

            if u.path == "/trash/restore":
                rel = body.get("path", "")
                ts = str(body.get("ts", ""))
                if not rel or not ts:
                    return self._json(400, {"error": "need path + ts (from /trash/list)"})
                _rp, _rr, _rc = resolve_guarded(rel)
                if _readonly_for(_rc):
                    return self._json(403, {"error": "read-only mode is active"})
                okr, mes = trash_restore(_rr, rel, ts)
                if not okr:
                    return self._json(409, {"error": mes})
                return self._json(200, {"ok": True, "restored": mes})

            if u.path == "/trash/purge":
                return self._json(403, {"error": "manual purge is a settings-page "
                                                 "action, not an API call"})

            if u.path == "/zip":
                _, _, cfg = resolve_guarded(".")
                code, resp = zip_create(root, cfg, body)
                return self._json(code, resp)

            if u.path == "/unzip":
                _, _, cfg = resolve_guarded(".")
                code, resp = zip_extract(root, cfg, body)
                return self._json(code, resp)

            if u.path == "/ocr_pdf":
                # searchable PDF: page images + invisible text layer
                if not TESSERACT_BIN:
                    return self._json(501, {"error": "OCR unavailable: tesseract not "
                                                     "found. Install tesseract or set "
                                                     "TESSERACT_CMD."})
                if not HAVE_PYMUPDF:
                    return self._json(501, {"error": "ocr_pdf needs the PDF add-on "
                                                     "(pip install pymupdf)"})
                rel = body.get("path", "")
                out_rel = body.get("out", "")
                if not rel or not out_rel:
                    return self._json(400, {"error": "need path (image or scanned PDF) "
                                                     "+ out (new .pdf, root-relative)"})
                if Path(out_rel).suffix.lower() != ".pdf":
                    return self._json(400, {"error": "out must end in .pdf"})
                p, _r, _c = resolve_guarded(unquote(rel))
                ext = p.suffix.lower()
                if ext not in ({".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"}
                               | {".pdf"}):
                    return self._json(400, {"error": f"ocr_pdf supports images and "
                                                     f"PDF, not {ext}"})
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {rel}"})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                # source read + result write → confirm BEFORE raster/OCR work
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "ocr_pdf", "path": rel, "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "ocr_pdf reads the source and writes a new PDF — "
                                 "needs confirmation",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                raw_lang = body.get("lang") or _get_ocr_lang()
                parts = [x for x in re.split(r"[\s,+]+", raw_lang)
                         if x and re.fullmatch(r"[a-zA-Z_]{2,8}", x)]
                lang = "+".join(parts) if parts else "eng"
                try:
                    dpi = max(72, min(int(body.get("dpi", 200)), 400))
                except (TypeError, ValueError):
                    dpi = 200
                try:
                    max_pages = max(1, min(int(body.get("max_pages", _OCR_PDF_MAX_PAGES)),
                                           _OCR_PDF_MAX_PAGES))
                except (TypeError, ValueError):
                    max_pages = _OCR_PDF_MAX_PAGES
                code, resp = ocr_pdf(p, op, oroot, lang, dpi, max_pages)
                return self._json(code, resp)

            if u.path == "/pdf_op":
                # split | merge | rotate (pymupdf, P3)
                op = body.get("op", "")
                if op not in ("split", "merge", "rotate"):
                    return self._json(400, {"error": "op must be split|merge|rotate"})
                if not HAVE_PYMUPDF:
                    return self._json(501, {"error": "pdf_op needs the PDF add-on "
                                                     "(pip install pymupdf)"})
                srcs_rel = body.get("paths") or ([body.get("path")] if body.get("path") else [])
                if not isinstance(srcs_rel, list) or not srcs_rel or len(srcs_rel) > 20:
                    return self._json(400, {"error": "paths must be a 1-20 list"})
                out_rel = body.get("out", "")
                if not out_rel or Path(out_rel).suffix.lower() != ".pdf":
                    return self._json(400, {"error": "need out (root-relative .pdf)"})
                if op != "merge" and len(srcs_rel) > 1:
                    return self._json(400, {"error": f"{op} takes exactly one input"})
                resolved = []
                for r in srcs_rel:
                    if not isinstance(r, str) or not r:
                        return self._json(400, {"error": f"bad path: {r!r}"})
                    sp, _sr, _sc = resolve_guarded(unquote(r))
                    if sp.suffix.lower() != ".pdf":
                        return self._json(400, {"error": f"input must be .pdf, not "
                                                         f"{sp.suffix}"})
                    if not sp.is_file():
                        return self._json(404, {"error": f"no such file: {r}"})
                    resolved.append(sp)
                opath, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                try:
                    angle = int(body.get("angle", 90)) % 360
                except (TypeError, ValueError):
                    angle = 90
                if op == "rotate" and angle == 0:
                    return self._json(400, {"error": "angle 0 does nothing"})
                spec_pages = str(body.get("pages", "") or "")
                # confirmation mirrors /write: only when OUT already exists
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "pdf_op", "pdfop": op,
                                  "paths": srcs_rel, "out": out_rel}
                if opath.exists() and not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "target exists — overwrite needs confirmation",
                        "confirmation_required": True, **iss})
                if confirm_token:
                    okc, err = confirmation_consume(confirm_token, confirm_params)
                    if not okc:
                        return self._json(400, {"error": err})
                code, resp = pdf_op(op, resolved, opath, oroot, spec_pages, angle)
                return self._json(code, resp)

            if u.path == "/docx_merge":
                # fill {{placeholders}} in a .docx template (openworker #454)
                rel = body.get("path", "")
                out_rel = body.get("out", "")
                if not rel or not out_rel:
                    return self._json(400, {"error": "need path (.docx template) + "
                                                     "out (root-relative .docx)"})
                if Path(out_rel).suffix.lower() != ".docx":
                    return self._json(400, {"error": "out must end in .docx"})
                p, _r, _c = resolve_guarded(unquote(rel))
                if p.suffix.lower() != ".docx":
                    return self._json(400, {"error": f"template must be .docx, "
                                                     f"not {p.suffix}"})
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {rel}"})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                values = body.get("values") or {}
                if not isinstance(values, dict):
                    return self._json(400, {"error": "values must be an object of "
                                                     "{placeholder: text}"})
                if len(json.dumps(values, default=str)) > MAX_READ:
                    return self._json(413, {"error": "values payload too large"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "docx_merge", "path": rel, "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "docx_merge writes a new file — needs confirmation",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                code, resp = docx_merge(p, op, oroot, values,
                                        bool(body.get("strict")))
                return self._json(code, resp)

            if u.path == "/pptx_from_template":
                # build a deck from a .potx/.pptx layout template (#454)
                rel = body.get("path", "")
                out_rel = body.get("out", "")
                if not rel or not out_rel:
                    return self._json(400, {"error": "need path (.potx/.pptx "
                                                     "template) + out (.pptx)"})
                if Path(out_rel).suffix.lower() != ".pptx":
                    return self._json(400, {"error": "out must end in .pptx"})
                p, _r, _c = resolve_guarded(unquote(rel))
                if p.suffix.lower() not in {".potx", ".pptx"}:
                    return self._json(400, {"error": f"template must be .potx or "
                                                     f".pptx, not {p.suffix}"})
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {rel}"})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                values = body.get("values") or {}
                slides = body.get("slides") or []
                if not isinstance(values, dict) or not isinstance(slides, list):
                    return self._json(400, {"error": "values must be an object, "
                                                     "slides a list of {layout, "
                                                     "title, body}"})
                if len(slides) > 100:
                    return self._json(400, {"error": "slides capped at 100"})
                # validate layout indices BEFORE the confirmation flow so bad
                # specs fail fast (no token burn on typos)
                for spec in slides:
                    if not isinstance(spec, dict):
                        return self._json(400, {"error": "each slide must be an object "
                                                         "{layout, title, body}"})
                    li = spec.get("layout", 1)
                    if not isinstance(li, int) or isinstance(li, bool) or not (0 <= li <= 11):
                        return self._json(400, {"error": f"bad layout index {li!r} — "
                                                          f"valid range depends on the "
                                                          f"template (usually 0-11)"})
                if len(json.dumps({"v": values, "s": slides}, default=str)) > MAX_READ:
                    return self._json(413, {"error": "payload too large"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "pptx_from_template", "path": rel,
                                  "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "pptx_from_template writes a new file — needs "
                                 "confirmation",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                code, resp = pptx_from_template(p, op, oroot, slides, values)
                return self._json(code, resp)

            if u.path == "/pdf_from_text":
                # bridge-native PDF authoring (P3) — no Pyodide shim needed
                out_rel = body.get("out", "")
                if not out_rel or Path(out_rel).suffix.lower() != ".pdf":
                    return self._json(400, {"error": "need out (root-relative "
                                                     ".pdf)"})
                blocks = body.get("blocks")
                if not isinstance(blocks, list) or not blocks:
                    return self._json(400, {"error": "blocks must be a list of "
                                                     "{style: title|h1|h2|body|"
                                                     "pagebreak, text}"})
                if len(blocks) > 5000:
                    return self._json(413, {"error": "blocks capped at 5000"})
                for b in blocks:
                    if not isinstance(b, dict) or \
                       (b.get("style", "body") not in
                            ("title", "h1", "h2", "body", "pagebreak")) or \
                       ("text" in b and not isinstance(b.get("text"), str)):
                        return self._json(400, {"error": "each block is "
                                                         "{style, text} (style "
                                                         "title|h1|h2|body|"
                                                         "pagebreak)"})
                if len(json.dumps(blocks, default=str)) > MAX_READ:
                    return self._json(413, {"error": "payload too large"})
                page_size = str(body.get("page_size", "a4")).lower()
                if page_size not in ("a4", "letter"):
                    return self._json(400, {"error": "page_size must be a4|letter"})
                title = body.get("title", "")
                if title and not isinstance(title, str):
                    return self._json(400, {"error": "title must be a string"})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "pdf_from_text", "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "pdf_from_text writes a new file — needs "
                                 "confirmation",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                code, resp = pdf_from_text(op, oroot, title, blocks, page_size)
                return self._json(code, resp)

            if u.path == "/docx_write":
                # structured Word authoring by sections (P3)
                out_rel = body.get("out", "")
                if not out_rel or Path(out_rel).suffix.lower() != ".docx":
                    return self._json(400, {"error": "need out (root-relative "
                                                     ".docx)"})
                sections = body.get("sections")
                if not isinstance(sections, list) or not sections:
                    return self._json(400, {"error": "sections must be a list of "
                                                     "{style: h1|h2|paragraph|list|"
                                                     "numbered|pagebreak, text, "
                                                     "items?}"})
                if len(sections) > 2000:
                    return self._json(413, {"error": "sections capped at 2000"})
                for s in sections:
                    if not isinstance(s, dict) or \
                       (s.get("style", "paragraph") not in
                            ("h1", "h2", "paragraph", "pagebreak") +
                            _DOCX_LIST_STYLES) or \
                       ("text" in s and not isinstance(s.get("text"), str)):
                        return self._json(400, {"error": "each section is {style, "
                                                         "text, items?}"})
                    if "items" in s and not isinstance(s.get("items"), list):
                        return self._json(400, {"error": "items must be a list of "
                                                         "strings"})
                if len(json.dumps(sections, default=str)) > MAX_READ:
                    return self._json(413, {"error": "payload too large"})
                title = body.get("title", "")
                if title and not isinstance(title, str):
                    return self._json(400, {"error": "title must be a string"})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "docx_write", "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "docx_write writes a new file — needs "
                                 "confirmation",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                code, resp = docx_write(op, oroot, title, sections)
                return self._json(code, resp)

            if u.path == "/xlsx_append":
                # append rows to an existing .xlsx or create it (P3)
                rel = body.get("path", "")
                if not rel or Path(rel).suffix.lower() != ".xlsx":
                    return self._json(400, {"error": "need path (root-relative "
                                                     ".xlsx; created if absent)"})
                rows = body.get("rows")
                if not isinstance(rows, list) or not rows:
                    return self._json(400, {"error": "rows must be a list of "
                                                     "lists (cell values)"})
                if len(rows) > 1000:
                    return self._json(413, {"error": "rows capped at 1000 per call"})
                for r in rows:
                    if not isinstance(r, list) or len(r) > 256 or \
                       any(isinstance(v, (dict, list)) for v in r):
                        return self._json(400, {"error": "each row is a flat list "
                                                         "of cell values (≤256)"})
                if len(json.dumps(rows, default=str)) > MAX_READ:
                    return self._json(413, {"error": "payload too large"})
                sheet = body.get("sheet") or None
                if sheet is not None and (not isinstance(sheet, str) or
                                          len(sheet) > 64):
                    return self._json(400, {"error": "sheet must be a short string"})
                header = body.get("header")
                if header is not None and (not isinstance(header, list) or
                                           len(header) > 256):
                    return self._json(400, {"error": "header must be a flat list "
                                                     "(used only when creating)"})
                p, proot, pcfg = resolve_guarded(unquote(rel), for_write=True)
                if _readonly_for(pcfg):
                    return self._json(403, {"error": "read-only mode is active"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "xlsx_append", "path": rel}
                if p.exists() and not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "target exists — appending modifies it, "
                                 "needs confirmation",
                        "confirmation_required": True, **iss})
                if confirm_token:
                    okc, err = confirmation_consume(confirm_token, confirm_params)
                    if not okc:
                        return self._json(400, {"error": err})
                code, resp = xlsx_append(p, rows, sheet, proot, header)
                return self._json(code, resp)


            if u.path == "/docx_mailmerge":
                # one document PER ROW from xlsx/csv/inline (P3, #454 family)
                rel = body.get("path", "")
                out_rel = body.get("out", "")
                if not rel or not out_rel:
                    return self._json(400, {"error": "need path (.docx template) "
                                                     "+ out (pattern like "
                                                     "'out/{{name}}-contract.docx' "
                                                     "or a .zip)"})
                if not (out_rel.lower().endswith(".zip") or
                        out_rel.lower().endswith(".docx")):
                    return self._json(400, {"error": "out must end in .docx "
                                                     "(name pattern) or .zip"})
                to_zip = out_rel.lower().endswith(".zip")
                p, _r, _c = resolve_guarded(unquote(rel))
                if p.suffix.lower() != ".docx":
                    return self._json(400, {"error": f"template must be .docx, "
                                                     f"not {p.suffix}"})
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {rel}"})
                rows_src = body.get("rows")
                src_path = None
                if isinstance(rows_src, list):
                    if not rows_src or len(rows_src) > 50:
                        return self._json(400, {"error": "rows: 1-50 objects"})
                    rows = rows_src
                elif isinstance(rows_src, str) and rows_src:
                    sp, _sr, _sc = resolve_guarded(unquote(rows_src))
                    if sp.suffix.lower() not in (".xlsx", ".csv"):
                        return self._json(400, {"error": "rows file must be "
                                                         ".xlsx or .csv"})
                    if not sp.is_file():
                        return self._json(404, {"error": f"no such file: "
                                                         f"{rows_src}"})
                    src_path = sp
                    rows = None  # loaded after cheap validation
                else:
                    return self._json(400, {"error": "rows must be a list of "
                                                     "objects or a path to "
                                                     ".xlsx/.csv"})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                if rows is None:
                    code, res = _mailmerge_rows(src_path, None)
                    if code != 200:
                        return self._json(code, res)
                    rows = res["rows"]
                    if not rows:
                        return self._json(400, {"error": "row source has no data "
                                                         "rows"})
                # validate the out pattern resolves per-row BEFORE confirm
                if not to_zip:
                    probes = []
                    for row in rows[:50]:
                        nm = _out_name_for(out_rel,
                                           {str(k): str(v) for k, v in
                                            (row.items() if isinstance(row, dict)
                                             else [])})
                        probes.append(nm)
                    if any("{{" in n for n in probes):
                        return self._json(400, {"error": "out pattern has "
                                                         "placeholders not "
                                                         "present in the rows "
                                                         "(e.g. {{name}})"})
                    if len(set(probes)) != len(probes):
                        return self._json(400, {"error": "out pattern collides "
                                                         "(two rows map to the "
                                                         "same filename) — add "
                                                         "a unique column to "
                                                         "the pattern"})
                if len(json.dumps(rows, default=str)) > MAX_READ:
                    return self._json(413, {"error": "payload too large"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "docx_mailmerge", "path": rel,
                                  "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": f"docx_mailmerge writes {len(rows)} files — "
                                 "needs confirmation",
                        "confirmation_required": True, "documents": len(rows),
                        **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                code, resp = docx_mailmerge(p, out_rel, oroot, _r, rows, to_zip)
                return self._json(code, resp)


            if u.path == "/convert":
                # LibreOffice headless format conversion (P3)
                rel = body.get("path", "")
                out_rel = body.get("out", "")
                if not rel or not out_rel:
                    return self._json(400, {"error": "need path + out "
                                                     "(root-relative, format "
                                                     "follows the extension)"})
                p, _r, _c = resolve_guarded(unquote(rel))
                if not p.is_file():
                    return self._json(404, {"error": f"no such file: {rel}"})
                to_ext = Path(out_rel).suffix.lower()
                src_ext = p.suffix.lower()
                if src_ext == to_ext:
                    return self._json(400, {"error": "out extension must differ "
                                                     "from the source format"})
                supported = set(_CONV_MATRIX.get(src_ext, {}))
                if to_ext not in supported:
                    targets = ", ".join(sorted(supported)) or "none"
                    return self._json(400, {"error": f"cannot convert {src_ext} "
                                                     f"-> {to_ext or '(none)'}",
                                            "supported_targets": targets})
                op, oroot, ocfg = resolve_guarded(unquote(out_rel), for_write=True)
                if _readonly_for(ocfg):
                    return self._json(403, {"error": "read-only mode is active"})
                confirm_token = body.get("confirmation_token")
                confirm_params = {"op": "convert", "path": rel, "out": out_rel}
                if not confirm_token:
                    iss = confirmation_issue(confirm_params)
                    return self._json(409, {
                        "error": "convert writes a new file — needs confirmation",
                        "confirmation_required": True, **iss})
                okc, err = confirmation_consume(confirm_token, confirm_params)
                if not okc:
                    return self._json(400, {"error": err})
                code, resp = convert_file(p, op, oroot, to_ext)
                return self._json(code, resp)


            if u.path == "/write_many":
                # batch writes; >5 files needs {"confirmed": true} (mass-edit
                # detection, roadmap P0b)
                items = body.get("items", [])
                if not isinstance(items, list) or not items or len(items) > 50:
                    return self._json(400, {"error": "items must be a 1-50 list"})
                if len(items) > 5 and not body.get("confirmed"):
                    return self._json(409, {
                        "error": f"batch of {len(items)} files needs explicit "
                                 f"confirmation",
                        "confirmation_required": True,
                        "confirm_op": "list the plan to the user; on approval "
                                      "re-send with confirmed:true",
                        "plan": [{"path": it.get("path"),
                                  "bytes": len(it.get("content", ""))}
                                 for it in items]})
                results = []
                for it in items:
                    path = it.get("path", "")
                    content = it.get("content", "")
                    try:
                        p, root, cfg = resolve_guarded(unquote(path), for_write=True)
                        if _readonly_for(cfg):
                            raise PermissionError("read-only mode")
                        okr, err = rate_check(len(content))
                        if not okr:
                            return self._json(429, {"error": err, "rate_limited": True,
                                                    "hint": "batch aborted — files "
                                                            "written so far are listed",
                                                    "results": results})
                        snap = snapshot_before_write(root, p) if p.exists() else None
                        atomic_write_text(p, content)
                        results.append({"path": path, "ok": True, "bytes": len(content),
                                        "snapshot": snap})
                    except (PermissionError, ExcludedPath) as e:
                        results.append({"path": path, "ok": False, "error": str(e)})
                return self._json(200, {"ok": all(r.get("ok") for r in results),
                                        "results": results})
        except ExcludedPath as e:
            return self._json(404, {"error": str(e), "excluded": True,
                                    "hint": "excluded by settings — tell the user; "
                                            "ignore patterns are editable in the "
                                            "File Bridge settings page"})
        except PermissionError as e:
            return self._json(400, {"error": str(e)})
        except Exception as e:
            return self._json(500, {"error": str(e)})
        return self._json(404, {"error": "unknown endpoint"})

    def _is_loopback(self):
        try:
            ip = self.client_address[0]
            return ip in ("127.0.0.1", "::1", "localhost")
        except Exception:
            return False

    # ---- local picker API (/api/root) ----

    def _do_api_root(self):
        # CSRF guard: browsers only skip the CORS preflight for "simple"
        # requests — application/json is NOT simple, so requiring it means
        # a cross-site page cannot reach this state-changing API with a
        # sneaky text/plain body (the request would need preflight, and
        # foreign origins get no CORS headers).
        ctype = (self.headers.get("Content-Type") or "").split(";")[0].strip().lower()
        if ctype != "application/json":
            return self._json(400, {"ok": False,
                                    "error": "Content-Type must be application/json"},
                              cors=False)
        length = int(self.headers.get("Content-Length", 0))
        try:
            body = json.loads(self.rfile.read(length) or b"{}")
        except Exception:
            return self._json(400, {"ok": False, "error": "invalid JSON"}, cors=False)

        self._audit("/api/root", args={k: "[redacted]" for k in body})  # settings changes: keys only

        if "ocr_lang" in body:
            # tesseract wants codes joined with '+' (e.g. swe+eng). Users
            # instinctively type 'eng, sv' — normalize separators, then
            # validate format + membership so typos fail HERE with a clear
            # message instead of at OCR time inside tesseract.
            raw = str(body["ocr_lang"]).strip()
            norm = re.sub(r"[\s,;+]+", "+", raw).strip("+") if raw else "eng"
            if not re.fullmatch(r"[A-Za-z0-9_]+(\+[A-Za-z0-9_]+)*", norm):
                return self._json(400, {"ok": False,
                                        "error": "OCR languages must be tesseract "
                                                 "codes joined with '+', e.g. "
                                                 "eng+swe"}, cors=False)
            avail = _ocr_langs_available()
            if avail:
                missing = [c for c in norm.split("+") if c not in avail]
                if missing:
                    return self._json(400, {"ok": False,
                                            "error": f"language not installed: "
                                                     f"{', '.join(missing)} "
                                                     f"(installed: {', '.join(avail)})"},
                                      cors=False)
            _set_ocr_lang(norm)

        # security settings ------------------------------------------------
        if "allowed_origin" in body:
            val = body["allowed_origin"]
            if val is None or (isinstance(val, str) and val.strip().lower() in ("", "none", "off")):
                set_allowed_origin(None)
            else:
                okn, norm = set_allowed_origin(str(val))
                if not okn:
                    return self._json(400, {"ok": False,
                                            "error": "invalid origin (want e.g. http://owui.internal:8080)"},
                                      cors=False)

        tok_action = body.get("token")
        new_token = None
        if isinstance(tok_action, dict):  # {"set": "..."} | {"clear": true} | {"generate": true}
            if tok_action.get("clear"):
                set_token(None)
            elif tok_action.get("generate"):
                new_token = generate_token()
            elif tok_action.get("set"):
                if not set_token(str(tok_action["set"])):
                    return self._json(400, {"ok": False, "error": "token must be 8-256 chars"},
                                      cors=False)
                new_token = str(tok_action["set"])  # show back once to the local admin
        elif tok_action is None and "token" in body:
            set_token(None)

        if isinstance(body.get("roots"), list):
            okr, err = set_roots(body["roots"])
            if not okr:
                return self._json(400, {"ok": False, "error": err}, cors=False)
        if "readonly" in body:
            _state_update(readonly=bool(body["readonly"]))
        if "allow_reveal" in body:
            _state_update(allow_reveal=bool(body["allow_reveal"]))
        if isinstance(body.get("ignore_global"), list):
            _state_update(ignore_global=[str(x)[:200] for x in body["ignore_global"]][:200])
        if body.get("root") and not isinstance(body.get("roots"), list):
            # legacy single-root setter (CLI arg path)
            p = Path(body["root"]).expanduser()
            if not p.is_dir():
                return self._json(400, {"ok": False, "error": f"not a folder: {p}"}, cors=False)
            save_root(p)
        ros = enabled_roots()
        resp = {"ok": True, "ocr_lang": _get_ocr_lang(),
                "roots": roots_config(), "root": ros[0]["path"] if ros else None}
        resp["allowed_origin"] = get_allowed_origin()
        resp["security"] = security_mode()
        if new_token:
            resp["token"] = new_token  # shown ONCE in the local picker only
        return self._json(200, resp, cors=False)

    # ---- local picker UI (served only to localhost browser) ----
    PICKER_HTML = """<!doctype html><html><head><meta charset="utf-8"><title>File Bridge</title>
<style>body{font-family:system-ui;max-width:680px;margin:50px auto;padding:0 20px;color:#222}
input:not([type=checkbox]),textarea{width:100%;padding:10px;font-size:16px;box-sizing:border-box}
textarea{font-family:ui-monospace,monospace}
#langbox label{white-space:nowrap;cursor:pointer}
#langbox input[type=checkbox]{width:auto;margin:0 6px 0 0;vertical-align:middle}
button{padding:10px 22px;font-size:16px;margin-top:12px;cursor:pointer}
button.small{margin-top:0;padding:4px 12px;font-size:13px}
.btnrow{display:flex;gap:8px;align-items:stretch}
h3{margin:0 0 8px}
.ok{color:#0a7d32;font-weight:bold}.hint{color:#666;font-size:14px}
.warn{color:#b00;font-weight:bold}.sec{background:#f4f6f8;padding:14px 18px;border-radius:8px;margin:16px 0}
.panel{background:#fff;border:1px solid #ddd;border-radius:6px;padding:10px 14px}
code{background:#eee;padding:2px 6px;border-radius:4px}</style></head><body>
<h2>📁 File Bridge <span id="beat" style="font-size:18px;color:#999">●</span></h2>
<p class="hint" id="beatinfo">checking whether the bridge is running…</p>
<p>This little service lets <b>Open WebUI in your browser</b> read &amp; write files
in <b>one folder you choose</b> on this computer. Nothing else is exposed.</p>
<div class="sec">
<h3>📁 Shared folder</h3>
<p class="hint">Only this folder is exposed — nothing else on this computer.</p>
<div class="btnrow">
<input id="root" placeholder="__ROOTPH__" value="__ROOT__" style="flex:1">
<button id="browsebtn" onclick="browse()" class="small" style="white-space:nowrap">Browse…</button>
</div>
<button onclick="setRoot()">Save folder</button>
<p class="ok" id="status"></p>
</div>
<div class="sec">
<h3>🔒 Security</h3>
<p class="hint">Lock the bridge to your Open WebUI address (required — until set,
file endpoints stay disabled):</p>
<input id="origin" placeholder="http://owui.yourcompany.com:8080" value="__ORIGIN__">
<button onclick="setOrigin()">Lock to this origin</button>
<p class="hint">Optional extra: org-wide token (Tier&nbsp;2) — matching requests must also
send it in an <code>X-Bridge-Token</code> header. Org flow (recommended): your OWUI admin
embeds a token in the public skill with <code>setup_owui.py --bridge-token</code>; paste
that token here so your bridge accepts exactly your organisation's requests. Or generate
a random one and give it to your admin to embed.</p>
<div class="btnrow">
<input id="token" placeholder="paste the org token from your OWUI admin" style="flex:1" autocomplete="off">
<button onclick="setToken()" class="small" style="white-space:nowrap">Set token</button>
</div>
<button onclick="genToken()">Generate random token</button>
<button onclick="clearToken()">Clear token</button>
<p class="hint" id="secstatus"></p>
</div>
<div class="sec">
<h3>🔤 OCR language</h3>
<p class="hint">For reading scanned PDFs / photos — tick one or more:</p>
<div id="langbox" class="panel" style="font-size:15px;display:flex;flex-wrap:wrap;gap:8px 22px;align-items:center"></div>
<p class="hint">Ticks combine automatically in tesseract syntax (e.g. <code>eng+swe</code>
— combining is free, mixing e.g. Swedish AND English fixes å/ä/ö and digits).
Need a code that is not listed? Type it below before saving.</p>
<input id="ocrlang" placeholder="eng+swe" value="__OCRLANG__" style="max-width:200px" oninput="syncBoxes(false)">
<button onclick="setLang()">Save language</button>
<p class="hint" id="langs"></p>
</div>
<div class="sec">
<h3>🚫 Ignore patterns</h3>
<p class="hint">Files matching these patterns are invisible to the AI — not
listed, not readable, and <b>writes to them are refused</b>. One pattern per
line, gitignore-style: <code>*.zip</code>, <code>secrets/</code> (that folder,
anywhere), <code>/build</code> (top level only); bare names match at any
depth. The preview below updates within seconds of saving.</p>
<textarea id="ignorepats" rows="4" placeholder="*.zip&#10;.secrets&#10;node_modules/">__IGNORE__</textarea>
<div class="btnrow" style="align-items:center;gap:10px;margin-top:6px">
<button onclick="setIgnore()" class="small" style="padding:6px 16px;font-size:14px">Save patterns</button>
<span class="hint" id="ignstat" style="margin:0"></span>
</div>
<p class="hint">Always ignored (built in, not editable):
<b>.DS_Store</b> · <b>._*</b> · <b>Thumbs.db</b> · <b>desktop.ini</b></p>
</div>
<div class="sec">
<div class="btnrow" style="align-items:center;gap:10px">
<h3 style="margin:0">👁 What the AI can see</h3>
<button onclick="renderPreview()" class="small">↻ Refresh</button>
</div>
<p class="hint">Exactly what Open WebUI's model sees when it lists your folders —
after ignore lists and security rules. Folders are collapsible; the view
auto-refreshes every 5&nbsp;s while this tab is visible. Nothing here is
editable; change the folder above instead.</p>
<div id="preview" class="panel hint" style="max-height:340px;overflow:auto;font-family:ui-monospace,monospace;
font-size:13px;color:#333;text-align:left"></div>
<p class="hint" id="previnfos"></p>
</div>
<hr>
<p class="hint">Status: __STATUS__<br>
You can close this browser tab — the service keeps running in the background.
It stops only when you stop it:</p>
<button onclick="stopBridge()" style="background:#fff1f0">Stop File Bridge</button>
<p class="hint">…or quit it like any app: right-click its icon in the Dock (macOS),
Ctrl+C in a terminal, or stop the service (Linux).</p>
<script>
const LANG_NAMES={eng:'English',swe:'Swedish',chi_sim:'Chinese (Simplified)',
chi_tra:'Chinese (Traditional)',dan:'Danish',nor:'Norwegian',fin:'Finnish',
isl:'Icelandic',deu:'German',fra:'French',spa:'Spanish',ita:'Italian',
por:'Portuguese',nld:'Dutch',rus:'Russian',ukr:'Ukrainian',pol:'Polish',
ces:'Czech',slk:'Slovak',slv:'Slovenian',hrv:'Croatian',srp:'Serbian',
bul:'Bulgarian',ell:'Greek',hun:'Hungarian',ron:'Romanian',tur:'Turkish',
ara:'Arabic',heb:'Hebrew',jpn:'Japanese',kor:'Korean',hin:'Hindi',tha:'Thai',
vie:'Vietnamese',ind:'Indonesian',cat:'Catalan',glg:'Galician',eus:'Basque',
lit:'Lithuanian',lav:'Latvian',est:'Estonian',osd:'auto-detect (page orientation)'};
function esc(s){return String(s).replace(/[&<>"']/g,m=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[m]));}
function fmtSize(n){if(n==null)return '';if(n<1024)return n+' B';if(n<1048576)return (n/1024).toFixed(1)+' KB';return (n/1048576).toFixed(1)+' MB';}
async function refresh(){const s=await (await fetch('/state')).json();
document.getElementById('root').value=s.root||'';
document.getElementById('origin').value=s.allowed_origin||'';
document.getElementById('secstatus').textContent='Security mode: '+s.security+
(String(s.security).indexOf('token')>=0
 ?' · a token IS configured (shown as dots below; paste a new one only to replace it)'
 :' · no token set');
document.getElementById('token').value=
 (String(s.security).indexOf('token')>=0?'••••••••••••':'');
const c=await (await fetch('/ocr/config')).json();
document.getElementById('ocrlang').value=c.lang||'eng';
renderLangs(c.available,c.lang);
document.getElementById('langs').textContent='engine: '+(c.engine||'?')+
 (c.user_dir?' — add languages: drop .traineddata files (tessdata_fast) into '+c.user_dir+' and restart the app':'');
renderPreview();}
function renderLangs(avail,cur){
 const box=document.getElementById('langbox');
 const sel=new Set(String(cur||'').split('+').filter(Boolean));
 box.innerHTML=(avail&&avail.length)?avail.map(c=>
  '<label>'+
  '<input type="checkbox" value="'+esc(c)+'"'+(sel.has(c)?' checked':'')+
  ' onchange="syncBoxes(true)"> '+esc(c)+
  (LANG_NAMES[c]?' — '+LANG_NAMES[c]:'')+'</label>').join(' ')
  :'<span class="hint">no installed language files found — type codes below instead</span>';}
function syncBoxes(fromBoxes){
 const inp=document.getElementById('ocrlang');
 if(fromBoxes){
  const v=[...document.querySelectorAll('#langbox input:checked')].map(x=>x.value);
  inp.value=v.join('+');
 }else{
  const sel=new Set(inp.value.split(/[+,\\s]+/).filter(Boolean));
  document.querySelectorAll('#langbox input').forEach(x=>x.checked=sel.has(x.value));
 }}
async function beat(){
 const dot=document.getElementById('beat'),info=document.getElementById('beatinfo');
 try{
  const h=await fetch('/health',{cache:'no-store'});
  if(!h.ok)throw new Error('http '+h.status);
  const j=await h.json();
  dot.style.color='#0a7d32';
  info.textContent='Running · v'+(j.version||'?')+' · security: '+(j.security||'?')+
   (j.ok?'':' · no folder chosen yet');
 }catch(e){
  dot.style.color='#b00';
  info.textContent='no response — the bridge has stopped';
 }}
setInterval(beat,5000);beat();
async function renderPreview(){
 if(window._pvBusy)return;window._pvBusy=true;
 const box=document.getElementById('preview'), info=document.getElementById('previnfos');
 const t0=performance.now();
 try{
  const h=await (await fetch('/health')).json();
  if(!h.ok){box.innerHTML='🔒 '+(h.hint||'no folder chosen yet');info.textContent='';return;}
  const t=await (await fetch('/api/preview')).json();
  if(!t.ok&&t.error&&!t.entries){box.innerHTML='🔒 '+esc(t.error);info.textContent='';return;}
  // remember which folders the user had open so auto-refresh doesn't collapse them
  const openPaths=new Set([...box.querySelectorAll('details[open]')].map(d=>d.dataset.p));
  let nf=0,nd=0,count=0;
  function row(node,depth,path){
   if(count>=500)return '';
   count++;
   const p=path+'/'+node.name;
   const size=node.size!=null?' <span style="color:#999">'+fmtSize(node.size)+'</span>':'';
   if(node.children){
    nd++;
    const n=node.children.length;
    const open=openPaths.has(p)||depth<1;
    const kids=node.children.map(ch=>row(ch,depth+1,p)).join('');
    return '<details data-p="'+esc(p)+'"'+(open?' open':'')+'>'+
     '<summary style="cursor:pointer">📁 '+esc(node.name)+
     ' <span style="color:#999">'+n+'</span></summary>'+
     '<div style="margin-left:14px">'+kids+'</div></details>';
   }
   nf++;
   return '<div>📄 '+esc(node.name)+size+'</div>';
  }
  const html=t.entries.map(n=>row(n,0,'')).join('');
  box.innerHTML=html||'(empty folder)';
  info.textContent=(nf)+' file'+(nf===1?'':'s')+' · '+nd+' folder'+(nd===1?'':'s')+
   (t.truncated?' · TRUNCATED at cap — the model sees the same limit':'')+
   ' · ignore lists applied · symlinks never shown';
 }catch(e){box.innerHTML='preview unavailable: '+esc(e.message||e);}
 finally{
  window._pvBusy=false;
  // adaptive cadence: a slow walk (huge folder / cold cache) means the
  // user mounted something big — poll gently (30 s) until it's fast again
  window._pvMs=performance.now()-t0;
 }
}
(function pvLoop(){
 setTimeout(async()=>{
  if(document.visibilityState==='visible')await renderPreview();
  const slow=(window._pvMs||0)>1500?30000:5000;
  pvLoop();
 },(window._pvMs||0)>1500?30000:5000);
})();
document.addEventListener('visibilitychange',()=>{if(document.visibilityState==='visible')renderPreview()});
async function setIgnore(){
 const pats=document.getElementById('ignorepats').value.split('\\n').map(s=>s.trim()).filter(Boolean);
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({ignore_global:pats})});
 const d=await res.json();
 const st=document.getElementById('ignstat');
 if(d.ok){st.textContent='✓ saved — '+(pats.length?pats.length+' pattern'+(pats.length>1?'s':''):'ignoring nothing extra');renderPreview();}
 else{st.textContent='✗ '+(d.error||'failed');}}
async function setLang(){
 const l=document.getElementById('ocrlang').value.trim();
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({ocr_lang:l})});
 const d=await res.json();
 document.getElementById('status').textContent=d.ok?'✓ OCR language: '+d.ocr_lang:'✗ '+(d.error||'failed');}
async function setRoot(){
 const r=document.getElementById('root').value.trim();
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({root:r})});
 const d=await res.json();
 document.getElementById('status').textContent=d.ok?'✓ Saved: '+d.root:'✗ '+d.error;
 refresh();}
async function browse(){
 const b=document.getElementById('browsebtn'),st=document.getElementById('status');
 b.disabled=true;b.textContent='Waiting for dialog…';
 try{
  const res=await fetch('/api/pick_folder',{method:'POST'});
  const d=await res.json();
  if(d.ok&&d.path){document.getElementById('root').value=d.path;st.textContent='';}
  else if(d.ok){st.textContent='';}                       // canceled — no drama
  else{st.textContent='✗ '+(d.error||'folder picker failed');}
 }catch(e){st.textContent='✗ folder picker failed: '+(e.message||e);}
 b.disabled=false;b.textContent='Browse…';}
async function setOrigin(){
 const o=document.getElementById('origin').value.trim();
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({allowed_origin:o})});
 const d=await res.json();
 document.getElementById('secstatus').textContent = d.ok?('Security mode: '+d.security+(d.error?' — '+d.error:'')):('✗ '+(d.error||'failed'));
 if(d.ok&&d.security==='UNLOCKED')document.getElementById('secstatus').textContent+=' ⚠ set an origin to unlock file serving';}
async function setToken(){
 const v=document.getElementById('token').value.trim();
 if(!v){document.getElementById('secstatus').textContent='✗ paste a token first (or use Generate)';return;}
 if(v=='••••••••••••'){document.getElementById('secstatus').textContent='Token unchanged — paste a NEW value to replace it';return;}
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({token:{set:v}})});
 const d=await res.json();
 document.getElementById('secstatus').textContent = d.ok?('Token set — Security mode: '+d.security):('✗ '+(d.error||'failed'));}
async function genToken(){
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({token:{generate:true}})});
 const d=await res.json();
 document.getElementById('secstatus').textContent = d.ok?('Token (copy now, shown once): '+d.token+' — mode: '+d.security):('✗ '+(d.error||'failed'));}
async function clearToken(){
 const res=await fetch('/api/root',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({token:{clear:true}})});
 const d=await res.json();
 document.getElementById('secstatus').textContent='Security mode: '+d.security;}
async function stopBridge(){
 if(!confirm('Stop the File Bridge service?\\nOpen WebUI will lose file access until you start it again.'))return;
 try{await fetch('/api/shutdown',{method:'POST'});}catch(e){}
 document.getElementById('beat').style.color='#b00';
 document.getElementById('beatinfo').textContent='stopping… this page goes offline now';}
refresh();
</script></body></html>"""


def _is_readonly() -> bool:
    if os.environ.get("FILE_BRIDGE_READONLY", "").lower() in ("1", "true", "yes"):
        return True
    return bool(_state_load().get("readonly"))


# --------------------------- stderr log + rotation (P2 rollout item)
# The audit log already rotates; this covers the console/request log
# (BaseHTTPRequestHandler writes every request line to sys.stderr). When
# running as a service the journal captures stderr anyway — this is for
# double-click / console / nohup launches. Disable with FILE_BRIDGE_NO_LOGFILE=1.

BRIDGE_LOG_MAX = int(os.environ.get("FILE_BRIDGE_LOG_MAX_BYTES",
                                    str(5 * 1024 * 1024)))


class _RotatingLog:
    """Minimal line-buffered stderr/stdout sink with size-based rotation
    (same .1 shift semantics as the audit log)."""

    def __init__(self, path: Path, max_bytes: int):
        self.path = path
        self.max = max(64 * 1024, max_bytes)
        self._fh = None
        self._writes = 0

    def _rotate(self):
        old = self.path.parent / (self.path.name + ".1")
        try:
            old.unlink()
        except OSError:
            pass
        try:
            self.path.rename(old)
        except OSError:
            pass

    def _open(self):
        try:
            if self.path.exists() and self.path.stat().st_size > self.max:
                self._rotate()
            self._fh = open(self.path, "a", encoding="utf-8", buffering=1)
            try:
                os.chmod(self.path, 0o600)
            except OSError:
                pass
        except OSError:
            self._fh = None

    def write(self, s):
        if self._fh is None:
            self._open()
            if self._fh is None:
                return
        self._writes += 1
        if self._writes % 64 == 0:
            try:
                if self.path.stat().st_size > self.max:
                    self._fh.close()
                    self._rotate()
                    self._fh = open(self.path, "a", encoding="utf-8", buffering=1)
            except OSError:
                pass
        self._fh.write(s)

    def flush(self):
        if self._fh:
            self._fh.flush()


# ------------------------------------- native folder picker (P3, stdlib-only)
# The local settings page's "Browse…" button. Uses each OS's built-in
# choose-folder dialog via subprocess — no GUI toolkit dependency:
#   macOS:   osascript "choose folder"          (always present)
#   Windows: PowerShell FolderBrowserDialog      (always present, -STA)
#   Linux:   zenity, else kdialog               (best effort)
# Only reachable from the loopback picker API — the local user clicking
# the button IS the consent, same stance as /reveal but local-initiated.

_ACTIVE_DIALOGS: set = set()
_DIALOGS_LOCK = threading.Lock()


def _run_dialog(cmd: list) -> tuple[int, bytes, bytes]:
    """Run a native dialog subprocess, registered so a bridge shutdown can
    kill it (Stop button / Dock Quit while a dialog is open must not leave
    an orphaned dialog on screen). Returns (returncode, stdout, stderr)."""
    p = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    with _DIALOGS_LOCK:
        _ACTIVE_DIALOGS.add(p)
    try:
        try:
            out, err = p.communicate(timeout=600)
        except subprocess.TimeoutExpired:
            p.kill()
            out, err = p.communicate()
            raise
    finally:
        with _DIALOGS_LOCK:
            _ACTIVE_DIALOGS.discard(p)
    return p.returncode, out or b"", err or b""


def kill_active_dialogs() -> None:
    """Close any still-open native dialogs (called on bridge shutdown)."""
    with _DIALOGS_LOCK:
        for p in list(_ACTIVE_DIALOGS):
            try:
                p.kill()
            except Exception:
                pass


def pick_folder_dialog(prompt: str = "File Bridge — choose the folder to share"
                       ) -> tuple[str | None, str | None]:
    """Returns (path, error). (None, None) → the user canceled the dialog."""
    try:
        if sys.platform == "darwin":
            script = ("POSIX path of (choose folder with prompt \""
                      + prompt.replace("\\", "\\\\").replace('"', '\\"') + "\")")
            try:
                rc, out_b, err_b = _run_dialog(["osascript", "-e", script])
            except subprocess.TimeoutExpired:
                return None, "folder dialog timed out (10 min) — try again or type the path"
            err = err_b.decode("utf-8", "replace")
            if rc != 0:
                if "cancel" in err.lower():
                    return None, None
                return None, err.strip()[:300] or f"osascript exited {rc}"
            out = out_b.decode("utf-8", "replace").strip()
            if out.endswith("/") and len(out) > 1:   # POSIX path keeps a /
                out = out[:-1]
            return out or None, None
        if _IS_WINDOWS:
            ps = ("& {Add-Type -AssemblyName System.Windows.Forms; "
                  "$d = New-Object System.Windows.Forms.FolderBrowserDialog; "
                  "$d.Description = 'File Bridge: choose the folder to share'; "
                  "if ($d.ShowDialog() -eq [System.Windows.Forms.DialogResult]::OK) "
                  "{ [Console]::OutputEncoding = [System.Text.Encoding]::UTF8; "
                  "[Console]::Out.Write($d.SelectedPath) }")
            try:
                rc, out_b, err_b = _run_dialog(["powershell", "-NoProfile", "-STA",
                                                "-Command", ps])
            except subprocess.TimeoutExpired:
                return None, "folder dialog timed out (10 min) — try again or type the path"
            if rc == 0:
                sel = out_b.decode("utf-8", "replace").strip()
                if sel:
                    return sel, None
                return None, None   # dialog closed without a choice
            return None, (err_b.decode("utf-8", "replace").strip()[:300]
                          or f"powershell exited {rc}")
        import shutil as _sh
        home = os.path.expanduser("~")
        if _sh.which("zenity"):
            try:
                rc, out_b, err_b = _run_dialog(
                    ["zenity", "--file-selection", "--directory",
                     "--filename", home + "/", "--title", prompt])
            except subprocess.TimeoutExpired:
                return None, "folder dialog timed out (10 min) — try again or type the path"
            sel = out_b.decode("utf-8", "replace").strip()
            if rc == 0 and sel:
                return sel, None
            if rc in (0, 1):
                return None, None   # 1 = canceled, empty = canceled
            return None, (err_b.decode("utf-8", "replace").strip()[:300]
                          or f"zenity exited {rc}")
        if _sh.which("kdialog"):
            try:
                rc, out_b, _err_b = _run_dialog(
                    ["kdialog", "--getexistingdirectory", home])
            except subprocess.TimeoutExpired:
                return None, "folder dialog timed out (10 min) — try again or type the path"
            sel = out_b.decode("utf-8", "replace").strip()
            if rc == 0 and sel:
                return sel, None
            return None, None
        return None, ("no folder dialog found — install zenity, or type the "
                      "folder path manually")
    except OSError as e:
        return None, f"could not launch the folder dialog: {e}"


# ------------------------------------- macOS Dock presence (P3, stdlib-only)
# A windowed PyInstaller .app never touches AppKit, so LaunchServices
# classifies the process BackgroundOnly — no Dock icon, no visible sign
# the bridge is running (user feedback 2026-08-28). Bootstrapping
# NSApplication via ctypes fixes that without adding a dependency:
#   - activation policy Regular → the app shows in the Dock while running
#   - Dock / right-click → Quit works: NSApplication's default terminate
#     handles the quit Apple event (verified empirically on arm64)
#   - [NSApp run] parks the MAIN thread; the HTTP server moves to a
#     daemon thread. The settings-page Stop button (and any future need
#     to stop programmatically) breaks the loop with stop: plus
#     CFRunLoopStop(main) — thread-safe, no NSEvent struct marshalling.
#   - an app delegate answers applicationShouldHandleReopen: — clicking
#     the Dock icon (or re-opening the .app) while running pops the
#     settings/status page in the browser (2.6.3, user request).
# ctypes pitfall (cost a segfault in testing): objc_getClass /
# sel_registerName / objc_msgSend MUST have restype/argtypes declared,
# or ctypes truncates 64-bit pointers to c_int.
# Every step is guarded — if anything fails, main() falls back to plain
# serve_forever() and the bridge works exactly as before.

_IS_MAC_FROZEN = sys.platform == "darwin" and bool(getattr(sys, "frozen", False))


class CocoaDock:
    """AppKit bootstrap for the frozen macOS app. .ok False → disabled."""

    def __init__(self):
        self.ok = False
        self._app = None
        self._reopen_fn = None        # set via set_reopen() from main()
        self._reopen_imp = None       # CFUNCTYPE must outlive the delegate
        self._reopen_delegate = None
        self._reopen_last = 0.0       # debounce: a dock double-click = 1 tab
        if not _IS_MAC_FROZEN:
            return
        try:
            import ctypes
            self._objc = ctypes.CDLL("/usr/lib/libobjc.dylib")
            ctypes.CDLL("/System/Library/Frameworks/AppKit.framework/AppKit")
            ob = self._objc
            ob.objc_getClass.restype = ctypes.c_void_p
            ob.objc_getClass.argtypes = [ctypes.c_char_p]
            ob.sel_registerName.restype = ctypes.c_void_p
            ob.sel_registerName.argtypes = [ctypes.c_char_p]
            ob.objc_msgSend.restype = ctypes.c_void_p
            ob.objc_msgSend.argtypes = [ctypes.c_void_p, ctypes.c_void_p]
            app = ob.objc_msgSend(ob.objc_getClass(b"NSApplication"),
                                  ob.sel_registerName(b"sharedApplication"))
            if not app:
                return
            set_pol = ob.objc_msgSend
            set_pol.restype = None
            set_pol.argtypes = [ctypes.c_void_p, ctypes.c_void_p, ctypes.c_long]
            set_pol(app, ob.sel_registerName(b"setActivationPolicy:"), 0)  # Regular
            self._app = app
            self.ok = True
        except Exception:
            self.ok = False
        # Dock-click → settings page (user request 2026-08-29): add-on, never
        # fatal — a failure here leaves the plain 2.4 dock (icon + Quit).
        if self.ok:
            self._install_reopen_delegate()

    # ---- Dock click / Finder double-click on the RUNNING app ------------
    # LaunchServices does not spawn a second process for an already-running
    # bundle; it activates the instance and sends it a "reopen" Apple event,
    # delivered to the application delegate as
    # applicationShouldHandleReopen:hasVisibleWindows:. Without a delegate
    # the click does nothing at all. We build one from raw ctypes — same
    # stdlib-only discipline as the rest of CocoaDock:
    #   NSObject subclass via objc_allocateClassPair + class_addMethod(IMP)
    #   → [NSApp setDelegate:]. The IMP runs ON the AppKit main thread (the
    #   run loop dispatches it), so it must not block or fork there: it
    #   spawns a Python thread and returns YES immediately.

    def set_reopen(self, fn):
        """fn() is called (off the main thread) when the user clicks the
        app's Dock icon or re-opens the .app while it is running."""
        self._reopen_fn = fn

    def _install_reopen_delegate(self):
        try:
            import ctypes
            ob = ctypes.CDLL("/usr/lib/libobjc.dylib")
            ob.objc_getClass.restype = ctypes.c_void_p
            ob.objc_getClass.argtypes = [ctypes.c_char_p]
            ob.sel_registerName.restype = ctypes.c_void_p
            ob.sel_registerName.argtypes = [ctypes.c_char_p]
            sel = ob.sel_registerName
            ob.objc_allocateClassPair.restype = ctypes.c_void_p
            ob.objc_allocateClassPair.argtypes = [ctypes.c_void_p,
                                                  ctypes.c_char_p,
                                                  ctypes.c_size_t]
            ob.class_addMethod.restype = ctypes.c_bool
            ob.class_addMethod.argtypes = [ctypes.c_void_p, ctypes.c_void_p,
                                           ctypes.c_void_p, ctypes.c_char_p]
            ob.objc_registerClassPair.restype = None
            ob.objc_registerClassPair.argtypes = [ctypes.c_void_p]

            # BOOL imp(id self, SEL _cmd, id sender, BOOL hasVisibleWindows)
            REOPEN_IMP = ctypes.CFUNCTYPE(
                ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p,
                ctypes.c_void_p, ctypes.c_bool)

            def _imp(_self, _cmd, _sender, _has_vis):
                now = time.monotonic()
                if now - self._reopen_last < 2.0:   # double-click = 1 tab
                    return True
                self._reopen_last = now
                fn = self._reopen_fn
                print("Dock icon clicked — opening the File Bridge page")
                if fn is not None:
                    threading.Thread(target=fn, daemon=True).start()
                return True

            imp = REOPEN_IMP(_imp)
            self._reopen_imp = imp   # keep the C trampoline alive (GC-proof)

            cls = ob.objc_allocateClassPair(ob.objc_getClass(b"NSObject"),
                                            b"FileBridgeReopenDelegate", 0)
            if not cls:
                return
            if not ob.class_addMethod(
                    cls,
                    sel(b"applicationShouldHandleReopen:hasVisibleWindows:"),
                    ctypes.cast(imp, ctypes.c_void_p), b"B@:@B"):
                return
            ob.objc_registerClassPair(cls)
            msg = ob.objc_msgSend
            msg.restype = ctypes.c_void_p
            msg.argtypes = [ctypes.c_void_p, ctypes.c_void_p]
            delegate = msg(msg(cls, sel(b"alloc")), sel(b"init"))
            if not delegate:
                return
            set_d = ob.objc_msgSend
            set_d.restype = None
            set_d.argtypes = [ctypes.c_void_p, ctypes.c_void_p, ctypes.c_void_p]
            set_d(self._app, sel(b"setDelegate:"), delegate)
            self._reopen_delegate = delegate   # NSApplication does NOT
            # retain its delegate; holding the raw pointer plus zero
            # releases anywhere keeps it alive for the process lifetime.
        except Exception:
            self._reopen_delegate = None

    def run_forever(self):
        """Park the main thread in the AppKit event loop until stopped."""
        if not self.ok:
            return
        import ctypes
        ob = ctypes.CDLL("/usr/lib/libobjc.dylib")   # fresh handle: argtypes
        ob.objc_getClass.restype = ctypes.c_void_p   # mutations must not race
        ob.sel_registerName.restype = ctypes.c_void_p
        f = ob.objc_msgSend
        f.restype = None
        f.argtypes = [ctypes.c_void_p, ctypes.c_void_p]
        f(self._app, ob.sel_registerName(b"run"))

    def request_stop(self):
        """Break the AppKit loop (safe from any thread). No-op if not ok.
        The canonical pyobjc recipe: [NSApp stop:] sets a flag that the
        run loop only checks WHEN IT PROCESSES AN EVENT — an idle loop
        would sleep forever (verified: performSelectorOnMainThread alone
        does not wake it). So: stop: + performSelectorOnMainThread for
        good measure, then an application-defined no-op NSEvent to wake
        the loop and make it observe the flag."""
        if not self.ok:
            return
        try:
            import ctypes

            class _NSPoint(ctypes.Structure):
                _fields_ = [("x", ctypes.c_double), ("y", ctypes.c_double)]

            ob = ctypes.CDLL("/usr/lib/libobjc.dylib")
            ob.objc_getClass.restype = ctypes.c_void_p
            ob.sel_registerName.restype = ctypes.c_void_p
            sel = ob.sel_registerName
            # 1. set the stop flag (idempotent, cheap ivar write)
            f_stop = ob.objc_msgSend
            f_stop.restype = None
            f_stop.argtypes = [ctypes.c_void_p, ctypes.c_void_p, ctypes.c_void_p]
            f_stop(self._app, sel(b"stop:"), None)
            # 2. also run stop: on the main thread (documented-safe path)
            f_perf = ob.objc_msgSend
            f_perf.restype = None
            f_perf.argtypes = [ctypes.c_void_p, ctypes.c_void_p,
                               ctypes.c_void_p, ctypes.c_void_p, ctypes.c_bool]
            f_perf(self._app,
                   sel(b"performSelectorOnMainThread:withObject:waitUntilDone:"),
                   sel(b"stop:"), None, False)
            # 3. wake the loop with a no-op event so it observes the flag
            #    (+[NSEvent otherEventWithType:...] is a CLASS method —
            #    calling it on an alloc'd instance silently no-ops)
            f_ev = ob.objc_msgSend
            f_ev.restype = ctypes.c_void_p
            f_ev.argtypes = [ctypes.c_void_p, ctypes.c_void_p,
                             ctypes.c_longlong, _NSPoint, ctypes.c_ulonglong,
                             ctypes.c_double, ctypes.c_longlong, ctypes.c_void_p,
                             ctypes.c_short, ctypes.c_longlong, ctypes.c_longlong]
            ev = f_ev(ob.objc_getClass(b"NSEvent"),
                      sel(b"otherEventWithType:location:modifierFlags:"
                          b"timestamp:windowNumber:context:subtype:"
                          b"data1:data2:"),
                      15, _NSPoint(0, 0), 0, 0.0, 0, None, 0, 0, 0)  # AppDefined
            if not ev:
                return
            f_post = ob.objc_msgSend
            f_post.restype = None
            f_post.argtypes = [ctypes.c_void_p, ctypes.c_void_p,
                               ctypes.c_void_p, ctypes.c_bool]
            f_post(self._app, sel(b"postEvent:atStart:"), ev, True)
        except Exception:
            pass


_SHUTDOWN_FN = None   # set in main(): called from a server thread to stop


def _is_file_bridge_url(url: str) -> bool:
    """True if the listener at url answers /version like a File Bridge
    (token-free endpoint) — tells a live bridge apart from whatever else
    might hold the port, before we open a browser tab at it."""
    try:
        import urllib.request
        with urllib.request.urlopen(f"{url.rstrip('/')}/version",
                                    timeout=1.5) as r:
            return r.status == 200 and b'"bridge"' in r.read(4096)
    except Exception:
        return False


def _user_alert(message: str, title: str = "File Bridge") -> None:
    """Visible warning for startup aborts that would otherwise be silent.

    Packaged launches (.app / windowed exe) have no console, so print()
    only reaches state_dir/bridge.log — the user clicks the icon and
    NOTHING happens, with the reason buried in a log they don't know
    exists (user request 2026-08-29: warn, don't abort silently). Native,
    best-effort, never fatal:
      macOS   osascript dialog (always present; registered via
              _run_dialog so a shutdown can't orphan it on screen)
      Windows MessageBoxW through ctypes (stdlib — no GUI dep, same
              discipline as CocoaDock)
      Linux   notify-send / zenity / kdialog, first one installed
    Skipped when the message is already visible on a terminal (source
    runs), and under FILE_BRIDGE_NO_UI=1 (services must not throw modal
    dialogs at a login screen).
    """
    if not getattr(sys, "frozen", False):
        return
    if os.environ.get("FILE_BRIDGE_NO_UI"):
        return
    try:
        if _IS_WINDOWS:
            import ctypes
            MB_ICONWARNING, MB_SETFOREGROUND = 0x30, 0x10000
            ctypes.windll.user32.MessageBoxW(
                0, message, title, MB_ICONWARNING | MB_SETFOREGROUND)
        elif sys.platform == "darwin":
            # json.dumps produces a legal AppleScript string literal
            # (quotes/backslashes escaped), \n included.
            _run_dialog(["osascript", "-e",
                         f'display dialog {json.dumps(message)} '
                         f'with title {json.dumps(title)} '
                         f'buttons {{"OK"}} default button "OK" '
                         f'with icon caution'])
        else:
            for cmd in (["notify-send", "-a", title, message],
                        ["zenity", "--error", f"--title={title}",
                         f"--text={message}"],
                        ["kdialog", "--title", title, "--error", message]):
                try:
                    subprocess.run(cmd, timeout=60)
                    break
                except (OSError, subprocess.SubprocessError):
                    continue
    except Exception:
        pass   # the bridge.log line is the fallback record


def main():
    # PyInstaller --windowed (shipped Windows exe; console=False) leaves
    # sys.stdout/sys.stderr as None on Windows when there is no console —
    # any print()/isatty() on None would crash at startup (and a windowed
    # traceback dialog then hangs the process). This must run BEFORE the
    # first print below. macOS .app launches get /dev/null fds instead, so
    # only Windows hits the None case.
    if sys.stdout is None:
        sys.stdout = open(os.devnull, "w")
    if sys.stderr is None:
        sys.stderr = open(os.devnull, "w")

    folder = sys.argv[1] if len(sys.argv) > 1 else None
    if folder:
        p = Path(folder).expanduser().resolve()
        if not p.is_dir():
            print(f"error: {folder} is not a folder")
            _user_alert(f"The folder given to File Bridge does not exist:\n"
                        f"{folder}\n\nStart File Bridge again and pick a "
                        f"folder on its settings page.")
            sys.exit(1)
        save_root(p)
        print(f"Sharing: {p}")

    root = load_root()

    # console/request log → state_dir/bridge.log when running non-interactive
    # (service / nohup); interactive terminal keeps plain stderr.
    if (not sys.stdout.isatty() and not sys.stderr.isatty()
            and not os.environ.get("FILE_BRIDGE_NO_LOGFILE")):
        logf = _RotatingLog(STATE_DIR / "bridge.log", BRIDGE_LOG_MAX)
        class _T:
            def __init__(self, sink, orig):
                self.sink, self.orig = sink, orig
            def write(self, s):
                try:
                    self.sink.write(s)
                except Exception:
                    pass
                try:
                    self.orig.write(s)
                except Exception:
                    pass
            def flush(self):
                try:
                    self.sink.flush()
                except Exception:
                    pass
                try:
                    self.orig.flush()
                except Exception:
                    pass
            def isatty(self):
                return False
        sys.stderr = _T(logf, sys.stderr)
        sys.stdout = _T(logf, sys.stdout)

    # serve picker page (local setup UI)
    orig_do_GET = Handler.do_GET
    def do_GET_picker(self):
        if urlparse(self.path).path in ("/", "/picker"):
            root = load_root()
            html = Handler.PICKER_HTML.replace("__ROOT__", str(root) if root else "")
            html = html.replace("__ORIGIN__", get_allowed_origin() or "")
            html = html.replace("__OCRLANG__", _get_ocr_lang())
            html = html.replace("__IGNORE__", _hesc("\n".join(_global_ignore())))
            html = html.replace("__STATUS__",
                                f"sharing {root} · security {security_mode()}" if root
                                else f"no folder chosen yet · security {security_mode()}")
            # placeholder shows a path shaped like THIS os (user feedback:
            # C:\Users\... on a Mac reads as a Windows-only tool)
            if _IS_WINDOWS:
                rootph = "C:\\Users\\you\\Documents\\my-folder"
            elif sys.platform == "darwin":
                rootph = "/Users/you/Documents/my-folder"
            else:
                rootph = "/home/you/Documents/my-folder"
            html = html.replace("__ROOTPH__", rootph)
            body = html.encode()
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)
            return
        return orig_do_GET(self)
    Handler.do_GET = do_GET_picker

    # second start while a bridge is already listening. Windows: THE path —
    # the exe has no tray icon, so "click the app icon again" IS a second
    # process (macOS rarely gets here: LaunchServices activates the running
    # app instead, handled by the Dock reopen delegate). We must never bind
    # anyway: with SO_REUSEADDR a second instance can end up sharing the
    # listen queue and connections land on a process that never serves them
    # (verified: two frozen instances → /version answers, /list hangs
    # forever). So: verify the listener is one of ours (/version), pop its
    # settings page, exit 0. Only a foreign listener still errors.
    import socket as _sock
    try:
        _probe = _sock.create_connection(("127.0.0.1", PORT), timeout=1)
        _probe.close()
        _already = True
    except OSError:
        _already = False
    if _already:
        url = f"http://127.0.0.1:{PORT}"
        if _is_file_bridge_url(url):
            if os.environ.get("FILE_BRIDGE_NO_UI"):
                print(f"File Bridge already running at {url} — this copy "
                      f"exits (browser suppressed: FILE_BRIDGE_NO_UI).")
            else:
                print(f"File Bridge is already running at {url} — opening "
                      f"its page in your browser; this copy exits (the "
                      f"running one keeps serving).")
                try:
                    webbrowser.open(url)
                except Exception:
                    pass
            sys.exit(0)
        print(f"error: port {PORT} is held by something that is not a File "
              f"Bridge — set FILE_BRIDGE_PORT to move the bridge elsewhere")
        _user_alert(f"File Bridge cannot start: port {PORT} is being used "
                    f"by another program.\n\nClose that program and start "
                    f"File Bridge again, or move File Bridge to another "
                    f"port (FILE_BRIDGE_PORT) and re-run "
                    f"scripts/setup_owui.py so the skill follows.")
        sys.exit(1)

    socketserver.ThreadingTCPServer.allow_reuse_address = True
    # daemon request threads: Stop/Dock-Quit must exit immediately instead of
    # hanging in server_close() joining a thread that is parked on a
    # 10-minute native folder dialog (ThreadingHTTPServer does the same).
    socketserver.ThreadingTCPServer.daemon_threads = True
    try:
        httpd = socketserver.ThreadingTCPServer(("127.0.0.1", PORT), Handler)
    except OSError:
        # rare race only (the port was lost between the probe and the
        # bind): re-verify who holds it — one of ours → hand the user to
        # the live settings page; anything else → the foreign-holder
        # warning above, not a windowed traceback dialog.
        url0 = f"http://127.0.0.1:{PORT}"
        if _is_file_bridge_url(url0):
            print(f"File Bridge is already running at {url0} — "
                  f"this second copy exits.")
            if not os.environ.get("FILE_BRIDGE_NO_UI"):
                try:
                    webbrowser.open(url0)
                except Exception:
                    pass
            return
        print(f"error: port {PORT} was taken by a non-File-Bridge program "
              f"at bind time — set FILE_BRIDGE_PORT to move the bridge")
        _user_alert(f"File Bridge cannot start: port {PORT} is being used "
                    f"by another program.\n\nClose that program and start "
                    f"File Bridge again, or move File Bridge to another "
                    f"port (FILE_BRIDGE_PORT) and re-run "
                    f"scripts/setup_owui.py so the skill follows.")
        sys.exit(1)

    # macOS .app: Dock presence via the CocoaDock bootstrap — the server
    # moves to a daemon thread and the main thread parks in the AppKit
    # loop (Dock icon + working Quit). Everywhere else: serve_forever on
    # the main thread exactly as before. Both modes stop cleanly via
    # /api/shutdown (settings page Stop button).
    global _SHUTDOWN_FN
    dock = CocoaDock()
    url = f"http://127.0.0.1:{PORT}"
    mode = security_mode()
    dock.set_reopen(lambda: webbrowser.open(url))   # Dock/Finder click → UI
    if dock.ok:

        def _stop_cocoa():
            dock.request_stop()      # break the AppKit loop (main thread)
            httpd.shutdown()         # stop serving (Timer thread, ≤0.5 s)
        _SHUTDOWN_FN = _stop_cocoa
        threading.Thread(target=httpd.serve_forever, daemon=True).start()
    else:
        _SHUTDOWN_FN = httpd.shutdown

    print(f"File Bridge v{VERSION} running at {url}  (Ctrl+C to stop)")
    print(f"Security mode: {mode}")
    if mode == "UNLOCKED":
        print("⚠ UNLOCKED: file endpoints are DISABLED until you set the allowed")
        print("  Open WebUI origin (and optionally a token) in the settings page.")
    # Clicking the app icon should always surface the UI (user request
    # 2026-08-29), not only on first run with no folder configured. Scoped
    # to icon-style launches: a folder ARGUMENT means a scripted/CLI start
    # (skill, tests, shortcuts that pin a folder know what they opened),
    # and services set FILE_BRIDGE_NO_UI=1 so a login autostart stays
    # silent. Cold launch opens here; a click on the ALREADY-running app
    # arrives as the Dock reopen event (see CocoaDock.set_reopen).
    _ui_launch = (bool(getattr(sys, "frozen", False)) and not folder
                  and not os.environ.get("FILE_BRIDGE_NO_UI"))
    if not root:
        print("No folder set — opening the setup page in your browser...")
    elif _ui_launch:
        print("App started from its icon — opening the File Bridge page...")
    if not root or _ui_launch:
        threading.Timer(0.5, lambda: webbrowser.open(url)).start()
    if dock.ok:
        print("Running with a Dock icon — click it any time to open the File"
              " Bridge page; Quit from the Dock, or the Stop button on the"
              " settings page, ends the bridge.")
        dock.run_forever()
    else:
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            pass
    kill_active_dialogs()   # no orphaned choose-folder dialog after exit
    httpd.server_close()
    print("\nFile Bridge stopped.")


if __name__ == "__main__":
    main()
